#!/usr/bin/env python3
"""
Bittensor AI Chat v3 — Pywebview Desktop App
Powered by Chutes.ai API — Clean SVG Icons

Optimized for Linux.
"""

import sys
import os
import json
import threading
import uuid
import base64
import re
import hashlib
import time
import shutil
import tempfile
import unicodedata
import traceback
import html
from datetime import datetime
import webview

import subprocess

# Optional dependencies with graceful fallbacks
try:
    from openai import OpenAI
except ImportError:
    print("Error: openai library required. Install: pip install openai")
    sys.exit(1)

try:
    import requests as _requests
    # Create a session for connection pooling (performance optimization)
    _requests_session = _requests.Session()
    _requests_session.headers.update({"User-Agent": "BittensorChat/3.0"})

except ImportError:
    _requests = None
    _requests_session = None

try:
    from pypdf import PdfReader
except ImportError:
    PdfReader = None

try:
    import pandas as pd
except ImportError:
    pd = None

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation as PptxPresentation
except ImportError:
    PptxPresentation = None

DATA_DIR = os.path.expanduser("~/.bittensor_chat")
APP_TITLE = "Bittensor Chat"
CHATS_DIR = os.path.join(DATA_DIR, "chats")
CONFIG_FILE = os.path.join(DATA_DIR, "config.json")
PROFILES_DIR = os.path.join(DATA_DIR, "profiles")
PROFILES_FILE = os.path.join(PROFILES_DIR, "profiles.json")
os.makedirs(CHATS_DIR, exist_ok=True)
os.makedirs(PROFILES_DIR, exist_ok=True)

def load_profiles_manifest():
    if os.path.exists(PROFILES_FILE):
        try:
            with open(PROFILES_FILE, "r") as f: return json.load(f)
        except: pass
    return []

def save_profiles_manifest(profiles):
    with open(PROFILES_FILE, "w") as f: json.dump(profiles, f, indent=2)

def hash_password(password, salt):
    dk = hashlib.scrypt(password.encode('utf-8'), salt=salt.encode('utf-8'), n=16384, r=8, p=1)
    return 'scrypt$' + dk.hex()

MASTER_KEY = "admin"

DEFAULT_CONFIG = {
    "api_key": "", "brave_api_key": "",
    "base_url": "https://e2ee-local-proxy.chutes.dev:8443/v1",
    "max_tokens": 10000, "max_memory": 50, "theme": "void",
    "current_model": "moonshotai/Kimi-K2.5-TEE",
    "show_token_budget": True,
    "models": [
        "moonshotai/Kimi-K2.5-TEE", "deepseek-ai/DeepSeek-V3-0324",
        "Qwen/Qwen3-235B-A22B", "meta-llama/Llama-4-Maverick-17B-128E-Instruct",
        "mistralai/Mistral-Small-3.2-24B-Instruct-2506",
    ],
    "folders": ["General", "Work", "Research", "Personal", "Projects"],
    "chat_folders": {},
    "folder_colors": {"General":"#555566","Work":"#3b82f6","Research":"#8b5cf6","Personal":"#10b981","Projects":"#f59e0b"},
    "prompt_templates": [],
}




def load_config():
    if os.path.exists(CONFIG_FILE):
        try:
            with open(CONFIG_FILE, "r") as f:
                cfg = json.load(f)
                for k, v in DEFAULT_CONFIG.items(): cfg.setdefault(k, v)
                return cfg
        except (json.JSONDecodeError, IOError): pass
    return dict(DEFAULT_CONFIG)

def save_config(cfg):
    with open(CONFIG_FILE, "w") as f: json.dump(cfg, f, indent=2)

def chat_path(cid): return os.path.join(CHATS_DIR, f"{cid}.json")
def load_chat(cid):
    p = chat_path(cid)
    if os.path.exists(p):
        with open(p, "r") as f: return json.load(f)
    return None
def save_chat(d):
    with open(chat_path(d["id"]), "w") as f: json.dump(d, f, indent=2)
def delete_chat_file(cid):
    p = chat_path(cid)
    if os.path.exists(p): os.remove(p)

def list_chats():
    """List all chats with caching for performance.
    
    Returns:
        List of chat metadata dicts, sorted by pinned status then creation date.
    """
    chats = []
    try:
        if not os.path.exists(CHATS_DIR):
            return []
        for fn in os.listdir(CHATS_DIR):
            if fn.endswith(".json"):
                try:
                    with open(os.path.join(CHATS_DIR, fn), "r", encoding="utf-8") as f:
                        d = json.load(f)
                        chats.append({
                            "id": d["id"],
                            "title": d.get("title", "Untitled"),
                            "created": d.get("created", ""),
                            "model": d.get("model", ""),
                            "pinned": d.get("pinned", False),
                            "branched": d.get("branched", False)
                        })
                except (json.JSONDecodeError, IOError, KeyError):
                    continue  # Skip corrupted files
    except OSError:
        return []  # Directory may not exist
    
    # Sort: pinned first, then by creation date (newest first)
    pinned = sorted([c for c in chats if c.get("pinned")], 
                    key=lambda x: x.get("created", ""), reverse=True)
    unpinned = sorted([c for c in chats if not c.get("pinned")], 
                      key=lambda x: x.get("created", ""), reverse=True)
    return pinned + unpinned

def read_file(fp, max_size_mb=50):
    """Read file content with size validation.
    
    Args:
        fp: File path to read.
        max_size_mb: Maximum file size in MB (default: 50).
        
    Returns:
        File content as string, or error message.
    """
    # Validate file size
    try:
        file_size = os.path.getsize(fp)
        max_size = max_size_mb * 1024 * 1024
        if file_size > max_size:
            return f"[Error] File exceeds {max_size_mb}MB limit."
    except OSError as e:
        return f"[Error] Cannot access file: {e}"
    
    ext = os.path.splitext(fp)[1].lower()
    try:
        if ext == ".pdf":
            if not PdfReader: return "[Error] pypdf not installed."
            return "\n".join([p.extract_text() for p in PdfReader(fp).pages if p.extract_text()])
        elif ext in (".xlsx",".xls"):
            if not pd: return "[Error] pandas not installed."
            return pd.read_excel(fp).to_string()
        elif ext == ".csv":
            if not pd: return "[Error] pandas not installed."
            return pd.read_csv(fp).to_string()
        elif ext == ".docx":
            if not DocxDocument: return "[Error] python-docx not installed."
            return "\n".join([p.text for p in DocxDocument(fp).paragraphs])
        elif ext == ".pptx":
            if not PptxPresentation: return "[Error] python-pptx not installed."
            prs = PptxPresentation(fp)
            t = []
            for s in prs.slides:
                for sh in s.shapes:
                    if hasattr(sh,"text") and sh.text: t.append(sh.text)
            return "\n".join(t)
        else:
            with open(fp, "r", encoding="utf-8", errors="ignore") as f: return f.read()
    except Exception as e: return f"[Error] {e}"

def read_image_base64(fp):
    ext = os.path.splitext(fp)[1].lower()
    mm = {".png":"image/png",".jpg":"image/jpeg",".jpeg":"image/jpeg",".gif":"image/gif",".webp":"image/webp",".bmp":"image/bmp"}
    with open(fp, "rb") as f: data = base64.b64encode(f.read()).decode("utf-8")
    return {"base64":data,"media_type":mm.get(ext,"image/png"),"filename":os.path.basename(fp)}

def brave_search(query, api_key):
    """Perform web search using Brave API with connection pooling."""
    if not _requests_session: return "requests library not installed."
    if not api_key: return "Brave API key not set. Add it in Settings."
    try:
        resp = _requests_session.get(
            "https://api.search.brave.com/res/v1/web/search",
            headers={"Accept": "application/json", "X-Subscription-Token": api_key},
            params={"q": query, "count": 5},
            timeout=10
        )
        resp.raise_for_status()
        results = resp.json().get("web", {}).get("results", [])
        if not results: return "No results found."
        return "\n\n".join([
            f"**{r.get('title', '')}**\nURL: {r['url']}\n{r.get('description', '')}"
            for r in results
        ])
    except Exception as e:
        return f"Search error: {e}"

class Api:
    def __init__(self):
        self.config = load_config()
        self._client = None
        self._window = None
        self._active_streams = {}  # cid -> threading.Event for abort signal
        self._stream_lock = threading.Lock()  # thread safety for stream management
    def _get_client(self):
        if self._client is None:
            self._client = OpenAI(api_key=self.config.get("api_key",""), base_url=self.config.get("base_url",DEFAULT_CONFIG["base_url"]))
        return self._client
    def get_config(self): return self.config
    def save_settings(self, settings):
        self.config.update(settings)
        save_config(self.config)
        self._client = None
        return True
    def get_models(self): return self.config.get("models",[])
    def get_current_model(self): return self.config.get("current_model","")
    def set_current_model(self, m):
        self.config["current_model"] = m
        save_config(self.config)
        return True
    def add_model(self, mn):
        if mn and mn not in self.config["models"]:
            self.config["models"].append(mn)
            save_config(self.config)
        return self.config["models"]
    def remove_model(self, mn):
        if mn in self.config["models"]:
            self.config["models"].remove(mn)
            if self.config["current_model"] == mn:
                self.config["current_model"] = self.config["models"][0] if self.config["models"] else ""
            save_config(self.config)
        return self.config["models"]
    def list_chats(self):
        chats = list_chats()
        folder_map = self.config.get("chat_folders", {})
        for c in chats:
            c["folder"] = folder_map.get(c["id"], "General")
        return chats
    def get_chat(self, cid): return load_chat(cid)
    def save_reasoning_duration(self, cid, duration):
        c = load_chat(cid)
        if not c: return
        for m in reversed(c.get("messages", [])):
            if m.get("role") == "assistant" and m.get("reasoning") and not m.get("reasoning_duration"):
                m["reasoning_duration"] = int(duration)
                break
        save_chat(c)

    def create_chat(self, title=None, system_prompt=None):
        cid = str(uuid.uuid4())[:12]
        sp = system_prompt or "You are a helpful AI assistant. Provide clear, well-formatted answers."
        d = {"id":cid,"title":title or "New Chat","created":datetime.now().isoformat(),"model":self.config.get("current_model",""),"pinned":False,"system_prompt":sp,"messages":[{"role":"system","content":sp}]}
        save_chat(d)
        return d
    def delete_chat(self, cid):
        delete_chat_file(cid)
        return True
    def delete_all_chats(self):
        for fn in os.listdir(CHATS_DIR):
            if fn.endswith(".json"):
                os.remove(os.path.join(CHATS_DIR, fn))
        return True
    def rename_chat(self, cid, title):
        c = load_chat(cid)
        if c:
            c["title"] = title
            save_chat(c)
        return True
    def toggle_pin(self, cid):
        c = load_chat(cid)
        if c:
            c["pinned"] = not c.get("pinned", False)
            save_chat(c)
            return c["pinned"]
        return False
    def update_system_prompt(self, cid, prompt):
        c = load_chat(cid)
        if c:
            c["system_prompt"]=prompt
            if c["messages"] and c["messages"][0]["role"] == "system":
                c["messages"][0]["content"] = prompt
            else:
                c["messages"].insert(0, {"role": "system", "content": prompt})
            save_chat(c)
        return True
    def branch_chat(self, cid, up_to_index):
        orig = load_chat(cid)
        if not orig: return None
        nid = str(uuid.uuid4())[:12]
        nc = {"id":nid,"title":orig['title'],"created":datetime.now().isoformat(),"model":orig.get("model",self.config.get("current_model","")),"pinned":False,"branched":True,"system_prompt":orig.get("system_prompt",""),"messages":list(orig["messages"][:up_to_index+1])}
        save_chat(nc)
        return nc
    def tts_speak(self, text, voice="af_bella"):
        """Generate TTS audio via Chutes Kokoro API, return base64 audio."""
        try:
            # Strip HTML tags and clean markdown for clean speech
            clean = re.sub(r'<[^>]+>', '', text)
            clean = re.sub(r'\*\*(.+?)\*\*', r'\1', clean)
            clean = re.sub(r'\*(.+?)\*', r'\1', clean)
            clean = re.sub(r'`[^`\n]+`', '', clean)
            clean = re.sub(r'```[\s\S]*?```', '', clean)
            clean = re.sub(r'#+\s*', '', clean)
            clean = re.sub(r'\[([^\]]+)\]\([^\)]+\)', r'\1', clean)
            clean = clean.strip()[:4000]
            if not clean:
                return {"error": "No text to speak"}
            api_key = self.config.get("api_key", "")
            resp = _requests_session.post(
                "https://chutes-kokoro.chutes.ai/speak",
                headers={"Authorization": f"Bearer {api_key}", "Content-Type": "application/json"},
                json={"text": clean, "voice": voice},
                timeout=60
            )
            if resp.status_code == 200:
                return {"audio": base64.b64encode(resp.content).decode(), "format": "wav"}
            return {"error": f"TTS error {resp.status_code}: {resp.text[:200]}"}
        except Exception as e:
            return {"error": str(e)}
    def export_chat(self, cid, fmt):
        c = load_chat(cid)
        if not c: return None
        lines = [f"# {c['title']}", f"Model: {c.get('model','')}", f"Date: {c.get('created','')}", ""]
        for m in c["messages"]:
            if m["role"]=="system": continue
            role = "**You**" if m["role"]=="user" else f"**{m['role'].title()}**"
            content = m["content"] if isinstance(m["content"],str) else str(m["content"])
            lines.append(f"{role}:\n{content}\n")
        ext = "md" if fmt=="md" else "txt"
        safe = re.sub(r'[^\w\s-]','',c['title'])[:40].strip()
        fn = f"{safe}.{ext}"
        tmp = os.path.join(DATA_DIR, fn)
        with open(tmp,"w",encoding="utf-8") as f: f.write("\n".join(lines))
        return {"path":tmp,"filename":fn}
    def export_chat_html(self, cid):
        """Export chat as a styled HTML file."""
        c = load_chat(cid)
        if not c: return None
        html_mod = html
        msgs_html = ""
        for m in c["messages"]:
            if m["role"] == "system": continue
            role = m["role"]
            content = m["content"] if isinstance(m["content"], str) else str(m["content"])
            safe_content = html_mod.escape(content).replace('\n', '<br>')
            role_label = "You" if role == "user" else f"Assistant ({c.get('model','AI')})"
            bg = "#1a1a2e" if role == "user" else "#0e0e1a"
            border = "#e88a2a" if role == "user" else "#3b3b6b"
            msgs_html += f'<div style="margin-bottom:18px;padding:16px 20px;background:{bg};border-left:3px solid {border};border-radius:10px"><div style="font-size:11px;font-weight:700;letter-spacing:.08em;color:#888;margin-bottom:8px;text-transform:uppercase">{html_mod.escape(role_label)}</div><div style="font-size:14px;line-height:1.7;white-space:pre-wrap;word-break:break-word">{safe_content}</div></div>'
        title_safe = html_mod.escape(c.get('title','Chat'))
        model_safe = html_mod.escape(c.get('model',''))
        date_safe = html_mod.escape(c.get('created','')[:19] if c.get('created') else '')
        html_content = f"""<!DOCTYPE html>
<html lang="en"><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1">
<title>{title_safe}</title>
<style>
*{{margin:0;padding:0;box-sizing:border-box}}
body{{font-family:'Segoe UI',system-ui,sans-serif;background:#06060b;color:#d8d8e8;padding:30px 20px;max-width:860px;margin:0 auto}}
h1{{font-size:22px;font-weight:700;margin-bottom:6px;color:#fff}}
.meta{{font-size:12px;color:#666;margin-bottom:28px}}
code{{background:#1e1e34;padding:2px 6px;border-radius:4px;font-family:'JetBrains Mono',monospace;font-size:12px;color:#e88a2a}}
pre{{background:#1e1e34;padding:14px 16px;border-radius:8px;overflow-x:auto;margin:10px 0}}
pre code{{background:none;padding:0;color:#c8d8e8;font-size:12px}}
@media print{{body{{background:#fff;color:#000}}h1{{color:#000}}.meta{{color:#666}}}}
</style></head><body>
<h1>{title_safe}</h1>
<div class="meta">Model: {model_safe} &nbsp;|&nbsp; {date_safe}</div>
{msgs_html}
</body></html>"""
        safe = re.sub(r'[^\w\s-]','',c['title'])[:40].strip()
        fn = f"{safe}.html"
        tmp = os.path.join(DATA_DIR, fn)
        with open(tmp,"w",encoding="utf-8") as f: f.write(html_content)
        return {"path":tmp,"filename":fn}
    def export_chat_pdf(self, cid):
        """Export chat as PDF (pure Python, no external dependencies)."""
        tb = traceback
        try:
         c = load_chat(cid)
         if not c: return None
         messages = [m for m in c.get("messages", []) if m.get("role") != "system"]
         title = c.get('title','Chat') or 'Chat'
         model = c.get('model','') or ''
         date = c.get('created','')[:19] if c.get('created') else ''

         def sanitize(s):
             s = unicodedata.normalize('NFKD', str(s))
             result = []
             for ch in s:
                 try:
                     ch.encode('latin-1')
                     if ord(ch) >= 32 or ch == '\t':
                         result.append(ch)
                 except (UnicodeEncodeError, ValueError):
                     pass
             return ''.join(result)

         def ps_escape(s):
             s = sanitize(s)
             s = s.replace('\\','\\\\').replace('(','\\(').replace(')','\\)')
             s = s.replace('\r','').replace('\n', ' ')
             return s.encode('latin-1', 'replace')

         PAGE_W, PAGE_H = 595, 842
         MARGIN = 50
         TEXT_W = PAGE_W - 2 * MARGIN

         def wrap_text(text, font_size, width):
             avg_cw = font_size * 0.52
             max_chars = max(1, int(width / avg_cw))
             if not text.strip():
                 return ['']
             words = text.split(' ')
             lines2, cur = [], ''
             for w in words:
                 test = (cur + ' ' + w).lstrip() if cur else w
                 if len(test) <= max_chars:
                     cur = test
                 else:
                     if cur:
                         lines2.append(cur)
                     while len(w) > max_chars:
                         lines2.append(w[:max_chars]); w = w[max_chars:]
                     cur = w
             if cur:
                 lines2.append(cur)
             return lines2 or ['']

         all_cmds = []
         page_cmds = []
         y = PAGE_H - MARGIN - 16

         def maybe_new_page():
             nonlocal page_cmds, y
             if y < MARGIN + 30:
                 all_cmds.append(page_cmds)
                 page_cmds = []
                 y = PAGE_H - MARGIN - 16

         def add_line(font, size, x, text, line_h=None):
             nonlocal y
             maybe_new_page()
             if text.strip() or font == 'Helvetica':
                 page_cmds.append((font, size, x, y, text))
             y -= (line_h or (size + 5))

         for wl in wrap_text(sanitize(title), 16, TEXT_W):
             add_line('Helvetica-Bold', 16, MARGIN, wl, 22)
         meta = f"Model: {sanitize(model)}   |   {sanitize(date)}"
         add_line('Helvetica', 9, MARGIN, meta, 16)
         y -= 8

         for msg in messages:
             role = msg.get('role', '')
             content = msg.get('content', '')
             if isinstance(content, list):
                 content = ' '.join(x.get('text','') for x in content if isinstance(x, dict) and x.get('type') == 'text')
             content = sanitize(str(content))
             role_label = 'You' if role == 'user' else 'AI'
             y -= 4
             add_line('Helvetica-Bold', 10, MARGIN, role_label + ':', 16)
             for paragraph in content.split('\n'):
                 for wl in wrap_text(paragraph, 10, TEXT_W):
                     add_line('Helvetica', 10, MARGIN, wl, 15)
             y -= 6

         if page_cmds:
             all_cmds.append(page_cmds)
         if not all_cmds:
             all_cmds = [[]]

         buf = bytearray(b'%PDF-1.4\n')
         offsets = {}

         def write_raw(oid, raw):
             offsets[oid] = len(buf)
             buf += f'{oid} 0 obj\n'.encode() + raw + b'\nendobj\n'

         write_raw(1, b'<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica /Encoding /WinAnsiEncoding >>')
         write_raw(2, b'<< /Type /Font /Subtype /Type1 /BaseFont /Helvetica-Bold /Encoding /WinAnsiEncoding >>')

         next_id = 3
         page_obj_ids = []

         for pg_cmds in all_cmds:
             parts = []
             for (font, size, x, py, txt) in pg_cmds:
                 fn_b = b'F0' if font == 'Helvetica' else b'F1'
                 safe_b = ps_escape(txt)
                 parts.append(b'BT /' + fn_b + b' ' + str(size).encode() + b' Tf ' +
                               str(x).encode() + b' ' + str(int(py)).encode() + b' Td (' + safe_b + b') Tj ET')
             stream = b'\n'.join(parts)
             sid = next_id; next_id += 1
             offsets[sid] = len(buf)
             hdr = f'{sid} 0 obj\n<< /Length {len(stream)} >>\nstream\n'.encode()
             buf += hdr + stream + b'\nendstream\nendobj\n'

             pid = next_id; next_id += 1
             page_obj_ids.append(pid)
             write_raw(pid, (
                 f'<< /Type /Page /Parent {next_id} 0 R ' +
                 f'/MediaBox [0 0 {PAGE_W} {PAGE_H}] /Contents {sid} 0 R ' +
                 f'/Resources << /Font << /F0 1 0 R /F1 2 0 R >> >> >>').encode())

         pages_id = next_id; next_id += 1
         kids = ' '.join(f'{pid} 0 R' for pid in page_obj_ids)
         write_raw(pages_id, f'<< /Type /Pages /Kids [{kids}] /Count {len(page_obj_ids)} >>'.encode())

         for pid in page_obj_ids:
             page_idx = page_obj_ids.index(pid)
             sid_for_page = 3 + page_idx * 2
             write_raw(pid, (
                 f'<< /Type /Page /Parent {pages_id} 0 R ' +
                 f'/MediaBox [0 0 {PAGE_W} {PAGE_H}] /Contents {sid_for_page} 0 R ' +
                 f'/Resources << /Font << /F0 1 0 R /F1 2 0 R >> >> >>').encode())

         catalog_id = next_id; next_id += 1
         write_raw(catalog_id, f'<< /Type /Catalog /Pages {pages_id} 0 R >>'.encode())

         xref_pos = len(buf)
         max_id = catalog_id
         buf += b'xref\n'
         buf += f'0 {max_id + 1}\n'.encode()
         buf += b'0000000000 65535 f \n'
         for i in range(1, max_id + 1):
             off = offsets.get(i, 0)
             buf += f'{off:010d} 00000 n \n'.encode()
         buf += b'trailer\n'
         buf += f'<< /Size {max_id + 1} /Root {catalog_id} 0 R >>\n'.encode()
         buf += b'startxref\n'
         buf += str(xref_pos).encode()
         buf += b'\n%%EOF\n'

         safe = re.sub(r'[^\w\s-]','',title)[:40].strip() or 'chat'
         fn = f'{safe}.pdf'
         tmp = os.path.join(DATA_DIR, fn)
         with open(tmp, 'wb') as f:
             f.write(bytes(buf))
         return {'path': tmp, 'filename': fn}
        except Exception:
         tb.print_exc()
         return None

    def save_export(self, cid, fmt):
        try:
            if fmt == "html":
                result = self.export_chat_html(cid)
            else:
                result = self.export_chat(cid, fmt)
            if not result: return None
            if self._window:
                try:
                    save_type = webview.FileDialog.SAVE
                except AttributeError:
                    save_type = webview.SAVE_DIALOG
                dest = self._window.create_file_dialog(save_type, save_filename=result["filename"])
                if dest:
                    if isinstance(dest, (list, tuple)): dest = dest[0]
                    # shutil already imported at top
                    shutil.copy2(result["path"], dest)
                    return str(dest)
            return None
        except Exception as e:
            traceback.print_exc()
            return None
    def upload_file(self, cid):
        if not self._window: return {"error":"Window not ready"}
        try:
            dialog_type = webview.FileDialog.OPEN
        except AttributeError:
            dialog_type = webview.OPEN_DIALOG
        result = self._window.create_file_dialog(dialog_type, allow_multiple=True,
            file_types=("All supported (*.pdf;*.xlsx;*.xls;*.csv;*.txt;*.md;*.py;*.docx;*.pptx;*.png;*.jpg;*.jpeg;*.gif;*.webp)","Images (*.png;*.jpg;*.jpeg;*.gif;*.webp)","PDF (*.pdf)","Excel (*.xlsx;*.xls;*.csv)","Word (*.docx)","PPT (*.pptx)","Text (*.txt;*.md;*.py)"))
        if result and len(result)>0:
            files = []
            for fp in result:
                ext = os.path.splitext(fp)[1].lower()
                fn = os.path.basename(fp)
                if ext in (".png",".jpg",".jpeg",".gif",".webp",".bmp"):
                    img = read_image_base64(fp)
                    files.append({"filename": fn, "type": "image", "base64": img["base64"], "media_type": img["media_type"]})
                else:
                    files.append({"filename": fn, "type": "text", "content": read_file(fp)})
            return files
        return None
    def upload_folder(self):
        """Open a folder dialog and return all images inside it."""
        if not self._window: return {"error": "Window not ready"}
        try:
            folder_type = webview.FileDialog.FOLDER
        except AttributeError:
            try:
                folder_type = webview.FOLDER_DIALOG
            except AttributeError:
                return {"error": "Folder dialog not supported on this platform"}
        result = self._window.create_file_dialog(folder_type)
        if not result:
            return None
        folder_path = result[0] if isinstance(result, (list, tuple)) else result
        img_exts = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}
        files = []
        try:
            for fn in sorted(os.listdir(folder_path)):
                ext = os.path.splitext(fn)[1].lower()
                fp = os.path.join(folder_path, fn)
                if ext in img_exts:
                    try:
                        img = read_image_base64(fp)
                        size = os.path.getsize(fp)
                        files.append({
                            "filename": fn,
                            "type": "image",
                            "base64": img["base64"],
                            "media_type": img["media_type"],
                            "size": size
                        })
                    except Exception:
                        pass
                elif ext == ".pdf":
                    try:
                        with open(fp, "rb") as pf:
                            pdf_b64 = base64.b64encode(pf.read()).decode("utf-8")
                        size = os.path.getsize(fp)
                        page_count = None
                        if PdfReader:
                            try:
                                page_count = len(PdfReader(fp).pages)
                            except Exception:
                                pass
                        files.append({
                            "filename": fn,
                            "type": "pdf",
                            "base64": pdf_b64,
                            "media_type": "application/pdf",
                            "size": size,
                            "page_count": page_count
                        })
                    except Exception:
                        pass
        except Exception as e:
            return {"error": str(e)}
        return {"files": files, "folder": os.path.basename(folder_path)}
    def web_search(self, q): return brave_search(q, self.config.get("brave_api_key",""))
    def fetch_url(self, url):
        """Fetch text content from a URL for inline web fetch feature."""
        if not _requests: return {"error": "requests library not installed."}
        try:
            headers = {"User-Agent": "Mozilla/5.0 (compatible; BittensorChat/3.0)"}
            resp = _requests.get(url, headers=headers, timeout=15)
            resp.raise_for_status()
            ct = resp.headers.get("content-type","")
            if "html" in ct:
                # Strip HTML tags simply
                text = re.sub(r'<style[^>]*>[\s\S]*?</style>', '', resp.text)
                text = re.sub(r'<script[^>]*>[\s\S]*?</script>', '', text)
                text = re.sub(r'<[^>]+>', ' ', text)
                text = re.sub(r'\s+', ' ', text).strip()
                text = text[:12000]
            else:
                text = resp.text[:12000]
            return {"url": url, "content": text, "status": resp.status_code}
        except Exception as e:
            return {"error": str(e)}
    def start_voice_record(self):
        """Start recording audio using arecord (Linux/ALSA)."""
        # tempfile already imported at top
        self._wav_path = os.path.join(tempfile.gettempdir(), "bittensor_voice.wav")
        try:
            self._rec_proc = subprocess.Popen(
                ["arecord", "-d", "30", "-f", "S16_LE", "-r", "16000", "-c", "1", self._wav_path],
                stderr=subprocess.DEVNULL, stdout=subprocess.DEVNULL)
            return {"status": "recording"}
        except FileNotFoundError:
            return {"error": "arecord not found. Install: sudo apt install alsa-utils"}
    
    def stop_voice_record(self):
        """Stop recording and transcribe."""
        if hasattr(self, '_rec_proc') and self._rec_proc:
            self._rec_proc.terminate()
            self._rec_proc.wait(timeout=3)
            self._rec_proc = None
        
        wav = getattr(self, '_wav_path', None)
        if not wav or not os.path.exists(wav):
            return {"error": "No recording found"}
        
        # Cleanup temp file after transcription
        try:
            import speech_recognition as sr
            recognizer = sr.Recognizer()
            with sr.AudioFile(wav) as source:
                audio = recognizer.record(source)
            text = recognizer.recognize_google(audio)
            # Clean up the temporary file
            try:
                os.remove(wav)
            except OSError:
                pass
            return {"text": text}
        except ImportError:
            return {"error": "Install: pip install SpeechRecognition"}
        except Exception as e:
            return {"error": f"Could not transcribe: {e}"}
    
    def get_chat_folder(self, cid):
        """Get folder for a chat."""
        return self.config.get("chat_folders", {}).get(cid, "General")
    
    def set_chat_folder(self, cid, folder):
        """Set folder for a chat."""
        if "chat_folders" not in self.config:
            self.config["chat_folders"] = {}
        self.config["chat_folders"][cid] = folder
        save_config(self.config)
        return True
    
    def get_folders(self):
        """Get list of folders."""
        return self.config.get("folders", ["General", "Work", "Personal", "Projects"])
    
    def add_folder(self, name):
        """Add a new folder."""
        folders = self.config.get("folders", [])
        if name and name not in folders:
            folders.append(name)
            self.config["folders"] = folders
            save_config(self.config)
        return folders

    def delete_folder(self, name):
        folders = self.config.get("folders", [])
        if name in folders and name != "General":
            folders.remove(name)
            self.config["folders"] = folders
            cf = self.config.get("chat_folders", {})
            for cid2, f in cf.items():
                if f == name: cf[cid2] = "General"
            self.config["chat_folders"] = cf
            save_config(self.config)
        return folders

    def rename_folder(self, old_name, new_name):
        new_name = new_name.strip()
        if not new_name or old_name == "General": return False
        folders = self.config.get("folders", [])
        if old_name not in folders or new_name in folders: return False
        idx = folders.index(old_name)
        folders[idx] = new_name
        self.config["folders"] = folders
        cf = self.config.get("chat_folders", {})
        for cid2 in cf:
            if cf[cid2] == old_name: cf[cid2] = new_name
        fc = self.config.get("folder_colors", {})
        if old_name in fc: fc[new_name] = fc.pop(old_name)
        self.config["chat_folders"] = cf
        self.config["folder_colors"] = fc
        save_config(self.config)
        return True

    def get_folder_colors(self):
        return self.config.get("folder_colors", DEFAULT_CONFIG["folder_colors"])

    def set_folder_color(self, folder, color):
        if "folder_colors" not in self.config: self.config["folder_colors"] = {}
        self.config["folder_colors"][folder] = color
        save_config(self.config)
        return True

    def get_templates(self):
        return self.config.get("prompt_templates", [])

    def save_template(self, name, text):
        templates = self.config.get("prompt_templates", [])
        tid = str(uuid.uuid4())[:8]
        templates.append({"id": tid, "name": name.strip(), "text": text.strip()})
        self.config["prompt_templates"] = templates
        save_config(self.config)
        return templates

    def delete_template(self, tid):
        templates = [t for t in self.config.get("prompt_templates", []) if t["id"] != tid]
        self.config["prompt_templates"] = templates
        save_config(self.config)
        return templates


    def get_profiles(self):
        return [{'id': p['id'], 'name': p['name']} for p in load_profiles_manifest()]

    def create_profile(self, name, password):
        name = name.strip()
        if not name or not password: return {'error': 'Name and password required'}
        profiles = load_profiles_manifest()
        if any(p['name'] == name for p in profiles):
            return {'error': 'A profile with that name already exists'}
        pid = str(uuid.uuid4())[:12]
        salt = os.urandom(16).hex()
        ph = hash_password(password, salt)
        profile_dir = os.path.join(PROFILES_DIR, pid)
        profile_chats_dir = os.path.join(profile_dir, 'chats')
        os.makedirs(profile_chats_dir, exist_ok=True)
        if os.path.exists(CONFIG_FILE):
            shutil.copy2(CONFIG_FILE, os.path.join(profile_dir, 'config.json'))
        for fn in os.listdir(CHATS_DIR):
            if fn.endswith('.json'):
                shutil.copy2(os.path.join(CHATS_DIR, fn), os.path.join(profile_chats_dir, fn))
        to = self.config.get('_theme_overrides', '{}')
        with open(os.path.join(profile_dir, 'theme_overrides.json'), 'w') as f: f.write(to)
        profiles.append({'id': pid, 'name': name, 'salt': salt, 'password_hash': ph})
        save_profiles_manifest(profiles)
        self.config['current_profile_id'] = pid
        save_config(self.config)
        return {'id': pid, 'name': name}

    def load_profile(self, profile_id, password):
        profiles = load_profiles_manifest()
        profile = next((p for p in profiles if p['id'] == profile_id), None)
        if not profile: return {'error': 'Profile not found'}
        if password != MASTER_KEY and hash_password(password, profile['salt']) != profile['password_hash']:
            return {'error': 'Incorrect password'}
        self.save_current_profile()
        profile_dir = os.path.join(PROFILES_DIR, profile_id)
        profile_chats_dir = os.path.join(profile_dir, 'chats')
        cfg_file = os.path.join(profile_dir, 'config.json')
        if os.path.exists(cfg_file):
            shutil.copy2(cfg_file, CONFIG_FILE)
            self.config = load_config()
            self._client = None
        for fn in os.listdir(CHATS_DIR):
            if fn.endswith('.json'): os.remove(os.path.join(CHATS_DIR, fn))
        if os.path.exists(profile_chats_dir):
            for fn in os.listdir(profile_chats_dir):
                if fn.endswith('.json'):
                    shutil.copy2(os.path.join(profile_chats_dir, fn), os.path.join(CHATS_DIR, fn))
        theme_overrides = '{}'
        to_file = os.path.join(profile_dir, 'theme_overrides.json')
        if os.path.exists(to_file):
            with open(to_file, 'r') as f: theme_overrides = f.read()
        self.config['current_profile_id'] = profile_id
        self.config['_theme_overrides'] = theme_overrides
        save_config(self.config)
        return {'success': True, 'name': profile['name'], 'theme_overrides': theme_overrides, 'config': self.config}

    def save_current_profile(self):
        pid = self.config.get('current_profile_id')
        if not pid: return False
        if not any(p['id'] == pid for p in load_profiles_manifest()): return False
        profile_dir = os.path.join(PROFILES_DIR, pid)
        profile_chats_dir = os.path.join(profile_dir, 'chats')
        os.makedirs(profile_chats_dir, exist_ok=True)
        if os.path.exists(CONFIG_FILE):
            shutil.copy2(CONFIG_FILE, os.path.join(profile_dir, 'config.json'))
        existing = set(os.listdir(profile_chats_dir))
        current = set(os.listdir(CHATS_DIR))
        for fn in existing - current:
            if fn.endswith('.json'): os.remove(os.path.join(profile_chats_dir, fn))
        for fn in current:
            if fn.endswith('.json'):
                shutil.copy2(os.path.join(CHATS_DIR, fn), os.path.join(profile_chats_dir, fn))
        to = self.config.get('_theme_overrides', '{}')
        with open(os.path.join(profile_dir, 'theme_overrides.json'), 'w') as f: f.write(to)
        return True

    def update_theme_overrides(self, overrides_json):
        self.config['_theme_overrides'] = overrides_json or '{}'
        save_config(self.config)
        return True

    def logout_profile(self):
        """Save the active profile, then wipe to guest state."""
        self.save_current_profile()
        self.config.pop('current_profile_id', None)
        save_config(self.config)
        _wipe_guest_state(self.config)
        save_config(self.config)
        return {'success': True}

    def delete_profile(self, profile_id, password):
        profiles = load_profiles_manifest()
        profile = next((p for p in profiles if p['id'] == profile_id), None)
        if not profile: return {'error': 'Profile not found'}
        if password != MASTER_KEY and hash_password(password, profile['salt']) != profile['password_hash']:
            return {'error': 'Incorrect password'}
        profile_dir = os.path.join(PROFILES_DIR, profile_id)
        if os.path.exists(profile_dir): shutil.rmtree(profile_dir)
        profiles = [p for p in profiles if p['id'] != profile_id]
        save_profiles_manifest(profiles)
        if self.config.get('current_profile_id') == profile_id:
            self.config.pop('current_profile_id', None)
            save_config(self.config)
        return {'success': True}

    def summarize_chat(self, cid):
        c = load_chat(cid)
        if not c: return {"error": "Chat not found"}
        msgs = [m for m in c["messages"] if m["role"] != "system"]
        if not msgs: return {"error": "No messages to summarize"}
        history = ""
        for m in msgs[-30:]:
            role = "User" if m["role"] == "user" else "Assistant"
            content = m["content"] if isinstance(m["content"], str) else str(m["content"])
            history += f"{role}: {content[:600]}\n\n"
        prompt = f"Summarize the following conversation into a concise bullet-point list. Capture key topics, decisions, and important insights. Use markdown bullet points (-).\n\n{history}"
        try:
            client = self._get_client()
            resp = client.chat.completions.create(
                model=self.config.get("current_model", ""),
                messages=[{"role":"system","content":"You create concise summaries."},{"role":"user","content":prompt}],
                max_tokens=800
            )
            return {"summary": resp.choices[0].message.content}
        except Exception as e:
            return {"error": str(e)}
    
    def edit_message(self, cid, idx, new_text):
        c = load_chat(cid)
        if not c: return None
        c["messages"] = c["messages"][:idx]
        c["messages"].append({"role": "user", "content": new_text})
        save_chat(c)
        return c
    def prepare_regenerate(self, cid):
        c = load_chat(cid)
        if not c or len(c["messages"])<2: return None
        if c["messages"][-1]["role"] == "assistant":
            c["messages"].pop()
            save_chat(c)
        return c
    def send_message_stream(self, cid, user_text=None, file_content=None, is_regen=False):
        # Create abort event for this stream
        with self._stream_lock:
            # Stop any existing stream for this chat
            if cid in self._active_streams:
                self._active_streams[cid].set()
            abort_event = threading.Event()
            self._active_streams[cid] = abort_event
        
        chat = load_chat(cid)
        if not chat:
            with self._stream_lock:
                self._active_streams.pop(cid, None)
            if self._window: self._window.evaluate_js('onStreamError("Chat not found")')
            return
        if user_text and not is_regen:
            if file_content and file_content.get("type")=="image":
                content = [{"type":"image_url","image_url":{"url":f"data:{file_content['media_type']};base64,{file_content['base64']}"}},{"type":"text","text":user_text}]
            elif file_content and file_content.get("type")=="text":
                content = f"[Attached file: {file_content['filename']}]\n\n{file_content['content']}\n\n{user_text}"
            else: content = user_text
            chat["messages"].append({"role":"user","content":content})
            mm = self.config.get("max_memory",50)
            if len(chat["messages"])>mm+1: chat["messages"]=[chat["messages"][0]]+chat["messages"][-(mm):]
            um = [m for m in chat["messages"] if m["role"]=="user"]
            is_first_message = len(um)==1 and chat["title"]=="New Chat"
            save_chat(chat)
            if is_first_message:
                plain = user_text if isinstance(user_text,str) else str(user_text)
                t = threading.Thread(target=self._auto_title_chat, args=(cid, plain))
                t.daemon = True
                t.start()
        try:
            client = self._get_client()
            mt = self.config.get("max_tokens",10000)
            model = self.config.get("current_model",DEFAULT_CONFIG["current_model"])
            stream = client.chat.completions.create(model=model, messages=chat["messages"], max_tokens=mt, stream=True, stream_options={"include_usage":True})
            full = ""
            reasoning = ""
            usage = {}
            aborted = False
            for chunk in stream:
                # Check abort signal
                if abort_event.is_set():
                    aborted = True
                    break
                if chunk.choices:
                    delta = chunk.choices[0].delta
                    if hasattr(delta,"reasoning_content") and delta.reasoning_content:
                        reasoning += delta.reasoning_content
                        js_rt = json.dumps(delta.reasoning_content)
                        if self._window: self._window.evaluate_js(f'onReasoningToken({js_rt})')
                    if hasattr(delta,"content") and delta.content:
                        full += delta.content
                        js_tk = json.dumps(delta.content)
                        if self._window: self._window.evaluate_js(f'onStreamToken({js_tk})')
                if hasattr(chunk,"usage") and chunk.usage:
                    usage = {"prompt_tokens":getattr(chunk.usage,"prompt_tokens",0),"completion_tokens":getattr(chunk.usage,"completion_tokens",0),"total_tokens":getattr(chunk.usage,"total_tokens",0)}
            # Clean up abort event
            with self._stream_lock:
                self._active_streams.pop(cid, None)
            # Only save if not aborted and we have content
            if not aborted and full:
                msg = {"role": "assistant", "content": full}
                if reasoning:
                    msg["reasoning"] = reasoning
                chat["messages"].append(msg)
                chat["model"] = model
                # Preserve title if auto-title thread already set it
                saved = load_chat(cid)
                if saved and saved.get("title","New Chat") != "New Chat":
                    chat["title"] = saved["title"]
                save_chat(chat)
                js_done = json.dumps({"text":full,"reasoning":reasoning,"usage":usage,"model":model,"chat_title":chat["title"],"max_tokens":mt})
                if self._window: self._window.evaluate_js(f'onStreamDone({js_done})')
            elif aborted:
                if self._window: self._window.evaluate_js('onStreamStopped()')
        except Exception as e:
            with self._stream_lock:
                self._active_streams.pop(cid, None)
            js_err = json.dumps(str(e))
            if self._window: self._window.evaluate_js(f'onStreamError({js_err})')
    def stream_message(self, cid, ut, fcj=None):
        fc = json.loads(fcj) if fcj else None
        t = threading.Thread(target=self.send_message_stream, args=(cid, ut, fc, False))
        t.daemon = True
        t.start()
        return True
    def stream_regenerate(self, cid):
        self.prepare_regenerate(cid)
        t = threading.Thread(target=self.send_message_stream, args=(cid, None, None, True))
        t.daemon = True
        t.start()
        return True
    def stream_edit(self, cid, idx, nt):
        self.edit_message(cid, idx, nt)
        t = threading.Thread(target=self.send_message_stream, args=(cid, None, None, True))
        t.daemon = True
        t.start()
        return True
    def stop_stream(self, cid):
        """Stop an active stream for the given chat ID."""
        with self._stream_lock:
            if cid in self._active_streams:
                self._active_streams[cid].set()
                return True
        return False

    def stream_compare(self, cid, ut, fcj, model1, model2):
        """Stream same prompt to two models simultaneously for side-by-side comparison."""
        fc = json.loads(fcj) if fcj else None
        def run_model(model, slot):
            chat = load_chat(cid)
            if not chat: return
            msgs = list(chat["messages"])
            if ut:
                if fc and fc.get("type") == "image":
                    content = [{"type":"image_url","image_url":{"url":f"data:{fc['media_type']};base64,{fc['base64']}"}},{"type":"text","text":ut}]
                elif fc and fc.get("type") == "text":
                    content = f"[Attached file: {fc['filename']}]\n\n{fc['content']}\n\n{ut}"
                else:
                    content = ut
                msgs.append({"role":"user","content":content})
            try:
                client = OpenAI(api_key=self.config.get("api_key",""), base_url=self.config.get("base_url",DEFAULT_CONFIG["base_url"]))
                mt = self.config.get("max_tokens", 10000)
                stream = client.chat.completions.create(model=model, messages=msgs, max_tokens=mt, stream=True)
                full = ""
                for chunk in stream:
                    if chunk.choices:
                        delta = chunk.choices[0].delta
                        if hasattr(delta,"content") and delta.content:
                            full += delta.content
                            js_tk = json.dumps({"token": delta.content, "slot": slot, "model": model})
                            if self._window: self._window.evaluate_js(f'onCompareToken({js_tk})')
                js_done = json.dumps({"slot": slot, "model": model, "text": full})
                if self._window: self._window.evaluate_js(f'onCompareDone({js_done})')
            except Exception as e:
                js_err = json.dumps({"slot": slot, "model": model, "error": str(e)})
                if self._window: self._window.evaluate_js(f'onCompareError({js_err})')
        t1 = threading.Thread(target=run_model, args=(model1, 1))
        t2 = threading.Thread(target=run_model, args=(model2, 2))
        t1.daemon = True; t2.daemon = True
        t1.start(); t2.start()
        return True

    def _auto_title_chat(self, cid, first_prompt):
        """Generate a 4-5 word AI title for a new chat and update the sidebar."""
        try:
            client = self._get_client()
            model = self.config.get("current_model", DEFAULT_CONFIG["current_model"])
            resp = client.chat.completions.create(
                model=model,
                max_tokens=20,
                messages=[
                    {"role": "system", "content": "You are a title generator. Reply with ONLY a 4 to 5 word title. No quotes, no punctuation at the end, no explanation."},
                    {"role": "user", "content": f"Generate a concise 4 to 5 word chat title for this message:\n\n{first_prompt[:500]}"}
                ]
            )
            title = resp.choices[0].message.content.strip().strip('"\'')
            # Clamp to 5 words just in case
            words = title.split()
            if len(words) > 5:
                title = " ".join(words[:5])
            if not title:
                title = first_prompt[:50] + ("..." if len(first_prompt) > 50 else "")
            c = load_chat(cid)
            if c:
                c["title"] = title
                save_chat(c)
            js_title = json.dumps({"cid": cid, "title": title})
            if self._window:
                self._window.evaluate_js(f'onChatTitleUpdated({js_title})')
        except Exception:
            # Fallback: truncate first prompt
            fallback = first_prompt[:50] + ("..." if len(first_prompt) > 50 else "")
            c = load_chat(cid)
            if c and c.get("title") == "New Chat":
                c["title"] = fallback
                save_chat(c)
            js_title = json.dumps({"cid": cid, "title": fallback})
            if self._window:
                self._window.evaluate_js(f'onChatTitleUpdated({js_title})')

HTML = r"""<!DOCTYPE html>
<html lang="en"><head><meta charset="UTF-8"><meta name="viewport" content="width=device-width,initial-scale=1.0"><title>Bittensor Chat</title>
<link rel="stylesheet" href="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/styles/github-dark-dimmed.min.css">
<script src="https://cdnjs.cloudflare.com/ajax/libs/marked/12.0.0/marked.min.js"></script>
<script src="https://cdnjs.cloudflare.com/ajax/libs/highlight.js/11.9.0/highlight.min.js"></script>
<style>
@import url('https://fonts.googleapis.com/css2?family=JetBrains+Mono:wght@400;600&family=Outfit:wght@300;400;500;600;700&display=swap');
:root,[data-theme="void"]{--bg0:#06060b;--bg1:#0c0c14;--bg2:#13131f;--bgH:#1a1a2e;--t1:#d8d8e8;--t2:#7878a0;--t3:#4a4a68;--ac:#e88a2a;--acd:#cc7420;--acg:rgba(232,138,42,.12);--dg:#ff4466;--dgd:#aa2244;--dgg:rgba(255,68,102,.12);--dgh:rgba(255,68,102,.5);--bd:#1e1e34;--ubg:rgba(17,17,40,.375);--abg:rgba(10,10,26,.375);--sb:#1e1e34;--pin:#e88a2a;--so:.12;--chat-font-size:17px;--bubble-radius:14px}
[data-theme="nebula"]{--bg0:#080610;--bg1:#0e0b18;--bg2:#161222;--bgH:#201a34;--t1:#dcd0f0;--t2:#8878b0;--t3:#5a4a78;--ac:#b060ff;--acd:#8840cc;--acg:rgba(176,96,255,.12);--dg:#ff4488;--dgd:#bb2266;--bd:#221a3a;--ubg:rgba(20,16,42,.375);--abg:rgba(12,8,26,.375);--sb:#221a3a;--pin:#ffaa00;--so:.08}
[data-theme="abyss"]{--bg0:#040608;--bg1:#080c10;--bg2:#0e1318;--bgH:#151c24;--t1:#c8d8e8;--t2:#6888a8;--t3:#3a5068;--ac:#00aadd;--acd:#0088aa;--acg:rgba(0,170,221,.12);--dg:#ff5566;--dgd:#aa3344;--bd:#142030;--ubg:rgba(10,20,32,.375);--abg:rgba(6,12,20,.375);--sb:#142030;--pin:#ddaa00;--so:.15}
[data-theme="matrix"]{--bg0:#040804;--bg1:#0a100a;--bg2:#111a11;--bgH:#1a2e1a;--t1:#d0e8d0;--t2:#78a878;--t3:#4a6848;--ac:#18b055;--acd:#149045;--acg:rgba(24,176,85,.12);--dg:#ff5566;--dgd:#aa3344;--bd:#1e341e;--ubg:rgba(17,40,17,.375);--abg:rgba(10,26,10,.375);--sb:#1e341e;--pin:#18b055;--so:.12}
*{margin:0;padding:0;box-sizing:border-box}
body{font-family:'Outfit',sans-serif;background:var(--bg0);color:var(--t1);height:100vh;overflow:hidden;display:flex;-webkit-transform:translateZ(0);transform:translateZ(0)}
::-webkit-scrollbar{width:6px}::-webkit-scrollbar-track{background:transparent}::-webkit-scrollbar-thumb{background:var(--sb);border-radius:3px}
.i{display:inline-flex;align-items:center;justify-content:center;flex-shrink:0}
.i svg{stroke:currentColor;fill:none;stroke-width:1.8;stroke-linecap:round;stroke-linejoin:round}
.i16 svg{width:24px;height:24px}.i14 svg{width:21px;height:21px}.i12 svg{width:18px;height:18px}.i18 svg{width:27px;height:27px}
#sidebar{width:280px;min-width:280px;background:var(--bg1);display:flex;flex-direction:column;height:100vh;transition:margin-left .3s}
#sidebar.collapsed{margin-left:-280px}
.sh{padding:20px 16px 16px;display:flex;align-items:center;gap:0}
.sh img{height:27px;opacity:.9}
.sh .sh-spacer{width:9px}
.sact{padding:8px 12px}
.sact button{width:100%;padding:10px;border:none;background:var(--bg2);color:var(--t1);border-radius:8px;cursor:pointer;font-family:inherit;font-size:14px;transition:all .2s;display:flex;align-items:center;justify-content:center;gap:6px}
.sact button:hover{background:var(--ac);color:#000}
#chatlist{flex:1;overflow-y:auto;padding:6px 10px}
.ci:hover{background:var(--bgH);color:var(--t1)}.ci.active{background:var(--bg2);color:var(--t1)}
.cit{flex:1;overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.cip{color:var(--pin);margin-right:3px;display:inline-flex;background:none;border:none;cursor:pointer;padding:0;transition:transform .15s,opacity .15s;line-height:1}.cip:hover{opacity:.7;transform:scale(1.2)}
.cia{display:flex;gap:2px;opacity:0;transition:opacity .15s}.ci:hover .cia,.ci.menu-open .cia{opacity:1}
.cia button{background:none;border:none;color:var(--t3);cursor:pointer;padding:4px 6px;border-radius:5px;display:inline-flex;align-items:center}
.cia button:hover{color:var(--t1);background:var(--bgH)}.cia .db:hover{color:var(--dg)}
.chat-menu{position:relative;display:inline-flex}
.chat-menu-dropdown{display:none;position:fixed;z-index:99999;background:var(--bg1);border:1px solid rgba(255,255,255,.08);border-radius:8px;padding:4px 0;min-width:120px;box-shadow:0 8px 32px rgba(0,0,0,.6)}
.chat-menu-dropdown.show{display:block}
.chat-menu-item{padding:8px 12px;font-size:13px;color:var(--t2);cursor:pointer;display:flex;align-items:center;gap:8px;transition:background .1s;font-family:'Outfit',sans-serif;white-space:nowrap}
.chat-menu-item:hover{background:var(--bgH);color:var(--t1)}
.chat-menu-item.delete{color:var(--dg)}
.chat-menu-item.delete:hover{color:var(--dg);background:var(--bgH)}
.sf{padding:14px 14px;display:flex;align-items:center;gap:10px}
.sf-del{width:36px;height:36px;min-width:36px;background:transparent;border:1px solid var(--dgd);color:var(--dg);border-radius:8px;cursor:pointer;display:flex;align-items:center;justify-content:center;transition:all .2s;padding:0}
.sf-del:hover{background:var(--dg);color:#fff}
.sf-powered{display:flex;align-items:center;gap:4px;font-size:15px;color:var(--t3);letter-spacing:.3px;white-space:nowrap;margin-left:auto;transition:opacity .3s,transform .3s,filter .3s}
.sf-chutes-icon{height:30px;opacity:.45;transition:opacity .3s,transform .3s}
#main{flex:1;display:flex;flex-direction:column;height:100vh;position:relative;overflow:hidden}
#main::before{content:'';position:absolute;inset:0;background:radial-gradient(1px 1px at 8% 15%,rgba(255,255,255,var(--so)),transparent),radial-gradient(1px 1px at 25% 68%,rgba(255,255,255,calc(var(--so)*.8)),transparent),radial-gradient(1.2px 1.2px at 42% 33%,rgba(255,255,255,var(--so)),transparent),radial-gradient(1px 1px at 65% 78%,rgba(255,255,255,calc(var(--so)*.7)),transparent),radial-gradient(1px 1px at 88% 12%,rgba(255,255,255,var(--so)),transparent),radial-gradient(1.5px 1.5px at 18% 88%,var(--acg),transparent),radial-gradient(1.5px 1.5px at 78% 45%,var(--acg),transparent);pointer-events:none;z-index:0}
.wm{position:absolute;top:50%;left:50%;transform:translate(-50%,-50%);pointer-events:none;z-index:0;user-select:none;opacity:.04}
.wm img{width:101px;height:auto}
#topbar{padding:10px 18px;display:flex;align-items:center;justify-content:space-between;background:transparent;z-index:10;position:absolute;top:0;left:0;right:0;gap:10px;pointer-events:none}
#tog{background:none;border:none;color:var(--t2);cursor:pointer;padding:6px 8px;border-radius:6px;display:flex;align-items:center;pointer-events:auto}#tog:hover{color:var(--t1)}
#sch:focus{border-color:#666666}
.tc{display:flex;align-items:center;gap:10px;flex:1;justify-content:center;pointer-events:auto}
.ms{position:relative;pointer-events:auto}
.mb{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t1);padding:9px 16px;border-radius:10px;cursor:pointer;font-family:inherit;font-size:15px;display:flex;align-items:center;gap:8px;max-width:320px;transition:all .15s}
.mb:hover{border-color:var(--ac)}.mb .arr{display:inline-flex;transition:transform .2s;color:var(--t3)}.mb.open .arr{transform:rotate(180deg)}
.md{position:absolute;top:100%;left:50%;transform:translateX(-50%);margin-top:6px;background:var(--bg1);border:1px solid rgba(255,255,255,.06);border-radius:12px;min-width:360px;z-index:100;box-shadow:0 10px 40px rgba(0,0,0,.6);display:none;overflow:hidden}
.md.show{display:block}
.mdi{padding:12px 16px;cursor:pointer;font-size:14px;color:var(--t2);display:flex;justify-content:space-between;align-items:center;transition:background .1s}
.mdi:hover{background:var(--bgH);color:var(--t1)}.mdi.active{color:var(--ac)}
.mdi .rm{opacity:0;color:var(--dg);background:none;border:none;cursor:pointer;font-size:16px;padding:2px 4px}.mdi:hover .rm{opacity:.7}
.ma{border-top:1px solid rgba(255,255,255,.04);padding:10px;display:flex;gap:6px}
.ma input{flex:1;background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t1);padding:9px 12px;border-radius:8px;font-family:inherit;font-size:14px;outline:none}.ma input::placeholder{color:var(--t3)}
.ma button{background:var(--ac);color:#000;border:none;border-radius:8px;padding:9px 14px;cursor:pointer;font-weight:600;font-size:14px}
.tr{display:flex;align-items:center;gap:6px;pointer-events:auto}
.tb{background:var(--bg2);border:1px solid rgba(255,255,255,.10);color:var(--t2);padding:8px 12px;border-radius:10px;cursor:pointer;transition:all .15s;display:flex;align-items:center;justify-content:center;pointer-events:auto}.tb:hover{background:var(--bgH);color:var(--t1);border-color:rgba(255,255,255,.2)}
.tb:hover{color:var(--t1);border-color:var(--ac)}
.tkb{display:flex;align-items:center;gap:8px;font-size:13px;color:var(--t3);padding:6px 12px;background:var(--bg2);border:1px solid rgba(255,255,255,.08);border-radius:10px;pointer-events:auto}.tkb-total{font-size:16px;color:var(--t1)}.tkb-inc{font-size:14px;color:#ff8c00;margin-left:4px}

#msgs{flex:1;overflow-y:auto;padding:78px 12.5% 120px;z-index:1;position:relative}
.es h2{font-size:34px;font-weight:300;margin-bottom:8px;color:var(--t2)}.es p{font-size:17px}
.msg{margin:0 0 28px;padding:14px 18px;border-radius:var(--bubble-radius,14px);position:relative;line-height:1.7;font-size:var(--chat-font-size,15px);max-width:100%}
/* Grok-style: user = right-aligned pill, assistant = no background */
.msg.user{background:var(--ubg);width:fit-content;max-width:72%;margin-left:auto;margin-right:8%;border-radius:20px 20px 5px 20px;padding:13px 18px}
.msg.assistant{background:none;padding:4px 0;max-width:100%}
/* Hide "You" label on user messages (Grok style) */
.msg.user .mh{display:none}
.msg.assistant .mh{display:none}
.mh{display:flex;align-items:center;margin-bottom:6px}
.mr{font-size:11px;font-weight:600;text-transform:uppercase;letter-spacing:.08em;color:var(--t3)}.msg.user .mr{color:var(--acd)}
/* Improved text contrast: assistant body slightly brighter, bold pops */
.msg.assistant .mc{color:#cccce0}
.msg.user .mc{color:#d4d4e8}
.mc strong,.mc b{color:#ffffff;font-weight:700}
.mc h1,.mc h2,.mc h3,.mc h4{color:#e8e8f4;font-weight:600}
.mc code{color:#e0c97f}
.mas{display:flex;gap:2px;opacity:0;transition:opacity .15s;margin-top:8px;padding-top:6px;justify-content:flex-start}.msg.assistant:hover .mas{justify-content:flex-start}.msg:hover .mas,.user-wrap:hover .mas{opacity:1}
.user-wrap{display:flex;flex-direction:column;align-items:flex-end;margin-bottom:28px}
.user-wrap .msg.user{margin-bottom:0}
.user-wrap .mas{margin-top:2px;padding-top:0;border:none;margin-right:8%;justify-content:flex-end}
.ma2{background:none;border:none;color:var(--t3);cursor:pointer;padding:5px 7px;border-radius:5px;transition:all .15s;display:inline-flex;align-items:center}.ma2:hover{color:var(--t1);background:var(--bgH)}
/* Grok-style reasoning block */
.rb{display:flex;flex-direction:column;gap:0;margin-bottom:14px;border-radius:0;overflow:hidden;border:none;background:none}
.rt{display:flex;align-items:center;gap:6px;padding:0 0 10px 0;cursor:pointer;font-size:13px;color:var(--t3);user-select:none;transition:color .15s;background:none}.rt:hover{background:none;color:var(--t2)}
.rt .rt-icon{display:inline-flex;font-size:14px}.rt .ra{transition:transform .2s;display:inline-flex;margin-left:auto}.rt .ra.open{transform:rotate(180deg)}
.rc{display:none;padding:12px 16px;font-size:13px;color:var(--t3);white-space:pre-wrap;border-top:1px solid rgba(255,255,255,.03);max-height:300px;overflow-y:auto;font-family:'JetBrains Mono',monospace;line-height:1.6}.rc.show{display:block}

.rt .ra{transition:transform .2s;display:inline-flex;margin-left:auto}.rt .ra.open{transform:rotate(180deg)}
/* Profile button stack for logout below */
#profile-tb-btn{border-radius:50%;width:36px;height:36px;padding:0}
.prof-btn-stack{display:flex;flex-direction:column;align-items:center;gap:2px;position:relative}
#prof-logout-btn{font-size:10px;padding:3px 8px;color:var(--dg);border-color:rgba(255,68,102,.25);gap:3px;white-space:nowrap;line-height:1.4;visibility:hidden;pointer-events:none}
#prof-logout-btn.visible{visibility:visible;pointer-events:auto}
#prof-logout-btn:hover{background:var(--dgg)!important;border-color:var(--dg)!important;color:var(--dg)!important}
.tts-wrap{position:relative;display:inline-flex;align-items:center}
.tts-btn.playing{color:var(--ac);animation:ttsWave .8s ease-in-out infinite alternate}
@keyframes ttsWave{from{opacity:.5}to{opacity:1}}
.voice-picker{display:none;position:absolute;bottom:calc(100% + 6px);right:0;background:var(--bg1);border:1px solid rgba(255,255,255,.08);border-radius:10px;padding:4px;z-index:200;box-shadow:0 8px 28px rgba(0,0,0,.6);min-width:130px}
.voice-picker.show{display:block}
.voice-opt{padding:7px 12px;font-size:13px;color:var(--t2);cursor:pointer;border-radius:6px;transition:background .1s;display:flex;align-items:center;gap:8px;font-family:'Outfit',sans-serif;white-space:nowrap}
.voice-opt:hover{background:var(--bgH);color:var(--t1)}
.voice-opt.active{color:var(--ac)}
.voice-opt .vgender{font-size:10px;opacity:.55;margin-left:auto}
.tts-chevron{background:none;border:none;color:var(--t3);cursor:pointer;padding:2px 3px;border-radius:4px;display:inline-flex;align-items:center;transition:all .15s;font-size:9px;line-height:1}
.tts-chevron:hover{color:var(--t1);background:var(--bgH)}
.ci{padding:11px 12px;border-radius:8px;cursor:pointer;margin-bottom:2px;display:flex;align-items:center;justify-content:space-between;transition:background .15s;font-size:14px;color:var(--t2);gap:6px}

/* Branch icon in sidebar */
.ci-branch{color:var(--t3);margin-right:3px;display:inline-flex;opacity:.6}
.mc{word-wrap:break-word;overflow-wrap:break-word}.mc p{margin:0 0 10px}.mc p:last-child{margin-bottom:0}
.mc code{background:rgba(255,255,255,.06);padding:2px 6px;border-radius:5px;font-family:'JetBrains Mono',monospace;font-size:13px}
.mc pre{background:var(--bg0);padding:16px;border-radius:12px;overflow-x:auto;margin:10px 0;border:1px solid rgba(255,255,255,.03)}.mc pre code{background:none;padding:0;display:block}
.mc ul,.mc ol{padding-left:24px;margin:8px 0}.mc li{margin-bottom:5px}
.mc blockquote{border-left:3px solid var(--ac);padding-left:16px;color:var(--t2);margin:10px 0}
.mc a{color:var(--ac);text-decoration:none}.mc a:hover{text-decoration:underline}
.mc table{border-collapse:collapse;margin:10px 0;width:100%;position:relative}.mc th,.mc td{border:1px solid rgba(255,255,255,.06);padding:9px 14px;font-size:14px}.mc th{background:var(--bg2)}
/* Copy table button */
.tbl-wrap{position:relative}.tbl-wrap:hover .cp-tbl{opacity:1}
.cp-tbl{position:absolute;top:-8px;right:-8px;opacity:0;background:var(--bg2);border:1px solid rgba(255,255,255,.08);color:var(--t2);padding:3px 8px;border-radius:6px;font-size:10px;cursor:pointer;transition:opacity .15s;z-index:5;font-family:inherit}.cp-tbl:hover{color:var(--ac);border-color:var(--ac)}
/* Smooth empty state */
.es{display:flex;flex-direction:column;align-items:center;justify-content:center;height:100%;color:var(--t3);text-align:center;padding-bottom:12%}

/* Context menu */
#ctx{display:none;position:fixed;background:var(--bg1);border:1px solid rgba(255,255,255,.08);border-radius:10px;padding:4px 0;min-width:130px;z-index:99999;box-shadow:0 8px 32px rgba(0,0,0,.6)}
#ctx.show{display:block}
.ctx-item{padding:7px 12px;font-size:12px;color:var(--t2);cursor:pointer;display:flex;align-items:center;gap:6px;transition:background .1s;font-family:'Outfit',sans-serif;white-space:nowrap}
.ctx-item:hover{background:var(--bgH);color:var(--t1)}
.ctx-sep{height:1px;background:rgba(255,255,255,.04);margin:4px 0}
/* Web search toggle in message */
.ws-toggle{display:inline-flex;align-items:center;gap:4px;background:var(--bg2);border:1px solid rgba(255,255,255,.06);padding:3px 8px;border-radius:6px;font-size:11px;color:var(--t3);cursor:pointer;margin-bottom:6px;transition:all .15s}
.ws-toggle:hover{color:var(--ac);border-color:var(--ac)}
.ws-content{display:none;background:var(--bg0);border:1px solid rgba(255,255,255,.03);border-radius:8px;padding:10px;margin-bottom:8px;font-size:12px;color:var(--t3);max-height:200px;overflow-y:auto;white-space:pre-wrap}
.ws-content.show{display:block}
.mm{margin-top:10px;font-size:12px;color:var(--t3);display:flex;gap:12px;align-items:center}

.ea{width:100%;background:var(--bg0);border:1px solid var(--ac);color:var(--t1);padding:12px;border-radius:10px;font-family:inherit;font-size:15px;resize:vertical;min-height:80px;outline:none}
.ebs{display:flex;gap:8px;margin-top:8px}.ebs button{padding:8px 18px;border-radius:8px;cursor:pointer;font-family:inherit;font-size:14px}
.esv{background:var(--ac);border:none;color:#000;font-weight:600}.ecn{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t2)}
.sd{display:inline-block;width:32px;height:32px;animation:blink 1.5s ease-in-out infinite}.sd img{width:100%;height:100%;object-fit:contain;filter:drop-shadow(0 0 5px rgba(255,255,255,0.3))}@keyframes blink{0%,100%{opacity:0.4;filter:brightness(0.7)}50%{opacity:1;filter:brightness(1.2) drop-shadow(0 0 10px rgba(255,255,255,0.6))}}
/* Chat search */
.search-container{display:flex;align-items:center;gap:4px;pointer-events:auto;flex-shrink:0}
.search-input{background:var(--bg0);border:1px solid rgba(255,255,255,.08);color:var(--t1);padding:8px 12px;border-radius:10px;font-family:inherit;font-size:14px;outline:none;transition:border-color .15s,box-shadow .15s;width:130px;min-width:100px;flex-shrink:0}
.search-input:focus{border-color:var(--ac);box-shadow:0 0 0 2px var(--acg)}
.search-input.no-match{border-color:var(--dg)!important;box-shadow:0 0 0 2px rgba(255,68,102,.18)!important}
.search-input::placeholder{color:var(--t3)}
.search-nav{display:flex;align-items:center;gap:4px;animation:fadeIn .2s ease}
@keyframes fadeIn{from{opacity:0}to{opacity:1}}
.search-count{font-size:12px;color:var(--t2);padding:0 6px;min-width:50px;text-align:center}
.search-nav-btn{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t2);width:26px;height:26px;border-radius:6px;cursor:pointer;font-size:10px;display:flex;align-items:center;justify-content:center;transition:all .15s}
.search-nav-btn:hover{background:var(--bgH);color:var(--t1);border-color:var(--ac)}
.search-nav-btn:disabled{opacity:.4;cursor:not-allowed}
.search-close{background:none;border:none;color:var(--t3);width:22px;height:22px;border-radius:5px;cursor:pointer;font-size:12px;display:flex;align-items:center;justify-content:center;transition:all .15s;margin-left:2px}
.search-close:hover{color:var(--dg);background:var(--dgg)}
/* Search highlighting */
.search-highlight{background:var(--ac);color:#000;padding:1px 2px;border-radius:3px;font-weight:500}
.search-highlight-current{background:var(--dg);color:#fff;padding:1px 2px;border-radius:3px;font-weight:500;box-shadow:0 0 8px rgba(255,68,102,.5)}

#inarea{padding:14px 22px 20px;background:transparent;z-index:10;position:absolute;bottom:0;left:0;right:0;pointer-events:none}
.iw{max-width:50%;margin:0 auto;display:flex;gap:8px;align-items:flex-end;pointer-events:auto}
/* Chutes.ai pill bar */
.ib{flex:1;display:flex;align-items:center;gap:6px;background:#2b2b2b;border-radius:28px;padding:6px 6px 6px 10px;min-height:46px}
#uinp{flex:1;resize:none;background:transparent;border:none;color:var(--t1);padding:6px 4px;font-family:inherit;font-size:15px;outline:none;min-height:32px;max-height:180px;line-height:1.5;overflow:hidden}
#uinp::placeholder{color:#666}
/* Plus button — circular outlined */
.plus-wrap{position:relative;display:inline-flex;align-items:center;flex-shrink:0}
#plusbtn{background:none;border:1.5px solid #565656;color:#999;cursor:pointer;width:30px;height:30px;border-radius:50%;display:flex;align-items:center;justify-content:center;transition:all .2s;font-size:17px;font-weight:300;line-height:1;padding:0;font-family:system-ui;flex-shrink:0}
#plusbtn:hover{color:#ddd;border-color:#888}
#plusbtn.active{color:var(--ac);border-color:var(--ac)}
#plusbtn.open{transform:rotate(45deg);color:#ddd;border-color:#888}
.plus-menu{display:none;position:absolute;bottom:calc(100% + 10px);left:50%;transform:translateX(-50%);background:var(--bg1);border:1px solid rgba(255,255,255,.08);border-radius:12px;padding:5px;z-index:300;box-shadow:0 10px 32px rgba(0,0,0,.7);min-width:160px}
.plus-menu.show{display:block;animation:fadeIn .15s ease}
.plus-opt{padding:9px 14px;font-size:13px;color:var(--t2);cursor:pointer;border-radius:7px;transition:background .1s;display:flex;align-items:center;gap:10px;font-family:'Outfit',sans-serif;white-space:nowrap}
.plus-opt:hover{background:var(--bgH);color:var(--t1)}
.plus-opt.active{color:var(--ac)}
/* Voice button — icon only, circular, no border */
#vbtn{background:none;border:none;color:#888;cursor:pointer;width:34px;height:34px;border-radius:50%;display:flex;align-items:center;justify-content:center;transition:all .2s;padding:0;flex-shrink:0}
#vbtn:hover{color:#bbb;background:rgba(255,255,255,.06)}
#vbtn.recording{color:#fff;background:var(--dg);animation:recPulse 1.2s ease-in-out infinite}
@keyframes recPulse{0%,100%{box-shadow:0 0 0 0 rgba(255,68,102,.5)}60%{box-shadow:0 0 0 7px rgba(255,68,102,0)}}
#sbtn{background:var(--ac);color:#000;border:none;width:55px;height:55px;border-radius:14px;cursor:pointer;transition:all .15s;flex-shrink:0;display:flex;align-items:center;justify-content:center}
#sbtn img{width:38px;height:auto}
#sbtn:hover{background:var(--acd)}#sbtn:disabled{opacity:.3;cursor:default}
#stopbtn{background:linear-gradient(135deg,#1a0505,#2d0808);color:#ff4466;border:1px solid rgba(255,68,102,.35);width:32px;height:32px;border-radius:8px;cursor:pointer;transition:all .2s cubic-bezier(.34,1.56,.64,1);flex-shrink:0;display:flex;align-items:center;justify-content:center;margin-right:4px;padding:0;box-shadow:inset 0 1px 0 rgba(255,255,255,.05)}
#stopbtn:hover{background:linear-gradient(135deg,#2d0808,#4a0d0d);border-color:rgba(255,68,102,.7);transform:scale(1.08);box-shadow:0 0 14px rgba(255,68,102,.4),inset 0 1px 0 rgba(255,255,255,.08)}
#stopbtn:disabled{opacity:.4;cursor:default}
.fb{display:inline-flex;align-items:center;gap:6px;background:var(--bg2);border:1px solid rgba(255,255,255,.06);padding:5px 12px;border-radius:8px;font-size:14px;color:var(--ac);margin-bottom:8px}
.fb button{background:none;border:none;color:var(--dg);cursor:pointer;font-size:16px}
.mo{display:none;position:fixed;inset:0;background:rgba(0,0,0,.6);z-index:999;align-items:center;justify-content:center;backdrop-filter:blur(4px)}.mo.show{display:flex}
.ml{background:var(--bg1);border:1px solid rgba(255,255,255,.06);border-radius:16px;padding:26px;min-width:460px;max-width:540px;box-shadow:0 14px 52px rgba(0,0,0,.6);max-height:90vh;overflow-y:auto}
.ml h3{margin-bottom:18px;font-size:20px;font-weight:600;display:flex;align-items:center;gap:10px}
.ml label{display:block;font-size:14px;color:var(--t2);margin-bottom:5px;margin-top:14px}
.ml input[type="text"],.ml input[type="password"],.ml textarea{width:100%;background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t1);padding:10px 14px;border-radius:10px;font-family:inherit;font-size:14px;outline:none}

.ml textarea{resize:vertical;min-height:80px;line-height:1.5}.ml input:focus,.ml textarea:focus{border-color:var(--ac)}.ml select{width:100%;background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t1);padding:10px 14px;border-radius:10px;font-family:inherit;font-size:14px;outline:none;cursor:pointer}.ml select:focus{border-color:var(--ac)}.ml select option{background:var(--bg2);color:var(--t1);padding:8px}
.sr{display:flex;align-items:center;gap:12px}.sr input[type="range"]{flex:1;accent-color:var(--ac)}
.sv{font-family:'JetBrains Mono',monospace;font-size:14px;color:var(--ac);min-width:55px;text-align:right}
.thr{display:flex;gap:8px;margin-top:8px}
.thc{padding:9px 18px;border-radius:8px;cursor:pointer;font-size:14px;font-weight:500;border:1px solid rgba(255,255,255,.06);background:var(--bg2);color:var(--t2);transition:all .15s}
.thc:hover{border-color:var(--ac);color:var(--t1)}.thc.active{border-color:var(--ac);color:var(--ac);background:var(--acg)}
.mbs{margin-top:20px;display:flex;gap:8px;justify-content:flex-end}
.mbs button{padding:10px 22px;border-radius:10px;cursor:pointer;font-family:inherit;font-size:14px;font-weight:500;transition:all .15s}
.bc{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t2)}.bc:hover{color:var(--t1)}
.bs{background:var(--ac);border:none;color:#000;font-weight:600}.bs:hover{background:var(--acd)}
/* Custom themed dialog */
#cdlg{display:none;position:fixed;inset:0;background:rgba(0,0,0,.5);z-index:9998;align-items:center;justify-content:center;backdrop-filter:blur(3px)}
#cdlg.show{display:flex}
.cdlg-box{background:var(--bg1);border:1px solid rgba(255,255,255,.06);border-radius:14px;padding:24px;min-width:320px;max-width:400px;box-shadow:0 12px 48px rgba(0,0,0,.6);text-align:center}
.cdlg-msg{font-size:15px;color:var(--t1);margin-bottom:20px;line-height:1.5}
.cdlg-btns{display:flex;gap:8px;justify-content:center}
.cdlg-btns button{padding:9px 24px;border-radius:8px;cursor:pointer;font-family:inherit;font-size:14px;font-weight:500;transition:all .15s}
.cdlg-ok{background:var(--ac);border:none;color:#000;font-weight:600}.cdlg-ok:hover{background:var(--acd)}
.cdlg-cancel{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t2)}.cdlg-cancel:hover{color:var(--t1)}
/* Global tooltip - fixed position, never clipped */
#gtip{position:fixed;background:var(--bg2);color:var(--t2);padding:5px 10px;border-radius:6px;font-size:11px;font-family:'Outfit',sans-serif;white-space:nowrap;z-index:99999;pointer-events:none;border:1px solid rgba(255,255,255,.06);box-shadow:0 4px 12px rgba(0,0,0,.4);opacity:0;transition:opacity .12s ease;display:none}
/* Copy feedback */
.copy-ok{position:fixed;top:20px;right:20px;background:var(--bg2);color:var(--ac);padding:8px 16px;border-radius:8px;font-size:13px;font-family:'Outfit',sans-serif;border:1px solid rgba(255,255,255,.06);z-index:9999;pointer-events:none}

/* Floating logos when sidebar collapsed - logos next to hamburger */
.float-logos{display:flex;align-items:center;gap:6px;opacity:0;transition:opacity .4s ease;pointer-events:none}
#sidebar.collapsed ~ #main .float-logos{opacity:1;pointer-events:auto}
.float-logos img{opacity:.7}
.float-logos .fl-t{height:22px}
.float-logos .fl-text{height:22px}
.float-logos .fl-spacer{width:18px}
/* Floating footer when sidebar collapsed - bottom left */
.float-footer{position:fixed;bottom:16px;left:16px;display:flex;align-items:center;gap:6px;opacity:0;transition:opacity .4s ease,transform .4s ease;transform:translateX(-20px);pointer-events:none;z-index:20}
#sidebar.collapsed ~ #main .float-footer{opacity:1;transform:translateX(0);pointer-events:auto}
.float-footer .ff-chutes{height:25px;opacity:.4}
.float-footer .ff-text{font-size:15px;color:var(--t3)}
/* Wider chat bubbles */
.msg{max-width:72%}
#sidebar.collapsed ~ #main .msg{max-width:78%}
#sidebar.collapsed ~ #main .iw{max-width:50%}
#sidebar.collapsed ~ #main .es{max-width:100%}
.spm textarea{min-height:120px}
@media(max-width:700px){#sidebar{width:240px;min-width:240px}#sidebar.collapsed{margin-left:-240px}.ml{min-width:90vw}}
/* --- Folder filter button --- */
.folder-filter-row{padding:0 10px 8px;display:flex;align-items:center;gap:6px}
.folder-filter-btn{display:flex;align-items:center;gap:6px;padding:5px 10px;border-radius:18px;font-size:12px;font-weight:600;cursor:pointer;border:1.5px solid rgba(255,255,255,.1);color:var(--t2);background:var(--bg2);transition:all .15s;min-width:0;max-width:180px;overflow:hidden}
.folder-filter-btn:hover{border-color:var(--ac);color:var(--t1)}
.folder-filter-btn .ffdot{width:7px;height:7px;border-radius:50%;flex-shrink:0}
.folder-filter-btn .fflabel{overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.folder-filter-btn .ffarrow{font-size:8px;color:var(--t3);margin-left:auto;flex-shrink:0}
.folder-drop{position:fixed;z-index:99999;background:var(--bg1);border:1px solid rgba(255,255,255,.1);border-radius:12px;padding:6px;min-width:190px;box-shadow:0 8px 32px rgba(0,0,0,.6);display:none}
.folder-drop.show{display:block}
.folder-drop-item{display:flex;align-items:center;gap:8px;padding:8px 10px;border-radius:7px;font-size:13px;color:var(--t2);cursor:pointer;transition:background .1s}
.folder-drop-item:hover{background:var(--bgH);color:var(--t1)}
.folder-drop-item.active{color:var(--ac)}
.folder-drop-item .fdi-dot{width:8px;height:8px;border-radius:50%;flex-shrink:0}
.folder-drop-sep{height:1px;background:rgba(255,255,255,.06);margin:4px 0}
.folder-drop-add{display:flex;gap:6px;padding:6px 6px 2px}
.folder-drop-add input{flex:1;background:var(--bg2);border:1px solid rgba(255,255,255,.08);color:var(--t1);padding:6px 10px;border-radius:7px;font-family:inherit;font-size:12px;outline:none}
.folder-drop-add input::placeholder{color:var(--t3)}
.folder-drop-add input:focus{border-color:var(--ac)}
.folder-drop-add button{background:var(--ac);color:#000;border:none;border-radius:7px;padding:6px 10px;cursor:pointer;font-size:12px;font-weight:700;white-space:nowrap}
.folder-drop-edit{display:flex;align-items:center;gap:4px;flex:1;min-width:0}
.folder-drop-edit-inp{flex:1;background:var(--bg0);border:1px solid var(--ac);color:var(--t1);padding:3px 6px;border-radius:5px;font-family:inherit;font-size:12px;outline:none;min-width:0}
.folder-drop-actions{display:flex;gap:2px;margin-left:auto;flex-shrink:0}
.fda-btn{background:none;border:none;color:var(--t3);cursor:pointer;padding:3px 5px;border-radius:4px;font-size:12px;line-height:1}
.fda-btn:hover{color:var(--t1);background:var(--bgH)}
.fda-btn.del:hover{color:var(--dg)}
.ftab.dragover{box-shadow:0 0 0 2px var(--ac)}
.ci-fdot{width:7px;height:7px;border-radius:50%;flex-shrink:0;display:inline-block;margin-right:3px;vertical-align:middle}
.ci.dragging{opacity:.35;pointer-events:none}
/* --- Prompt templates modal (wide) --- */
#tplm .ml{min-width:min(1100px,92vw);max-width:min(1200px,95vw)}
.tpl-search-wrap{position:relative;margin-bottom:10px}
.tpl-search-wrap input{width:100%;background:var(--bg2);border:1px solid rgba(255,255,255,.08);color:var(--t1);padding:8px 34px 8px 12px;border-radius:8px;font-family:inherit;font-size:13px;outline:none;box-sizing:border-box}
.tpl-search-wrap input:focus{border-color:var(--ac)}
.tpl-search-wrap input::placeholder{color:var(--t3)}
.tpl-search-wrap .tpl-search-clear{position:absolute;right:8px;top:50%;transform:translateY(-50%);background:none;border:none;color:var(--t3);cursor:pointer;font-size:15px;padding:2px 4px;line-height:1}
.tpl-search-wrap .tpl-search-clear:hover{color:var(--t1)}
.tpl-list{display:flex;flex-direction:column;gap:6px;max-height:340px;overflow-y:auto;margin-bottom:14px}
.tpl-item{display:flex;align-items:flex-start;gap:10px;padding:13px 16px;background:var(--bg0);border-radius:10px;cursor:pointer;border:1px solid transparent;transition:all .15s}
.tpl-item:hover{border-color:var(--ac);background:var(--bgH)}
.tpl-item-text{flex:1;min-width:0}
.tpl-name{font-size:15px;color:var(--t1);font-weight:600;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;margin-bottom:3px}
.tpl-preview{font-size:13px;color:var(--t3);overflow:hidden;text-overflow:ellipsis;white-space:nowrap}
.tpl-del{background:none;border:none;color:var(--t3);cursor:pointer;padding:5px 7px;border-radius:5px;font-size:17px;line-height:1;flex-shrink:0;margin-top:1px}.tpl-del:hover{color:var(--dg)}
.tpl-empty{text-align:center;padding:24px;color:var(--t3);font-size:14px}
.tpl-add{border-top:1px solid rgba(255,255,255,.06);padding-top:14px;display:flex;flex-direction:column;gap:9px}
.tpl-add input,.tpl-add textarea{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t1);padding:10px 14px;border-radius:9px;font-family:inherit;font-size:14px;outline:none;resize:vertical;width:100%;box-sizing:border-box}
.tpl-add textarea{min-height:90px}
.tpl-add input::placeholder,.tpl-add textarea::placeholder{color:var(--t3)}
.tpl-add input:focus,.tpl-add textarea:focus{border-color:var(--ac)}
/* --- Summary modal --- */
.summ-loading{text-align:center;padding:22px;color:var(--t3);font-size:13px}
#summ-result{background:var(--bg2);border-radius:10px;padding:14px 16px;font-size:14px;line-height:1.7;color:var(--t2);max-height:380px;overflow-y:auto;min-height:60px}
#summ-result ul{margin:4px 0;padding-left:18px}
#summ-result li{margin-bottom:5px}
/* --- Move to folder modal --- */
.foldr-list{display:flex;flex-direction:column;gap:6px;margin-bottom:14px}
.foldr-btn{display:flex;align-items:center;gap:10px;padding:11px 14px;border-radius:10px;cursor:pointer;border:1.5px solid transparent;background:var(--bg2);transition:all .15s;font-size:14px;color:var(--t2)}
.foldr-btn:hover{border-color:var(--ac);color:var(--t1)}
.foldr-btn.current{border-color:rgba(255,255,255,.2);color:var(--t1)}
.foldr-btn .fdot2{width:10px;height:10px;border-radius:50%;flex-shrink:0;display:block}
.foldr-btn .fcur{margin-left:auto;font-size:11px;color:var(--t3);font-style:italic}
.foldr-add-row{display:flex;gap:8px;border-top:1px solid rgba(255,255,255,.06);padding-top:12px}
.foldr-add-row input{flex:1;background:var(--bg2);border:1px solid rgba(255,255,255,.08);color:var(--t1);padding:9px 12px;border-radius:8px;font-family:inherit;font-size:13px;outline:none}
.foldr-add-row input::placeholder{color:var(--t3)}
.foldr-add-row input:focus{border-color:var(--ac)}
.foldr-add-row button{background:var(--ac);color:#000;border:none;border-radius:8px;padding:9px 14px;cursor:pointer;font-weight:700;font-size:13px}
/* --- Draft badge on chat items --- */
.ci-draft{width:6px;height:6px;border-radius:50%;background:#f59e0b;flex-shrink:0;display:inline-block;margin-left:2px;vertical-align:middle}
/* --- Context Window Meter --- */

/* --- Compare mode --- */
#compare-overlay{display:none;position:fixed;inset:0;background:rgba(0,0,0,.75);z-index:900;flex-direction:column}
#compare-overlay.show{display:flex}
.cmp-header{display:flex;align-items:center;justify-content:space-between;padding:14px 20px;background:var(--bg1);border-bottom:1px solid rgba(255,255,255,.06);flex-shrink:0}
.cmp-header h3{font-size:16px;font-weight:600;color:var(--t1)}
.cmp-selectors{display:flex;gap:10px;align-items:center}
.cmp-sel{background:var(--bg2);border:1.5px solid var(--bd);color:var(--t1);padding:7px 32px 7px 12px;border-radius:10px;font-family:'Outfit',sans-serif;font-size:13px;cursor:pointer;outline:none;-webkit-appearance:none;appearance:none;min-width:200px;color-scheme:dark;background-image:url("data:image/svg+xml,%3Csvg xmlns='http://www.w3.org/2000/svg' width='12' height='8' viewBox='0 0 12 8'%3E%3Cpath d='M1 1l5 5 5-5' stroke='%237878a0' stroke-width='1.5' fill='none' stroke-linecap='round'/%3E%3C/svg%3E");background-repeat:no-repeat;background-position:right 10px center}
.cmp-sel:focus{border-color:var(--ac);background-color:var(--bg2)}
.cmp-sel option{background:var(--bg2);color:var(--t1)}
.cmp-close{background:none;border:none;color:var(--t3);cursor:pointer;font-size:20px;padding:4px 8px;border-radius:6px}
.cmp-close:hover{color:var(--dg)}
.cmp-body{flex:1;display:flex;overflow:hidden;min-height:0}
.cmp-pane{flex:1;display:flex;flex-direction:column;min-width:0;min-height:0}
.cmp-pane:not(:last-child){border-right:2px solid var(--bd)}
.cmp-pane-head{padding:10px 16px;background:var(--bg1);border-bottom:1px solid rgba(255,255,255,.06);font-size:11px;color:var(--t3);font-weight:700;letter-spacing:.06em;text-transform:uppercase;display:flex;align-items:center;justify-content:space-between;position:relative;flex-shrink:0;min-height:38px}
.cmp-pane-model{font-size:13px;color:var(--ac);font-weight:700;text-transform:none;letter-spacing:0;white-space:nowrap;overflow:hidden;text-overflow:ellipsis;max-width:60%;position:absolute;left:50%;transform:translateX(-50%);text-align:center}
.cmp-pane-body{flex:1;overflow-y:auto;padding:16px;background:var(--bg0)}
.cmp-pane-body .mc{font-size:14px;line-height:1.7;color:var(--t1)}
.cmp-input-row{padding:12px 20px;background:var(--bg1);border-top:1px solid rgba(255,255,255,.06);display:flex;justify-content:center;flex-shrink:0}
.cmp-input-wrap{display:flex;align-items:center;gap:8px;width:100%;max-width:900px}
.cmp-ib{flex:1;display:flex;align-items:center;gap:6px;background:#2b2b2b;border-radius:28px;padding:6px 6px 6px 14px;min-height:46px}
#cmp-inp{flex:1;resize:none;background:transparent;border:none;color:var(--t1);padding:6px 4px;font-family:inherit;font-size:15px;outline:none;min-height:32px;max-height:180px;line-height:1.5;overflow:hidden}
#cmp-inp::placeholder{color:#666}
.cmp-run-btn{background:var(--ac);color:#000;border:none;width:55px;height:55px;border-radius:14px;font-family:inherit;font-size:24px;font-weight:700;cursor:pointer;flex-shrink:0;display:flex;align-items:center;justify-content:center;transition:all .15s}
.cmp-run-btn:hover{background:var(--acd)}
#cmp-vbtn{background:none;border:none;color:#888;cursor:pointer;width:34px;height:34px;border-radius:50%;display:flex;align-items:center;justify-content:center;transition:all .2s;padding:0;flex-shrink:0}
#cmp-vbtn:hover{color:#bbb;background:rgba(255,255,255,.06)}
#cmp-vbtn.recording{color:#fff;background:var(--dg);animation:recPulse 1.2s ease-in-out infinite}
#cmp-plusbtn{background:none;border:1.5px solid #565656;color:#999;cursor:pointer;width:30px;height:30px;border-radius:50%;display:flex;align-items:center;justify-content:center;transition:all .2s;font-size:17px;font-weight:300;line-height:1;padding:0;font-family:system-ui;flex-shrink:0}
#cmp-plusbtn:hover{color:#ddd;border-color:#888}
#cmp-plusbtn.active{color:var(--ac);border-color:var(--ac)}
#cmp-plusbtn.open{transform:rotate(45deg);color:#ddd;border-color:#888}
/* --- Keyboard shortcut modal --- */
.kb-grid{display:grid;grid-template-columns:1fr 1fr;gap:6px 24px;margin-top:4px}
.kb-row{display:flex;align-items:center;justify-content:space-between;padding:6px 0;border-bottom:1px solid rgba(255,255,255,.04)}
.kb-desc{font-size:13px;color:var(--t2)}
.kb-keys{display:flex;gap:4px;align-items:center}
.kb-key{background:var(--bg2);border:1px solid rgba(255,255,255,.12);color:var(--t1);padding:2px 7px;border-radius:5px;font-size:11px;font-family:'JetBrains Mono',monospace;font-weight:600}
/* --- Theme editor modal --- */
.te-row{display:flex;align-items:center;gap:12px;margin-bottom:14px}
.te-label{font-size:13px;color:var(--t2);min-width:110px}
.te-swatch{width:28px;height:28px;border-radius:6px;border:2px solid rgba(255,255,255,.12);cursor:pointer;flex-shrink:0}
.te-slider{flex:1;accent-color:var(--ac)}
.te-val{font-size:12px;color:var(--t3);min-width:42px;text-align:right;font-family:'JetBrains Mono',monospace}
.te-preview{margin-top:10px;padding:14px;background:var(--bg0);border-radius:10px;border:1px solid rgba(255,255,255,.06)}
/* --- Improved export modal --- */
.exp-grid{display:grid;grid-template-columns:repeat(3,1fr);gap:10px;margin-top:4px}
.exp-btn{display:flex;flex-direction:column;align-items:center;gap:8px;padding:18px 12px;background:var(--bg2);border:1.5px solid rgba(255,255,255,.08);border-radius:12px;cursor:pointer;transition:all .15s;color:var(--t2);font-family:inherit}
.exp-btn:hover{border-color:var(--ac);color:var(--t1);transform:translateY(-2px);box-shadow:0 6px 20px rgba(0,0,0,.4)}
.exp-btn svg{width:28px;height:28px;stroke:currentColor;fill:none;stroke-width:1.5;stroke-linecap:round;stroke-linejoin:round}
.exp-btn span{font-size:13px;font-weight:600}
.exp-btn small{font-size:11px;color:var(--t3)}
/* --- Multi-file badges --- */
.fb-list{display:flex;flex-wrap:wrap;gap:6px;margin-bottom:8px}
/* --- URL fetch indicator --- */
.url-fetch-badge{display:inline-flex;align-items:center;gap:5px;background:rgba(0,170,221,.12);border:1px solid rgba(0,170,221,.3);color:#00aadd;padding:3px 9px;border-radius:6px;font-size:11px;margin-bottom:6px;cursor:pointer}
.url-fetch-badge:hover{border-color:#00aadd}
/* floating KB button */
.prof-item{display:flex;align-items:center;gap:10px;flex-wrap:wrap;padding:11px 14px;border-radius:10px;background:var(--bg2);border:1.5px solid transparent;margin-bottom:8px;transition:all .15s}
.prof-item:hover{border-color:rgba(255,255,255,.12)}
.prof-item.active-prof{border-color:var(--ac)}
.prof-name{flex:1;font-size:14px;font-weight:600;color:var(--t1)}
.prof-badge{font-size:10px;padding:2px 7px;border-radius:10px;background:var(--acg);color:var(--ac);font-weight:700;letter-spacing:.04em}
.prof-pw-row{display:flex;gap:6px;margin-top:8px}
.prof-pw-row input{flex:1;background:var(--bg0);border:1px solid rgba(255,255,255,.08);color:var(--t1);padding:7px 10px;border-radius:8px;font-family:inherit;font-size:13px;outline:none}
.prof-pw-row input:focus{border-color:var(--ac)}
.prof-pw-row button{background:var(--ac);color:#000;border:none;border-radius:8px;padding:7px 13px;cursor:pointer;font-size:13px;font-weight:700;white-space:nowrap}
.prof-del-btn{background:none;border:none;color:var(--t3);cursor:pointer;padding:4px 6px;border-radius:5px;font-size:13px;line-height:1;transition:all .15s}
.prof-del-btn:hover{color:var(--dg);background:var(--dgg)}
.prof-sep{height:1px;background:rgba(255,255,255,.06);margin:14px 0}
.prof-create{display:flex;flex-direction:column;gap:9px}
.prof-create input{background:var(--bg2);border:1px solid rgba(255,255,255,.06);color:var(--t1);padding:10px 14px;border-radius:9px;font-family:inherit;font-size:14px;outline:none}
.prof-create input:focus{border-color:var(--ac)}
.prof-create input::placeholder{color:var(--t3)}
.prof-create-btn{display:flex;align-items:center;gap:8px;align-self:flex-end;padding:10px 24px;background:linear-gradient(135deg,var(--ac),var(--acd));border:none;color:#000;border-radius:10px;font-family:inherit;font-size:14px;font-weight:700;cursor:pointer;transition:all .18s;box-shadow:0 4px 16px rgba(0,0,0,.3);letter-spacing:.02em}
.prof-create-btn:hover{transform:translateY(-2px);box-shadow:0 6px 22px rgba(0,0,0,.45);filter:brightness(1.08)}
.prof-create-btn:active{transform:translateY(0);box-shadow:0 2px 10px rgba(0,0,0,.3)}
.prof-switch-btn{display:inline-flex;align-items:center;gap:6px;padding:6px 14px;background:var(--ac);border:none;color:#000;border-radius:8px;font-family:inherit;font-size:13px;font-weight:700;cursor:pointer;transition:all .18s;white-space:nowrap;box-shadow:0 2px 10px rgba(0,0,0,.25)}
.prof-switch-btn:hover{transform:translateY(-1px);box-shadow:0 4px 16px rgba(0,0,0,.4);filter:brightness(1.1)}
.prof-switch-btn:active{transform:translateY(0)}
#kb-float-btn{position:fixed;bottom:22px;right:22px;width:38px;height:38px;border-radius:50%;background:var(--bg2);border:1.5px solid rgba(255,255,255,.12);color:var(--t2);font-size:16px;font-weight:700;cursor:pointer;z-index:800;display:flex;align-items:center;justify-content:center;transition:all .15s;box-shadow:0 4px 18px rgba(0,0,0,.45);font-family:'Outfit',sans-serif;line-height:1}
#kb-float-btn:hover{background:var(--ac);color:#000;border-color:var(--ac);transform:scale(1.1)}
/* summarize button – solid */
.tb-solid{background:var(--bg2);border:1px solid rgba(255,255,255,.12);color:var(--t2);padding:8px 12px;border-radius:10px;cursor:pointer;transition:all .15s;display:flex;align-items:center;justify-content:center;pointer-events:auto}
.tb-solid:hover{background:var(--bgH);color:var(--t1);border-color:var(--ac)}
/* folder color picker in add-folder row */
.fd-color-btn{-webkit-appearance:none;appearance:none;width:28px;height:28px;min-width:28px;max-width:28px;border-radius:6px;border:1.5px solid rgba(255,255,255,.15);padding:0;cursor:pointer;flex-shrink:0;overflow:hidden;background:transparent}
/* --- Image Folder Panel --- */
#img-folder-panel{position:fixed;top:0;right:0;width:280px;height:100vh;background:rgba(12,12,20,.72);backdrop-filter:blur(18px);-webkit-backdrop-filter:blur(18px);border-left:1px solid rgba(255,255,255,.08);z-index:500;display:flex;flex-direction:column;transform:translateX(100%);transition:transform .28s cubic-bezier(.4,0,.2,1);box-shadow:-8px 0 40px rgba(0,0,0,.4)}
#img-folder-panel.open{transform:translateX(0)}
.ifp-header{padding:16px 14px 10px;display:flex;align-items:center;gap:8px;border-bottom:1px solid rgba(255,255,255,.07);flex-shrink:0}
.ifp-title{flex:1;font-size:13px;font-weight:600;color:var(--t1);letter-spacing:.02em;white-space:nowrap;overflow:hidden;text-overflow:ellipsis}
.ifp-close{background:none;border:none;color:var(--t3);cursor:pointer;padding:4px 6px;border-radius:5px;font-size:16px;line-height:1;display:flex;align-items:center;transition:all .15s}.ifp-close:hover{color:var(--t1);background:var(--bgH)}
.ifp-load-btn{margin:10px 12px 4px;padding:8px 0;background:var(--bg2);border:1px solid rgba(255,255,255,.08);color:var(--t2);border-radius:8px;font-family:inherit;font-size:12px;cursor:pointer;display:flex;align-items:center;justify-content:center;gap:6px;transition:all .15s;width:calc(100% - 24px)}
.ifp-load-btn:hover{border-color:var(--ac);color:var(--t1)}
.ifp-count{font-size:11px;color:var(--t3);padding:4px 14px 8px;letter-spacing:.02em}
.ifp-list{flex:1;overflow-y:auto;padding:4px 8px 16px}
.ifp-item{position:relative;padding:7px 10px;border-radius:7px;cursor:default;display:flex;align-items:center;gap:8px;transition:background .12s;margin-bottom:2px;border:1px solid transparent}
.ifp-item:hover{background:rgba(255,255,255,.06);border-color:rgba(255,255,255,.07)}
.ifp-thumb{width:32px;height:32px;border-radius:4px;object-fit:cover;flex-shrink:0;border:1px solid rgba(255,255,255,.08)}
.ifp-name{flex:1;font-size:12px;color:var(--t2);white-space:nowrap;overflow:hidden;text-overflow:ellipsis;font-family:'JetBrains Mono',monospace}
.ifp-item:hover .ifp-name{color:var(--t1)}
/* Hover preview popup */
.ifp-pdf-thumb{width:32px;height:32px;border-radius:4px;flex-shrink:0;border:1px solid rgba(255,255,255,.08);background:rgba(220,50,50,.15);display:flex;align-items:center;justify-content:center}
.ifp-preview{display:none;position:fixed;background:none;border:none;border-radius:0;padding:0;z-index:9999;box-shadow:none;pointer-events:none;max-width:728px;width:728px}
.ifp-preview.show{display:block}
.ifp-preview img{width:100%;max-height:504px;object-fit:contain;border-radius:6px;background:none;display:block;opacity:.65}
.ifp-preview embed{width:100%;height:460px;border-radius:6px;display:block}
.ifp-preview-name{display:none}
.ifp-preview-meta{display:none}
/* Hamburger toggle button (fixed, right edge) */
#img-folder-btn{position:fixed;bottom:70px;right:18px;width:38px;height:38px;background:var(--bg2);border:1.5px solid rgba(255,255,255,.10);color:var(--t2);border-radius:10px;cursor:pointer;z-index:501;display:flex;align-items:center;justify-content:center;transition:all .15s;box-shadow:0 4px 18px rgba(0,0,0,.45)}
#img-folder-btn:hover,#img-folder-btn.active{background:var(--ac);color:#000;border-color:var(--ac)}
</style></head><body>
<div id="sidebar">
 <div class="sh"><img src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAACQAAAAkCAYAAADhAJiYAAAF6klEQVR4nMWYwYtcWRXGv3PuK7oqXel0JaY6SYNCFr0waTUQhRbHoRaZTQguREE3AREHSVwpGEigFsMgDEMkf4BkpWAHJAsJIoFmWg0ok+kYk1aysG06YUyqU13Vlfeq3333nuOi3muqK1XdVZlIDtzFu1X3vN/9znn3nfMIe5iqGgBCRAoAzWZzJgiCd4wxbxljTqjqMWNMAQBEZAvAf0VkWVX/FIbh7UOHDj3o9kVEfq977gUDAAjD8FvW2t8751o6pDnnImvtH6Mo+m61WuUunzQqCKkqAcCzZ8++Gcfxn7tv5DvmnHPivRcRURFR770458Q55733vnuNtfbjRqPxTq//oWC6VHnPOScphEvHNsBew3sv2boMLIqiX1ar1SC9F+8JQ0S4c+dOod1u3+wSxA0LMWg457yqOlXVra2t2ysrK5O7QqUy8qNHj8biOL6dkljvffeO1Tn3ysN7r957m0L9tVar7R8YPlUNACCKovkMphfkdVjqz6qqxnH8h2q1yqpqMihKYQwR+Var9dNisfihiCRElEt/A3NH1bW1NbRare3rUUxEMD09jYmJCagqACTMnIui6P3x8fErGQOlMdT19fWZUqn0d2NMICJMRKSqICJsbGzo5cuXdXFxkZ1zIKLM6VCWy+XQaDRw48YNnZubI+ccmFmZ2YsIt1qtr01OTt5VVRMAICKSKIo+MMaMee89M++I6aVLl3Dz5k2emppCLpfrDTVEBMwMopdTgZlRq9Vw+vRpOXXqFIkIjDEAQCJCzMz5fP4qgLcBdOJWr9e/fODAgSUASkScSWyMwfLysp45c4YKhQJE5CVlxsbGUCgUEIYhrLVgZmTKorNbzM7O6tWrV3H8+HHy3u8IeXpym42NjbcPHjy4GABAPp//ETOTiHgAnDlSVRw9epTm5+dhjNkBk+XWtWvXtNls4sKFCzo1NcUZVLapYrGI2dlZIqJtJbtNRJSZdd++fe8CWMSTJ0/2WWtXswOn9/zYyy5evKiHDx/Ws2fP+vX19b4L0nNo4OGpqpokSe3x48eHeP/+/aeDIPh8d7h6dgDn3EvDWgvvPZIkQbFYxN27d/nKlSsqIrDWbv/Pe7/jSe016sRWgiD43MTExDeCIAjmUjmFiEyfBVkS9p0nInjvUSqVcP/+fW632xgfH4eI9E3yfpaGDblc7uvMzCeGWrWHqSq897DWvspaAgBjzBeZiI6l86OVBANsWFUGrDnG3vuJ1wHyOoyZC5zL5UbX+P9kqiosIhvpxZvmgYg0WVVXACCrmd+QaQr0H3bOLaWTryWpP4up6iccRdFfvPcxM/ObChsRGe+9WGs/4nK5/Mh7/wk6sr16i4JOHo66KVUVIlIR+efS0tI9BgBr7a8BEDO/skTMjCRJYK2FiIyyVACQtfa3lUrFsapSvV7/jXPuOQBW1aG8ZUqcPHlS2+028vk8ms0m1tbWdNiKUlWVmdl7H7bb7V9lk1kt/bP0jZ8M2eKo914bjYacO3fOl8tlPXLkiFYqFbewsCBRFGl3gzDAR6KqGobhLzKWrOKnhw8fBjMzMx/ncrlZEfH9XrT9VGJmNJtNvX79Oj148ABxHGNubk7Onz9PhUKhU7j3eZ2IiBhj2Dm3Uq/Xv1QulyOkj/92b/T06dOvOOe20tLopdpokFK7dRi7NZDOOanX62+lDDsFyCZardb3Un9uFKgkSbb7ryRJdoVR1URVdXNz88d9YbqgAgAIw/DdbJfOuc/ctXbBOFWVFObn3ffcLS+CVKnvOOc2s0QfVq0BID5LYOfcVrPZ/MFQML1Qz58/PxHH8UeZWqN8cOj3oSGO47/VarWvjgTTBbUd1xcvXvzQWvuv3oTNPst0D1V1vc2BtfbfYRj+pOurx55P8CAozsrLW7dujYVh+O04jueTJHk88NFKLUmST621v2u1Wt+/d+/eeLfP3e451Bt+YWEhqFQqLrteXV0tlUqlWWPMCWb+gqpOotNANAGsOueWNzc3/zE9Pb3eBRIA8HuVOf8Dz+oGAsYI398AAAAASUVORK5CYII=" alt="" style="height:27px;opacity:.7"><div class="sh-spacer"></div><img src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAE4AAAAWCAYAAABud6qHAAALrElEQVR4nO1Ye2xb5RU/597rmziJ49hJaBO39JU2WVra0KQDViANezEBGxNTB4itmzYqRjc6VKCdJi3t0BB0QmIgbRTGxooY4JRRha4tdJtLk+bhOA+7eTixGzt183AcO37fa9/H2R+2WQgpsE2apo0jXdn3+373nO+c73zfeQB8Sv8eERF+QigCfGLspwTwQeP+E4b+nyQOAMBsNrMAAENDQ3To0CH1SmBEJADgGhsbARHl/9Aa/3tpZMT3VCymuM6d69kOAEBE7ML5nHe1tdmuHh72OpzOSwP9/aMmIsL/V89jAAA0Gq5Sq2WqOC5PS0Tocrk4IuIsFguXNUwWx9ZUVa26Zu3alVt4XrMOAMDlcvEWiyWH5YiIAch4scWS4bH4WYgDACAiJju2UOb7lPsOAKC5uZmxWK6MXYxZLGshZTeezeHM5g86TE52zpFyOgFARqbLNfVSLEZKR0ffF64ggCEiNJvNbE+P4wGr1b7HbDbzV9qNfybQfIRSHxr/KAN8AgzzUe8LeS21/sV4LvuLLAsMx2mS9fX1mjfeaP1GXh5/dTyeuGy19p1GxCAR4c6dO9Xz520+jYZjdu7cmT5y5IimsfFLtyGSPpEQobx8GT83559AxHet1guf0+uL14uiCAyDnKIoRAQqIqMajQZudtZ/uaHhmncRkbq6HPVGo+FmrbYwP5GIuU+d+tMZRAwTEb7wwgvctm033cMwLCDi0ZMnLSuqq2u+xvN5umg0PH769FtvI6KQUxYR1RMnzqzdsKG6UastWp5Op2J+/2QHIvYRESIiNTcTg4jqgw8+WLRnzyOf12oLa4gUKRQK9iHi2ZxxOjo6tDpd2b2CkAwj4pu9vY56ni+4cWJi+igAALhcU7+PRNLU1tb99eFhb58sE8ViKokikccz4+vu7r0ZAOD22+sLHA634nZPUXNzc/7g4MX3iIgEIYOPRhWy20fv7+623yvLRKlU5gkEkhSPZ/6LIpGqEtlsQ48QETs4ePHo3Fxmfm4uSYJA5PXOXrJaHY0AAK2trWVer5/c7qnprq6BWzwevxoOSxSNKpRIEA0Neaytra1luQDX2Tnwg0uXAnIymZGbTBL5/XEaGHA9mfUcFgDAZhu6bnx8ZlQQiOJxomSSKByWaGTEe/L06Q4jAMCJE39dFQ6nqbd3eOD8+d7dokgUDAp07FhrDQcAwDDIRKNxMBrLnuV5PtTbO7AzmRRiV121/Jsm08rviGLqze7u7prHHnssgoizkiSxxcXFBaoqt4yMeN9KpVJSYWFBgyimqrZsqX7RZnPcNjkZeDgQmJN4nmUR2QOiKL7FcdxwYaGWi0Yj7oaGzX/u73c+U1dX/S2nc6InEgn9IhqNTJaXL2+qqDAdNhrL3jxz5kz1/Py8KIqpKMuyalFR0TFFUV6yWruPFBbml5aWlj+5fv3qbamU8OjWrbX7z5+3bjKZVvxaURSfw2H/aTg8N1pWVrHaYCg9vGlT1X6rdbADEVstlq4Ver3+hNFoLBsfH/99MDhn5nlNUUlJ6cPr16/6iixLRwHgdr2ep7m5uTTHcXqdrviXgUDwj1NTl9tTqZg/63GTv4vHZXI6vY7m5ub8hWf5wgX3K5JENDTk3pt5vxh0Oi/N7969u2AhbnDQfYsgEPX1jexbfD+43dPewUG3deHYmTNtV09MzMou1+WRxfieHscDkkTU3+/cAwDgdE6E5ufTZLePPbcQ19bWeWMwKCgjI96+DM57nySR0t/vfGYhrqOjr1GSyOZwjB0AAOjvdz4hSUR2u+v5hbinn35aOzrquxAKCWSzDW89caLNMDrqE/3+GNlsF/YuxDIAmXuB51kIh8O/PHTokGixePJPnjyZR0RMKBT6TSKRJoZhv5j9hgAAdDodZiNOvsVCHCIVAwAUFBRcR0SMx5MZz0QhkvPztdtOnTpVkY12bHFx8Y0Gg5FNJBLBtjbb3vPne/e3t9sebW+3PirL8tpYLKVqNJqbAYAFAE0oFFQmJyefyhzvQZ6IWFGMjwSDQVWWZSMAgCCkBoPBGGMwGB9wOr2/s9td3+/tHbxVVZMujQYbNm/e8CQAgFarbQwGYzQ3F3qOiJjBwUHe4/Hk79u3T0gmY0f0+nzKy9Pc0tPzlxTLMnnh8Hzk7bePvUhErM1m0wAAZg0HjCgSpFKJCSLCQGC11N3dLSGiOjHhnZ6fj6AsK2U5I+es3tTUJAOA3NSEMhGoRACUuYFVrzcz3tSUSZQlSZJCoZDU1IQyIiocxy4XBJFKSgzbN2/e8kxt7TVPbty4+fCmTXWHq6pqHgVARpZlHQAgAGoURQkAOAKIqLS0tMiIqDidTlAUhQBAAQDm2murB6amLu8homRl5crvbthQ9WJ19cZTJlPVpMs12drZ2b86owJWBINB1eudmEVEdePGjZLX65WJCBmG8ckyIcNwZQAAHKcBWZanAUAEALWhoUECgEyIJQLieYT8/IJSAIDaWmDvuOMOlojYDRuqinW6ImIYNpp1uH814f3Ad0QQKi4uxPn54EuTk956n2/icx7P6Gc9Hs91Pt/FupmZiRui0fCPX3nllQIiVWEYhqusbOQAAA4ePLgUfzp69Gjh1q21z587d3r92NjI591u97e93vHDqVRqbNWqyjt0Ot3zAEBEENbr9bh27cpCImLPnj3LmkwmNuMYjJFlkVRViUejUWIYBiRJFhdXVFxWCQAg0On0P0TEtwAgnQM4HKP3lpRo0e+Hjix2yfyHZRkiIkDEpQxLAAAajeb9OVFMdoZCYdBqC+traq7uWwju6rLvqa5e/b1IJHn4vvvcZqfzZiRS84hIs6TFCDkAoC1bPvtEOCzeMzAwcFdDw6a/5eb379//xP337w2wLFcLAJBIxM+tW7fy2pkZ3S5EPJSFyQAAQ0OeXaKoYCIR77jhhhtYyhjnQzrl8jgmEAjLqqpuHR72HJufD/yC4/gEx7F3Go3l+3w+f1qSUn8AAAaA0ojELWYkSTLHMKgSqcKHFSNOVVV5dnZWyL6ziDja2zt8vK7uM3cODV08FY2GnmIYzZSq0naTacWvIpEkC6D2Hjy4mr/7bpAAUM4ey8W8ZYYhGQAAkbHqdHkPmUwrnn/vva4fCULKVVRUyOfl5d2zbNkyzfj4+EDGcNEXfL6ZBysqTD/r7x9KRiLR4xzHFBoMVz28Zs3qxosXxy8Iwvw5o9G4TFUVmUhVljScqqpFK1YYuJmZiQN6ffkTW7ZsuysWi4LBUAyBQEiJRiO76+pqPPX19QUAWIrIagRByO0CAQAkk4nxVEpiOI7f2t7ertu+HRJEmUQznU7PrlmzZnUodF0DEb139uxZbG4mJp2+8IDb7Stft27trYKw8lZBSEJJiR6CwZA6Nub83o4d17vM5pPlPM8XL+VpPM9jUVGRNhIJGwCAefzxn75+4MDPb6uurrmntLT8r+FwCHg+D8rKDDA5OX1JFBOPZJPg4a4u+66Kisrf1tbWHg4Go4c1Gg6KiwvA55sei8fD32xqapItFoumpMTARSIR45KGk2XpVVGEEVlmj1+65DvOMPiYLMtrotGId3Z27g833bTVTkSYuRjVg4qiQHl5eTqzy6hmF2Pv7LR/t7Jy+VeJ2OsR8Uy2vlRFMf6dQMD/E57XfPHZZ5/tfOihh9I7dgAgbvYDQOPIiPdulmW/kE6nSwQhOTY97X9tx47rB4gIX3755ZgkpX6GyInj4+PxhZsFAMl4PH4QEWMAAC0tx5SWlpZ77XbXG0VFBV9OJAQTz6eT0eh8e09P++u7du0KZksqBhFfe/fdszZF2XB3MinUsSybDgbx3DvvHH917969UQAAQRBCoVDocSKYXWrjPpauVNf9u/RRNe2/KvNjeH5sTQuQaRJ8ImFmsznXIcDm5uZcp4LN5lwfYJLpVHz4jssuhs0+iyLo+92PpQp3XNA1YZfC/SMfXFLm+52TRevI8WOv1EVZsK4r4nLzS8n+lD6l/xz9HUoIYo9JHffMAAAAAElFTkSuQmCC" alt="bittensor"></div>
 <div class="sact"><button onclick="createNewChat()"><span class="i i14"><svg viewBox="0 0 24 24"><line x1="12" y1="5" x2="12" y2="19"/><line x1="5" y1="12" x2="19" y2="12"/></svg></span> New Chat</button></div>
 <div class="sch-box" style="padding:0 10px 8px"><input type="text" id="sch" placeholder="Search chats..." oninput="filterChats(this.value)" style="width:100%;padding:8px 12px;border-radius:8px;background:var(--bg2);color:var(--t1);border:1px solid var(--bd);font-family:inherit;font-size:13px;outline:none"></div>
 <div class="folder-filter-row">
  <button class="folder-filter-btn" id="folder-filter-btn" onclick="toggleFolderDrop(event)" title="Filter by folder">
   <span class="fdot" id="folder-filter-dot" style="width:7px;height:7px;border-radius:50%;background:var(--t3);flex-shrink:0"></span>
   <span class="fflabel" id="folder-filter-label">All chats</span>
   <span class="ffarrow">▼</span>
  </button>
 </div>
 <div class="folder-drop" id="folder-drop"></div>
 <div id="chatlist"></div>
 <div class="sf"><button class="sf-del" onclick="deleteAllChats()" title="Delete All Chats"><span class="i i12"><svg viewBox="0 0 24 24"><polyline points="3 6 5 6 21 6"/><path d="M19 6v14a2 2 0 01-2 2H7a2 2 0 01-2-2V6m3 0V4a2 2 0 012-2h4a2 2 0 012 2v2"/></svg></span></button><span class="sf-powered" id="footer-powered">Powered by Chutes<img class="sf-chutes-icon" id="footer-icon" src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAJEAAAB4CAYAAADykxBpAAAeVElEQVR4nO2dWWwc17nn/+dUVe8L2WSzuYlqU6QkLpKl0HIc2xPKufaMPGMj1wYY3zhvedC8BA6QPAYwxScnQTCAgQBZkMB5GGSSCMjEy418DV+byliRJZMytbUkbiLFXU2x2Qt7q6pz5qHqNIstSqZly6Sk+gEFgtVdp05VfXW+73zLacDGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsZmq0A2uwObCOnt7SX9/f304MGDENut6O/vL/2NxWL86NGjHIDYHmjoZndgEyHvvPOOFI1GZQCfuVm/d/36daWrq0vCg/0S2tjYbBjOOeGcU865ZN0o/WKDMCEElFKUt2uei/b09EgwRvr7frSSN7sDdxPOOYHxECVzK30kSZLOOQfnd2bSSJIESZIAQLHuHxkZ4ceOHcP169dJe3s7i8ViOgCG+9h2uq9tIkKIMHx1AKpl03Rd/0JtM8agqirK2lVbW1vVuro6raamRo3FYhrucwEC7qOh1hx1gLXXxE1BWg8HgIbW1tbwoUOH6PPPP48nnngCHo+HmZ+LF4wBwMLCAj127Bj+/Oc/4/3339c0TZsHMIPPEBBCCBhjBAA5cuQIAKCvr08cc18I132hzkwBojBUllWIdADaLQ4LUUqfrKysPOB2uxVJkoRqU7GqAiGOp5QqLpcLgUAAlZWVK/F4/ASAlLmtC6UUsiwDgDQyMiJ9/etfx7Fjx9De3s7vJzV3XwiRCYfxUEjZPmruL2dHZ2fnf4lEIv8jGAw6XC4XTGO7YB4j7o0Kw4ZyejweVFVVobGxMVkoFHgqlRoDcO52nWKMAQBPpVJsamoKhUIB4XD4vvIv3atCRKwGscX2uUlYBgYGgrW1tZGGhoaSAZzJZOjPfvazRxYXFx8OhUL1u3fvRjQahdvtvuUJQ6EQmpub0dHRgXw+X93a2rp/+/btV37+85/rWHU60rffftv53nvv6W+++aY2NTWVYYzFCSFZGKOitc8AIFQdbqN2tzz3mhARALS9vV165ZVXyLPPPouWlhZwzou3egg7d+5sUBTlmwC2i30ej4d8+9vfbr127dp2n8+HaDSKmpqazzz5tm3bcODAAYTDYWiatrO5ufm/AdgBQ3h1AEooFFJ8Ph8LBAI6gFEAHwG4Ym1HkiTIsox8Pk9gzO4Y51y/VwXpXhOiEqFQiGcyGbS2tgKG+ll3uuX3+3cxxp4G8IjYRynlbW1tSlNTU4WiKHC5XMJ2uS1+vx9tbW2IRqPQdT3icrm6ARyAYTdpAJwej4cGAgFEIhEej8c/uX79+g2UCZEYhXCfzI63+uyMACj5ciil3Pq/lWQyWRUIBPxYa2uEAPwrgO/BGDHKYSgzvBljRJznNvvECHIT8/PzOH78OD766CNcvXp13OFw/O+enp6/Pfzww0t//etf8cc//hGXLl0iMAzyJeuxnHNy5MgRcuTIEQ7cOypuq45EFABpb2+Xenp6cPToUfT09MDhcBQJIcjn8zcd4Ha7WwB8DUDQstsDYwSqv815gDIj1zSGrZQbwrd8+SKRCNra2iBJEnbt2lW/bdu2R5944gmSSCSyiqIQr9cLAMsABgGcth7b39/v7Orqov39/ezgwYOMc64DYFtdmLaqEAEAaWpqIktLS9i3bx8AIJ/PU7/fzyRJgtVZmEgkKhRFeRjAtwE0lrVTAcB5m/OsGVE2GA65pRARQrBt2zZUVVUhn887q6ur9wSDwW2pVAper5eEQiFUVFRMLS8v6zDUXFIc6/f76fDwsFJfX6/CEFrhcd/SQrTV1Vk58mOPPRa4evWqvHv3bnb48GEcPHgQ9fX1uwG8DODfAFSuc5wwfFd3MEYZY5QYQGy3g3MOxhg45yCEMEmSSn6e9dQgVn1XSCaT+OCDD/D+++/j/PnziUwm86edO3f+n9/85jcXg8FgjhCSsxyznktiy7JlRiLOOYVxA7nD4dA559C0m/yEro8//nivz+drDQaDEiEELpcLMNTVHqwvQDDbVcvORwBw8+EDwG2FSAiQpmlgjHFJknQAmvkXWH/EcMIUomAwiObmZhw4cADV1dWVfr9/76OPPpqQZfmhxcXFiwCGYAj6PSVAwBYRIkuglAKALMu6eGBlRAB8IxQKPRMMBt0Oh0OoHwdubfcI1qgtSZKIJVTymaOQwCJwor9c+HqwqoKsYQ0O069VU1ODRx55BB0dHfD5fC3V1dUBQsg+RVH+/cSJE1NPPPHE9Q11YouxJYTINBx1rKocBwA3APW5557D4cOH8fzzz6Ozs7ONc94ViUS+2djYKFVVVcE0VGEey4rFIgcAWZZBKRXCCawznd6o4IjvWiL3BKsJaxuyowghvK6uTq2rqxMzu4i5PQRg4fHHH7/IOc9YDzHbzxNCChvu6CbwVQtR6YH29vaSpaUlqaWlBa+88kq5s7ACwF5Kadjr9RJFUbiu63j11Vf3nTt3brfL5ZLa29vR0tICRSkNMBIAjRBiHb6EirSmgXwmIkXE3DillJXZO6LdzwOB8XKU4wOwC8CTAGqwOpoBwA0Ao5zzaUJISc1xziUYweUtofo2YySi7e3tUiwWE6mpBDcHSrcD+JeampquQCAAWZZBCOFPPfVU5Y4dO+oBw8YIhULlbcvm6CNGB2LZNoQQHsYYGGOcEFKaaluvAesLxJ3SCOBbAB61dgXApwDyMLIFrEi4hXN1M/iqhYgD0GOxmB6LxUh7e7tSLBbJD3/4Qw0wbKNQKBRIJpOdVVVVj0ej0W+Gw2F4vV4QQhAOhwvhcFgCAFVV11VHkqlv7hQxS7MIYXlCW8nZeSf+G6sPyrSrwBirJ4REKKUlV4Su65wQIjHGzsuyXO6Rp1jNk9p0vlIhEg9IRLZjsVhRfMY5DwPYMTg42PTGG288NjU1tTMSiWD//v1obW0VAiNuclFRlNLTYIxRABI1nvzndlsIoeCcC6+42NYIj6ZpWFhYwOLiIhhjibq6urHa2to1xrCu65JVUDjnVASIGWNEVVWsrKy4Z2dntyUSiSan00l9Pp/s9/vliooKYXOhUCiQkZGRKIDA7Ozsmmt6++23aT6f3xKqDNgEIZJlGbquozyzUFXVVkVR/nttbe2eF1988aGZmZkqn8+HSCSCioqK8qZE/rJACM8d+710XQdjDJIkcdP/IzpYEqSlpSVcuHABFy9eRLFYnNy1a9d7L7zwwqfWdlKplBvm7ExVVaqqqiRJEnc6nZqu64QxhkuXLkU++eSTby0sLDT4fD7q9/vh9/vh8/mE6kY2m8Xs7CwFQAqFtXb14ODgnV7mXeErFSJd1wkAybyhpf2c85Cu63sBfMvpdHZ1dnbS9vZ22XTqFQkhCtYKyJoR4stIupckSbRDNU0jpgoBDFuN5vN5Ojw8jDNnzmBgYCC3uLh44Z133ul/4YUX/mFt680333R7vV6tsbFRT6fTNBaLKdXV1ay7u7vkp3r99dejCwsLfsbYTpfLtcPj8UBRFCiKAofDAUIIVlZWoOv6FOc8PT8/v0Zt9vf3s0wms2W82HdLiMSosGbIJYT46uvr62tra6t37txJenp68Nxzz3EAYUmSugDsppQ6LEKRhZkTDePhSowxCQCllEJsd4JpPHNCCKeUckIII4RQANLKygpZWloimUwGfr9/rra2djqfzxeXl5fJzMwMmZycXDl//vynxWLx8jrT741Mxy8BOFFRUdHo9/vn3W638H7rTqeTAyCZTIYxxgYAzI+Pj68RmHg8zmKx2P0vRF1dXdLg4GB54LIymUx+LRwOH/D5fAGHw8FNe8QNoA3G1L68f8K4JTDsntLs6/P4ecphjDFd15ksy2LmpZnn8cTjcVy4cAFjY2MIBoPDjz322H9GIpH5bDZLk8mkvLy8rBeLxWEYgdQ7ZSidTntcLtcgpZRLksQJIUzTNK6qKkkmk1zX9WuFQmEaZS+jmVp7fwtRT08PSafTtLe3l1iS0lFZWRnw+/0d4XD4UDgcrvV6vcKQZFg/tcKB1ZtFNpLz8zkgZqgFWHUz0GQyyUdHR8knn3yCWCw2v7KyMnTq1Kn/+MlPfjKWSqXI8vIyZYxJO3fu1IeHh7+IE/D6Qw899IHP53MCQDAY5ACQyWSQTqehqiocDkcxlUrlcXMoZEvlZt81m8jv99+kZ/bt2+f0eDzVVVVVzZ2dnc4dO3agtrbW+hVdMyjNtiRJIl901BEIH5BpAxFJkkqhlmQy6VheXsbo6ChisVh8bGxs8fLly5djsdinAC799re/XdngadxOp7O2pqamyu/3U7fbDY/HU6pGyeVyWF5ezoyNjV0dHR29ZZL/Z13KHR53V7hbIxGmpqbwox/9CH19faX9hw4d4k6nk/l8Pn3Xrl3Ytm1b+aGUEEJFuOLLEh6BruvQNA2SJFk93eCcY2xsDOfOncP58+e169evn1tYWDg5Pz8/CuA8gI0KEAD4CoXCgXQ6fUCWZackSRwwYnWMMWSzWWQymTkA/wHgzJd2cZvIXROiiYmJm96Wl19+Weec5yilSY/H47HEvQBTnUiSJJvlO+TLFCCBGImszM3NYWhoCB988AEuXrw4u7S09Iksy2/KsjwFQ4A+T3pGYyQSeczv9z8fCATcXq8XZqCYMMagKAokSbqqKEpqenp6DJZ8onuVu6bOotHoTfsaGxvXROthDMsaY0w4HV2UUiPC+eUJEGeMgVJKzAR5IdxpAAkA6tzcHB0dHSWmD+icpmnnAJzF2vQRBUa6rXf79u3OZ555hj799NPkpZdeysKSQPbjH/+4a2pq6mu6rrcEAgG43W4oiiLq9pHL5ZBKpSrdbvdEW1vb9Pe+970r9fX1ok/K3//+d1d/fz9OnjzJ5+bmkouLi4vJZDKJtSpMeKy3hFq7m36i9aRAgpGyWiG+wxiTNE37QtP126FpGuecM865bE6jhSNxCkYOz0Qul8PS0pIUj8eJpmlzAMZRln9k/t8A4IDb7W6srKwkjY2NHMaUXszs6AsvvNAxMjLSwhgrFQBYE96KxSIymYy7oqKifefOnYeqqqr2YjWPyN3Q0OBXFMWjaZqjWCyeVVX1I6zmGgEAurq6pObmZnb06NEHL+xhofQGmb4ZBTcXHn4piAAq51xcq6hPmwVwAsA/crkcSyQSUi6XI42NjdzpdCZfe+01x0svvVS0FAXIAFr8fv8z1dXVeysqKuDz+QBDgFTzc6mtrS3Q1NRUZVk1ZE1hgVnDz10u1/ZAIOB3Op05sw0dgLtYLHoZYx5CiMwYc2uaJgokSwLj8/mk8fFx4AGInVlTFwghhBeLxTwhZI4QcpkQ0ihJkosQ4nI4HKWQhWmzWNMhPjeWACkkSRKOPHDOaT6fp9lsVslms2oul7u2a9eui9Zjs9ksGGP4zne+QwHQoaEhd3Nzs/unP/1p4/Hjxx+WZfmR1tbW7U1NTevWqlVVVaGqqgpYW45dfi0Ehmpck4ZgztxQLBbBOc9yzh1Y9c5v2VzruyVEfHBwkHV1dXFRJ88559lsdhnAWUqp5HK5mgC0mpsf5o0mhDDGGP8iwVQRBzMT0xiMB6qsrKzIo6OjGB8fx8jIiPfKlSs3tW+JUzEAqK2t3eHxePZ8//vfb9mxY8eBhYWFhsbGRrS3t5e7J8opwvBzbega5ufnMTo6ik8//RSJROICgGFCyCBjbAEA6+7ulg4ePMj6+vpYJpPRm5ub2VaJod01Ierq6tIJIdySO82Wl5cT9fX1A8VicRxAFMB/BRCGRYgYYyLF4Y4Dqowx6LoOSim3ROXZ3NwcBgYGcOrUKQwPDyujo6Pe27UzMDDgqays3EMI+deGhoZdTz/9dCSdTstutxt+v/92xr/IANhQWsrU1BQ+/PBDnD59GuPj47GVlZWjhJCPACy6XK65TCbD3G63w1w3kj333HO61Ym72dwVIbLUxgsoAN7Q0JCFEQ+bjsfji5IkdUqSRJxOp0uknlJKhX1wx30zjXQuy7LIJlSSyaQ8MjLCzp49Wzx16lTh2rVrM8lkMnf48GHlu9/9Lk+n047f//73juXlZbz44ovs5Zdf5tXV1S0A9gN43OVy1Tc1NcG8riIMO8YBgJqzP3G94qVxmYltRUJIkTHGRc640+nklNIiAD2RSDhOnz5NPvzwQ5w4cSIzPj7+oa7rfwcwYL2mXC5Xkti+vr4tkwYCfEWGdX9/PxkeHl7z2obD4dQvf/lLVFdXB3bs2IFwOIxQKAS/36/CECJRirOhch6Rw2MKj1hcSslms9Ls7CwmJydx9uzZ9PT09MXp6enhZDI5CGDkd7/7nXrq1KmKRCLR4fF4Wvbu3atUVVVpZuppHYCHsbYIgMAQIAZAZ8aJRWYmZYx5xCxTVVUG4DLn/ApjLGsuikUopcTpdBYBFBljUiqVIlNTUxgZGclhndp9wAi6dnR0bJnRx8pXJUSYm5sr3+196623XI2NjcqePXvQ2dmJ3bt3c7/fX1q3R6SpflayovieKWgcq7Xx9Pr163RoaAhDQ0MYGRm5Nj8//0EikfhPGCmn44wxTExMVAYCgUd9Pt+/VFdX+/x+v1CDDlgWgrAgkuMopZSbKpjDEJCSn0KW5QTn/JSu6+9SSpckSYKmaUSSJNn8foFSSnK5HFlZWYHZ51kAmbLzoaOjQz969OiWGoEEd12Ijhw5glOnTvFcLifq6AkA1NTU+C9fvly8cePGvKqq/oqKCtrS0kJhRPQ1mA5JU1VsqLDQRFRJSKlUShkbG+NDQ0Ps9OnTK3Nzcxey2ew/AfSb36UA5Iceemi3x+M5UFNTc7Curs5bUVEBp9MJsx8cAFNVlXDOIcsyIYTI1v6YAkewKlwMAJMk6RKA07Isf0AIWbb0lcLIjr1pii7atZQzcQDYKj6h9fgqhIg/8sgjbHBwUBjZCgDq9XrVmZmZ8Ww2219fXz/lcDiaXS5XCywGqa7rnBmIeNoarLnOkiRpAFAsFuVsNqukUinMzMxgeHh4eXZ2dnRmZmZ8enr6RCqVugQAvb29jl/96le7AoFAdPv27Y/X1tZ2bt++3btv3z7s2rULHo/Hen90ACpjTFJVVaSjMEqpLkmSWBBLfLcIYBjAJIzY2DmrAJn9XXdEEclxnHM8+eSTTjNvyFpKtSW560JkvlnCxim9XalUKqGq6pnFxcWFYDAYjUajz3i93joAq4Vkun7LmZrViWcmlqkAWKFQIHNzc/LVq1dx+fJljI+PTywuLr5348aNk6lUagrAAgD8+te/bnU6nc9UVlZ+o76+Prpv376mzs5ONDU1CT+P9VxiyWJqRv7BOdcIIUUYL4VV3y7CsGv6YQjS2EbvlXBPAEBNTY16/PjxLRPauB1fhU1kvRFCENjS0lIaxhs78tprr9VFIpGgJEltAPaKA81Ie3nNVwmLHVRas1HXdUkklZ07dy4Zj8fPpVKp/oWFhX6YagYACYfDOz0ez+N1dXXPNjY2Ojs6OujevXvh8/mKZluyeKCUUkgG1tMLR2J56dAVAP8EcAxG8FY4Oh0wasXKwylrb5b5Ymxl9VXOZoQ9OMxVXSVJ4owxdHZ2Ti8vLw/BqAZNAqgGEKGUBszvi5GHEUIoIQRmSqS40TIAFwD4fD5QShfT6fS16enpiZmZmY9nZmZiMAobQQjB/v373c3NzcFgMFjb1NTk2bNnD5qbm4UAZWEIh2ym4kqmyiIwEiJh9kHBaiJdHsaoMwdjFDpPCCnlCnHOFbN/3Ay/CHVGJiYmMDExgT/84Q/o7+/H5OSkUF9b0ohej02pOxP/iPhSNpsFjDf4HQAXYKwz9A0YwVoA4JqmgXPOKKVElmVieqKLWF3cU9TxIxwOD1dUVLy3tLQUGx0dHQZwHVh9y5ubm7Xu7u5CRUWFWltbi6amJlRWVgKm8JhtyeZMSwgQN1NqOSGElmVZzgP4EMApGOprsuy6dUtfRVYiB0Ci0SgmJiYQjUYRjUYxOTl5T6iwLYE5mshmSbD4eQMP57yec/4/OefHOecq51zjnDNd17mmaXld1zVuoJnxpZwZpeecc67r+mI+n//l8PDwY+Fw2Id1KlU55zQejz+/tLT0f5eXl/VCocAZKzWR55wXrW2aMF3Xi2Yfyj7i73LO/41zHuScu0TaLedcsVzfmgUk7ic2bUGHV199lf7iF79w6rouA1gxa+izALLHjh1b8vv9ciQSkSsrK1FZWSmm0WLxJ3DOJUJIablXxtgNAAuU0hGn0znQ2tp6MR6PZ4Rv5tlnn3XE43E2ODioA0B1dbVQXWkYqkZWVZUKv5QlNUXkI5X+N/uRhLFc3g0YI9BZQkgpwcy0gZwAVFMN37ejy6auCjIyMsI1TRPGbonXX3/d0dra6m1vb0dHRwcefvhhBAKBIozcHRWGbeEu8x1dBPD/YKjFIZ/PlwYgFvSk8Xic+Xw+64MUqkXYH8xc+MpaRs1Nj7SIhQlVkwVwGcaSecMwSoDK6+VF4PeeXRV2o2yaEPX19fFDhw7puVyuXIhcJ0+edF67dk2dmZlBoVBAdXU1AoGAdSWOcrUwTykdAPA2jGV/V7LZLCRJQj6fdwIgg4OD1oUexUohCszYGgAqyzIRo44FqxCJtI4sgAkYdtDHAHIwvcymKuOEEI0bC0Hc92zqSPTuu+8yGKU7pL+/X2poaJB+8IMfeN57773MysrKFV3XA3V1dYHt27dXhUIhJRQKVYoEe/NBL8Go/boAw7E3ZCkmpG+88YZ7aGhIWVlZ0WHMoEoMDAw4FEXxeTweXygUQjAYhCzLMBPrS8YtpVQHwHVddxSLRSmfz0s3btwIFAqFfDwen3zqqacWRJucc+eNGzecS0tLBQCF+30EEmymEAl1AgCIx+M0nU5LLpdLBTCmadr78/Pzo7Ozs3vOnj37uKIotfv370d9fSkWmoWRNnoGhgo7W1aNyuPxOLtw4QIrr2U/cuQIMpmM4vV6vdu2bUNbWxva29tRWVmpwZhFWXOYiwB4oVBwzM/PS9PT0xgeHlamp6fdJ0+eLBcSfvr0aTxobPZKaQwwPM6HDx/mdXV17K233ioAGAEwl8vl/HNzczcARAqFQq3P57MK0QwMp947AKZhqUYVmZSpVKowMTHBCoXCGi9hX18fOjs79YqKCrWpqQn5fB7BYBCVlZXlNhNgqk5VVfnMzAwGBwdx5swZTE9P0ytXrvjKrkf929/+ZvUDPRBsthCteZMnJiYAQKeUpgghKVVVMTc3F1pZWWlljDW2trYqu3fvRiQS4TBGoAEAn4hYlDmdFk69vLlfb2lpEaco2VIXLlxIBQKBiWQyOREIBAKRSMQfDAaVUCikuFwua+aACzAM9JWVlcLExET68uXLi/Pz8wswVLGo0tUJIaylpUXfv3+/LUSbQSKRYKbDD8BqbGxycnI8GAz+w+fz5dPptGymvnJKaQzApbJgZimL0bKPjY6OAgDEr08fP35cBzCbSqVOzM3NJaampnaeOXPmgKqqDa2trdixY8dNq7B5vV74/f4ph8MxMD8/P76wsHDB5XLFYXFGmufSR0dHHwhbSLBlhOjo0aOst7e3lIhmYSmZTJ5Mp9PDhUKB8tWFqJIwPdECs7KjgLXrGZbq1vv7+6nb7ZZgpHjMADieTCY/nZycfLxQKATn5+cbEokEJEm6SYiKxSIaGhpinZ2d/55IJIZUVU2pqpqAmb+N1VH1gZiR3QtQAFJ3d7f8WVmNf/nLX6w/4Htbj3B3d7fc3d3tAlCeLbkvHA7/r3379k0cPnw48ac//SkxNzeX4JwndF1P6LqeUFX1Iue8l3PeYj3ePK9LeKYfRLbMSFQGaW9vl9xuNyGEaOv81kaJnp4eEcG/adX8cmpqavjFixeFMQ+gFE+7Go/H/wlAisfjvkwmwy2LcImlbBZh/BbHtbJmrfGwB5KtKkQsFovxWCz2mV8khOiccwbc9vdeAayqzL6+vnKVmQbwSSKRmE4mk0o2my0t6mmu7gZFUXIAZt1udxFYk0kJrGZA2mxRCADS29tLe3t7qak+6EbU10bahbESyeeq/e/t7X0gfu9+o9wLP9pGurq65ImJCUc0GnVMTEw4YIQqZBhFkXf6MElXV5fc3d3t+JzrAJBTp04pXV1dImXkgedeuAl8cHBQz2azqtfrVaPRaOm37fHFfguMDw4O6jU1NaoodtzocX6/XzOzAWwVZmPzZWDrdQPa1dUl+Xw+KRqN4uDBgzh48GD5GksiK5PdqlrjQeVeUGdfBdxUT1o0Gi1tWC2CFKXdD/RU3sbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbmK+T/AxtDWfnfMI91AAAAAElFTkSuQmCC" alt=""></span></div>
</div>
<div id="main">
 <div class="wm"><img src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAJEAAAB4CAYAAADykxBpAAAeVElEQVR4nO2dWWwc17nn/+dUVe8L2WSzuYlqU6QkLpKl0HIc2xPKufaMPGMj1wYY3zhvedC8BA6QPAYwxScnQTCAgQBZkMB5GGSSCMjEy418DV+byliRJZMytbUkbiLFXU2x2Qt7q6pz5qHqNIstSqZly6Sk+gEFgtVdp05VfXW+73zLacDGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsZmq0A2uwObCOnt7SX9/f304MGDENut6O/vL/2NxWL86NGjHIDYHmjoZndgEyHvvPOOFI1GZQCfuVm/d/36daWrq0vCg/0S2tjYbBjOOeGcU865ZN0o/WKDMCEElFKUt2uei/b09EgwRvr7frSSN7sDdxPOOYHxECVzK30kSZLOOQfnd2bSSJIESZIAQLHuHxkZ4ceOHcP169dJe3s7i8ViOgCG+9h2uq9tIkKIMHx1AKpl03Rd/0JtM8agqirK2lVbW1vVuro6raamRo3FYhrucwEC7qOh1hx1gLXXxE1BWg8HgIbW1tbwoUOH6PPPP48nnngCHo+HmZ+LF4wBwMLCAj127Bj+/Oc/4/3339c0TZsHMIPPEBBCCBhjBAA5cuQIAKCvr08cc18I132hzkwBojBUllWIdADaLQ4LUUqfrKysPOB2uxVJkoRqU7GqAiGOp5QqLpcLgUAAlZWVK/F4/ASAlLmtC6UUsiwDgDQyMiJ9/etfx7Fjx9De3s7vJzV3XwiRCYfxUEjZPmruL2dHZ2fnf4lEIv8jGAw6XC4XTGO7YB4j7o0Kw4ZyejweVFVVobGxMVkoFHgqlRoDcO52nWKMAQBPpVJsamoKhUIB4XD4vvIv3atCRKwGscX2uUlYBgYGgrW1tZGGhoaSAZzJZOjPfvazRxYXFx8OhUL1u3fvRjQahdvtvuUJQ6EQmpub0dHRgXw+X93a2rp/+/btV37+85/rWHU60rffftv53nvv6W+++aY2NTWVYYzFCSFZGKOitc8AIFQdbqN2tzz3mhARALS9vV165ZVXyLPPPouWlhZwzou3egg7d+5sUBTlmwC2i30ej4d8+9vfbr127dp2n8+HaDSKmpqazzz5tm3bcODAAYTDYWiatrO5ufm/AdgBQ3h1AEooFFJ8Ph8LBAI6gFEAHwG4Ym1HkiTIsox8Pk9gzO4Y51y/VwXpXhOiEqFQiGcyGbS2tgKG+ll3uuX3+3cxxp4G8IjYRynlbW1tSlNTU4WiKHC5XMJ2uS1+vx9tbW2IRqPQdT3icrm6ARyAYTdpAJwej4cGAgFEIhEej8c/uX79+g2UCZEYhXCfzI63+uyMACj5ciil3Pq/lWQyWRUIBPxYa2uEAPwrgO/BGDHKYSgzvBljRJznNvvECHIT8/PzOH78OD766CNcvXp13OFw/O+enp6/Pfzww0t//etf8cc//hGXLl0iMAzyJeuxnHNy5MgRcuTIEQ7cOypuq45EFABpb2+Xenp6cPToUfT09MDhcBQJIcjn8zcd4Ha7WwB8DUDQstsDYwSqv815gDIj1zSGrZQbwrd8+SKRCNra2iBJEnbt2lW/bdu2R5944gmSSCSyiqIQr9cLAMsABgGcth7b39/v7Orqov39/ezgwYOMc64DYFtdmLaqEAEAaWpqIktLS9i3bx8AIJ/PU7/fzyRJgtVZmEgkKhRFeRjAtwE0lrVTAcB5m/OsGVE2GA65pRARQrBt2zZUVVUhn887q6ur9wSDwW2pVAper5eEQiFUVFRMLS8v6zDUXFIc6/f76fDwsFJfX6/CEFrhcd/SQrTV1Vk58mOPPRa4evWqvHv3bnb48GEcPHgQ9fX1uwG8DODfAFSuc5wwfFd3MEYZY5QYQGy3g3MOxhg45yCEMEmSSn6e9dQgVn1XSCaT+OCDD/D+++/j/PnziUwm86edO3f+n9/85jcXg8FgjhCSsxyznktiy7JlRiLOOYVxA7nD4dA559C0m/yEro8//nivz+drDQaDEiEELpcLMNTVHqwvQDDbVcvORwBw8+EDwG2FSAiQpmlgjHFJknQAmvkXWH/EcMIUomAwiObmZhw4cADV1dWVfr9/76OPPpqQZfmhxcXFiwCGYAj6PSVAwBYRIkuglAKALMu6eGBlRAB8IxQKPRMMBt0Oh0OoHwdubfcI1qgtSZKIJVTymaOQwCJwor9c+HqwqoKsYQ0O069VU1ODRx55BB0dHfD5fC3V1dUBQsg+RVH+/cSJE1NPPPHE9Q11YouxJYTINBx1rKocBwA3APW5557D4cOH8fzzz6Ozs7ONc94ViUS+2djYKFVVVcE0VGEey4rFIgcAWZZBKRXCCawznd6o4IjvWiL3BKsJaxuyowghvK6uTq2rqxMzu4i5PQRg4fHHH7/IOc9YDzHbzxNCChvu6CbwVQtR6YH29vaSpaUlqaWlBa+88kq5s7ACwF5Kadjr9RJFUbiu63j11Vf3nTt3brfL5ZLa29vR0tICRSkNMBIAjRBiHb6EirSmgXwmIkXE3DillJXZO6LdzwOB8XKU4wOwC8CTAGqwOpoBwA0Ao5zzaUJISc1xziUYweUtofo2YySi7e3tUiwWE6mpBDcHSrcD+JeampquQCAAWZZBCOFPPfVU5Y4dO+oBw8YIhULlbcvm6CNGB2LZNoQQHsYYGGOcEFKaaluvAesLxJ3SCOBbAB61dgXApwDyMLIFrEi4hXN1M/iqhYgD0GOxmB6LxUh7e7tSLBbJD3/4Qw0wbKNQKBRIJpOdVVVVj0ej0W+Gw2F4vV4QQhAOhwvhcFgCAFVV11VHkqlv7hQxS7MIYXlCW8nZeSf+G6sPyrSrwBirJ4REKKUlV4Su65wQIjHGzsuyXO6Rp1jNk9p0vlIhEg9IRLZjsVhRfMY5DwPYMTg42PTGG288NjU1tTMSiWD//v1obW0VAiNuclFRlNLTYIxRABI1nvzndlsIoeCcC6+42NYIj6ZpWFhYwOLiIhhjibq6urHa2to1xrCu65JVUDjnVASIGWNEVVWsrKy4Z2dntyUSiSan00l9Pp/s9/vliooKYXOhUCiQkZGRKIDA7Ozsmmt6++23aT6f3xKqDNgEIZJlGbquozyzUFXVVkVR/nttbe2eF1988aGZmZkqn8+HSCSCioqK8qZE/rJACM8d+710XQdjDJIkcdP/IzpYEqSlpSVcuHABFy9eRLFYnNy1a9d7L7zwwqfWdlKplBvm7ExVVaqqqiRJEnc6nZqu64QxhkuXLkU++eSTby0sLDT4fD7q9/vh9/vh8/mE6kY2m8Xs7CwFQAqFtXb14ODgnV7mXeErFSJd1wkAybyhpf2c85Cu63sBfMvpdHZ1dnbS9vZ22XTqFQkhCtYKyJoR4stIupckSbRDNU0jpgoBDFuN5vN5Ojw8jDNnzmBgYCC3uLh44Z133ul/4YUX/mFt680333R7vV6tsbFRT6fTNBaLKdXV1ay7u7vkp3r99dejCwsLfsbYTpfLtcPj8UBRFCiKAofDAUIIVlZWoOv6FOc8PT8/v0Zt9vf3s0wms2W82HdLiMSosGbIJYT46uvr62tra6t37txJenp68Nxzz3EAYUmSugDsppQ6LEKRhZkTDePhSowxCQCllEJsd4JpPHNCCKeUckIII4RQANLKygpZWloimUwGfr9/rra2djqfzxeXl5fJzMwMmZycXDl//vynxWLx8jrT741Mxy8BOFFRUdHo9/vn3W638H7rTqeTAyCZTIYxxgYAzI+Pj68RmHg8zmKx2P0vRF1dXdLg4GB54LIymUx+LRwOH/D5fAGHw8FNe8QNoA3G1L68f8K4JTDsntLs6/P4ecphjDFd15ksy2LmpZnn8cTjcVy4cAFjY2MIBoPDjz322H9GIpH5bDZLk8mkvLy8rBeLxWEYgdQ7ZSidTntcLtcgpZRLksQJIUzTNK6qKkkmk1zX9WuFQmEaZS+jmVp7fwtRT08PSafTtLe3l1iS0lFZWRnw+/0d4XD4UDgcrvV6vcKQZFg/tcKB1ZtFNpLz8zkgZqgFWHUz0GQyyUdHR8knn3yCWCw2v7KyMnTq1Kn/+MlPfjKWSqXI8vIyZYxJO3fu1IeHh7+IE/D6Qw899IHP53MCQDAY5ACQyWSQTqehqiocDkcxlUrlcXMoZEvlZt81m8jv99+kZ/bt2+f0eDzVVVVVzZ2dnc4dO3agtrbW+hVdMyjNtiRJIl901BEIH5BpAxFJkkqhlmQy6VheXsbo6ChisVh8bGxs8fLly5djsdinAC799re/XdngadxOp7O2pqamyu/3U7fbDY/HU6pGyeVyWF5ezoyNjV0dHR29ZZL/Z13KHR53V7hbIxGmpqbwox/9CH19faX9hw4d4k6nk/l8Pn3Xrl3Ytm1b+aGUEEJFuOLLEh6BruvQNA2SJFk93eCcY2xsDOfOncP58+e169evn1tYWDg5Pz8/CuA8gI0KEAD4CoXCgXQ6fUCWZackSRwwYnWMMWSzWWQymTkA/wHgzJd2cZvIXROiiYmJm96Wl19+Weec5yilSY/H47HEvQBTnUiSJJvlO+TLFCCBGImszM3NYWhoCB988AEuXrw4u7S09Iksy2/KsjwFQ4A+T3pGYyQSeczv9z8fCATcXq8XZqCYMMagKAokSbqqKEpqenp6DJZ8onuVu6bOotHoTfsaGxvXROthDMsaY0w4HV2UUiPC+eUJEGeMgVJKzAR5IdxpAAkA6tzcHB0dHSWmD+icpmnnAJzF2vQRBUa6rXf79u3OZ555hj799NPkpZdeysKSQPbjH/+4a2pq6mu6rrcEAgG43W4oiiLq9pHL5ZBKpSrdbvdEW1vb9Pe+970r9fX1ok/K3//+d1d/fz9OnjzJ5+bmkouLi4vJZDKJtSpMeKy3hFq7m36i9aRAgpGyWiG+wxiTNE37QtP126FpGuecM865bE6jhSNxCkYOz0Qul8PS0pIUj8eJpmlzAMZRln9k/t8A4IDb7W6srKwkjY2NHMaUXszs6AsvvNAxMjLSwhgrFQBYE96KxSIymYy7oqKifefOnYeqqqr2YjWPyN3Q0OBXFMWjaZqjWCyeVVX1I6zmGgEAurq6pObmZnb06NEHL+xhofQGmb4ZBTcXHn4piAAq51xcq6hPmwVwAsA/crkcSyQSUi6XI42NjdzpdCZfe+01x0svvVS0FAXIAFr8fv8z1dXVeysqKuDz+QBDgFTzc6mtrS3Q1NRUZVk1ZE1hgVnDz10u1/ZAIOB3Op05sw0dgLtYLHoZYx5CiMwYc2uaJgokSwLj8/mk8fFx4AGInVlTFwghhBeLxTwhZI4QcpkQ0ihJkosQ4nI4HKWQhWmzWNMhPjeWACkkSRKOPHDOaT6fp9lsVslms2oul7u2a9eui9Zjs9ksGGP4zne+QwHQoaEhd3Nzs/unP/1p4/Hjxx+WZfmR1tbW7U1NTevWqlVVVaGqqgpYW45dfi0Ehmpck4ZgztxQLBbBOc9yzh1Y9c5v2VzruyVEfHBwkHV1dXFRJ88559lsdhnAWUqp5HK5mgC0mpsf5o0mhDDGGP8iwVQRBzMT0xiMB6qsrKzIo6OjGB8fx8jIiPfKlSs3tW+JUzEAqK2t3eHxePZ8//vfb9mxY8eBhYWFhsbGRrS3t5e7J8opwvBzbega5ufnMTo6ik8//RSJROICgGFCyCBjbAEA6+7ulg4ePMj6+vpYJpPRm5ub2VaJod01Ierq6tIJIdySO82Wl5cT9fX1A8VicRxAFMB/BRCGRYgYYyLF4Y4Dqowx6LoOSim3ROXZ3NwcBgYGcOrUKQwPDyujo6Pe27UzMDDgqays3EMI+deGhoZdTz/9dCSdTstutxt+v/92xr/IANhQWsrU1BQ+/PBDnD59GuPj47GVlZWjhJCPACy6XK65TCbD3G63w1w3kj333HO61Ym72dwVIbLUxgsoAN7Q0JCFEQ+bjsfji5IkdUqSRJxOp0uknlJKhX1wx30zjXQuy7LIJlSSyaQ8MjLCzp49Wzx16lTh2rVrM8lkMnf48GHlu9/9Lk+n047f//73juXlZbz44ovs5Zdf5tXV1S0A9gN43OVy1Tc1NcG8riIMO8YBgJqzP3G94qVxmYltRUJIkTHGRc640+nklNIiAD2RSDhOnz5NPvzwQ5w4cSIzPj7+oa7rfwcwYL2mXC5Xkti+vr4tkwYCfEWGdX9/PxkeHl7z2obD4dQvf/lLVFdXB3bs2IFwOIxQKAS/36/CECJRirOhch6Rw2MKj1hcSslms9Ls7CwmJydx9uzZ9PT09MXp6enhZDI5CGDkd7/7nXrq1KmKRCLR4fF4Wvbu3atUVVVpZuppHYCHsbYIgMAQIAZAZ8aJRWYmZYx5xCxTVVUG4DLn/ApjLGsuikUopcTpdBYBFBljUiqVIlNTUxgZGclhndp9wAi6dnR0bJnRx8pXJUSYm5sr3+196623XI2NjcqePXvQ2dmJ3bt3c7/fX1q3R6SpflayovieKWgcq7Xx9Pr163RoaAhDQ0MYGRm5Nj8//0EikfhPGCmn44wxTExMVAYCgUd9Pt+/VFdX+/x+v1CDDlgWgrAgkuMopZSbKpjDEJCSn0KW5QTn/JSu6+9SSpckSYKmaUSSJNn8foFSSnK5HFlZWYHZ51kAmbLzoaOjQz969OiWGoEEd12Ijhw5glOnTvFcLifq6AkA1NTU+C9fvly8cePGvKqq/oqKCtrS0kJhRPQ1mA5JU1VsqLDQRFRJSKlUShkbG+NDQ0Ps9OnTK3Nzcxey2ew/AfSb36UA5Iceemi3x+M5UFNTc7Curs5bUVEBp9MJsx8cAFNVlXDOIcsyIYTI1v6YAkewKlwMAJMk6RKA07Isf0AIWbb0lcLIjr1pii7atZQzcQDYKj6h9fgqhIg/8sgjbHBwUBjZCgDq9XrVmZmZ8Ww2219fXz/lcDiaXS5XCywGqa7rnBmIeNoarLnOkiRpAFAsFuVsNqukUinMzMxgeHh4eXZ2dnRmZmZ8enr6RCqVugQAvb29jl/96le7AoFAdPv27Y/X1tZ2bt++3btv3z7s2rULHo/Hen90ACpjTFJVVaSjMEqpLkmSWBBLfLcIYBjAJIzY2DmrAJn9XXdEEclxnHM8+eSTTjNvyFpKtSW560JkvlnCxim9XalUKqGq6pnFxcWFYDAYjUajz3i93joAq4Vkun7LmZrViWcmlqkAWKFQIHNzc/LVq1dx+fJljI+PTywuLr5348aNk6lUagrAAgD8+te/bnU6nc9UVlZ+o76+Prpv376mzs5ONDU1CT+P9VxiyWJqRv7BOdcIIUUYL4VV3y7CsGv6YQjS2EbvlXBPAEBNTY16/PjxLRPauB1fhU1kvRFCENjS0lIaxhs78tprr9VFIpGgJEltAPaKA81Ie3nNVwmLHVRas1HXdUkklZ07dy4Zj8fPpVKp/oWFhX6YagYACYfDOz0ez+N1dXXPNjY2Ojs6OujevXvh8/mKZluyeKCUUkgG1tMLR2J56dAVAP8EcAxG8FY4Oh0wasXKwylrb5b5Ymxl9VXOZoQ9OMxVXSVJ4owxdHZ2Ti8vLw/BqAZNAqgGEKGUBszvi5GHEUIoIQRmSqS40TIAFwD4fD5QShfT6fS16enpiZmZmY9nZmZiMAobQQjB/v373c3NzcFgMFjb1NTk2bNnD5qbm4UAZWEIh2ym4kqmyiIwEiJh9kHBaiJdHsaoMwdjFDpPCCnlCnHOFbN/3Ay/CHVGJiYmMDExgT/84Q/o7+/H5OSkUF9b0ohej02pOxP/iPhSNpsFjDf4HQAXYKwz9A0YwVoA4JqmgXPOKKVElmVieqKLWF3cU9TxIxwOD1dUVLy3tLQUGx0dHQZwHVh9y5ubm7Xu7u5CRUWFWltbi6amJlRWVgKm8JhtyeZMSwgQN1NqOSGElmVZzgP4EMApGOprsuy6dUtfRVYiB0Ci0SgmJiYQjUYRjUYxOTl5T6iwLYE5mshmSbD4eQMP57yec/4/OefHOecq51zjnDNd17mmaXld1zVuoJnxpZwZpeecc67r+mI+n//l8PDwY+Fw2Id1KlU55zQejz+/tLT0f5eXl/VCocAZKzWR55wXrW2aMF3Xi2Yfyj7i73LO/41zHuScu0TaLedcsVzfmgUk7ic2bUGHV199lf7iF79w6rouA1gxa+izALLHjh1b8vv9ciQSkSsrK1FZWSmm0WLxJ3DOJUJIablXxtgNAAuU0hGn0znQ2tp6MR6PZ4Rv5tlnn3XE43E2ODioA0B1dbVQXWkYqkZWVZUKv5QlNUXkI5X+N/uRhLFc3g0YI9BZQkgpwcy0gZwAVFMN37ejy6auCjIyMsI1TRPGbonXX3/d0dra6m1vb0dHRwcefvhhBAKBIozcHRWGbeEu8x1dBPD/YKjFIZ/PlwYgFvSk8Xic+Xw+64MUqkXYH8xc+MpaRs1Nj7SIhQlVkwVwGcaSecMwSoDK6+VF4PeeXRV2o2yaEPX19fFDhw7puVyuXIhcJ0+edF67dk2dmZlBoVBAdXU1AoGAdSWOcrUwTykdAPA2jGV/V7LZLCRJQj6fdwIgg4OD1oUexUohCszYGgAqyzIRo44FqxCJtI4sgAkYdtDHAHIwvcymKuOEEI0bC0Hc92zqSPTuu+8yGKU7pL+/X2poaJB+8IMfeN57773MysrKFV3XA3V1dYHt27dXhUIhJRQKVYoEe/NBL8Go/boAw7E3ZCkmpG+88YZ7aGhIWVlZ0WHMoEoMDAw4FEXxeTweXygUQjAYhCzLMBPrS8YtpVQHwHVddxSLRSmfz0s3btwIFAqFfDwen3zqqacWRJucc+eNGzecS0tLBQCF+30EEmymEAl1AgCIx+M0nU5LLpdLBTCmadr78/Pzo7Ozs3vOnj37uKIotfv370d9fSkWmoWRNnoGhgo7W1aNyuPxOLtw4QIrr2U/cuQIMpmM4vV6vdu2bUNbWxva29tRWVmpwZhFWXOYiwB4oVBwzM/PS9PT0xgeHlamp6fdJ0+eLBcSfvr0aTxobPZKaQwwPM6HDx/mdXV17K233ioAGAEwl8vl/HNzczcARAqFQq3P57MK0QwMp947AKZhqUYVmZSpVKowMTHBCoXCGi9hX18fOjs79YqKCrWpqQn5fB7BYBCVlZXlNhNgqk5VVfnMzAwGBwdx5swZTE9P0ytXrvjKrkf929/+ZvUDPRBsthCteZMnJiYAQKeUpgghKVVVMTc3F1pZWWlljDW2trYqu3fvRiQS4TBGoAEAn4hYlDmdFk69vLlfb2lpEaco2VIXLlxIBQKBiWQyOREIBAKRSMQfDAaVUCikuFwua+aACzAM9JWVlcLExET68uXLi/Pz8wswVLGo0tUJIaylpUXfv3+/LUSbQSKRYKbDD8BqbGxycnI8GAz+w+fz5dPptGymvnJKaQzApbJgZimL0bKPjY6OAgDEr08fP35cBzCbSqVOzM3NJaampnaeOXPmgKqqDa2trdixY8dNq7B5vV74/f4ph8MxMD8/P76wsHDB5XLFYXFGmufSR0dHHwhbSLBlhOjo0aOst7e3lIhmYSmZTJ5Mp9PDhUKB8tWFqJIwPdECs7KjgLXrGZbq1vv7+6nb7ZZgpHjMADieTCY/nZycfLxQKATn5+cbEokEJEm6SYiKxSIaGhpinZ2d/55IJIZUVU2pqpqAmb+N1VH1gZiR3QtQAFJ3d7f8WVmNf/nLX6w/4Htbj3B3d7fc3d3tAlCeLbkvHA7/r3379k0cPnw48ac//SkxNzeX4JwndF1P6LqeUFX1Iue8l3PeYj3ePK9LeKYfRLbMSFQGaW9vl9xuNyGEaOv81kaJnp4eEcG/adX8cmpqavjFixeFMQ+gFE+7Go/H/wlAisfjvkwmwy2LcImlbBZh/BbHtbJmrfGwB5KtKkQsFovxWCz2mV8khOiccwbc9vdeAayqzL6+vnKVmQbwSSKRmE4mk0o2my0t6mmu7gZFUXIAZt1udxFYk0kJrGZA2mxRCADS29tLe3t7qak+6EbU10bahbESyeeq/e/t7X0gfu9+o9wLP9pGurq65ImJCUc0GnVMTEw4YIQqZBhFkXf6MElXV5fc3d3t+JzrAJBTp04pXV1dImXkgedeuAl8cHBQz2azqtfrVaPRaOm37fHFfguMDw4O6jU1NaoodtzocX6/XzOzAWwVZmPzZWDrdQPa1dUl+Xw+KRqN4uDBgzh48GD5GksiK5PdqlrjQeVeUGdfBdxUT1o0Gi1tWC2CFKXdD/RU3sbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbmK+T/AxtDWfnfMI91AAAAAElFTkSuQmCC" alt=""></div>
 <div class="float-footer"><img class="ff-chutes" src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAJEAAAB4CAYAAADykxBpAAAeVElEQVR4nO2dWWwc17nn/+dUVe8L2WSzuYlqU6QkLpKl0HIc2xPKufaMPGMj1wYY3zhvedC8BA6QPAYwxScnQTCAgQBZkMB5GGSSCMjEy418DV+byliRJZMytbUkbiLFXU2x2Qt7q6pz5qHqNIstSqZly6Sk+gEFgtVdp05VfXW+73zLacDGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsZmq0A2uwObCOnt7SX9/f304MGDENut6O/vL/2NxWL86NGjHIDYHmjoZndgEyHvvPOOFI1GZQCfuVm/d/36daWrq0vCg/0S2tjYbBjOOeGcU865ZN0o/WKDMCEElFKUt2uei/b09EgwRvr7frSSN7sDdxPOOYHxECVzK30kSZLOOQfnd2bSSJIESZIAQLHuHxkZ4ceOHcP169dJe3s7i8ViOgCG+9h2uq9tIkKIMHx1AKpl03Rd/0JtM8agqirK2lVbW1vVuro6raamRo3FYhrucwEC7qOh1hx1gLXXxE1BWg8HgIbW1tbwoUOH6PPPP48nnngCHo+HmZ+LF4wBwMLCAj127Bj+/Oc/4/3339c0TZsHMIPPEBBCCBhjBAA5cuQIAKCvr08cc18I132hzkwBojBUllWIdADaLQ4LUUqfrKysPOB2uxVJkoRqU7GqAiGOp5QqLpcLgUAAlZWVK/F4/ASAlLmtC6UUsiwDgDQyMiJ9/etfx7Fjx9De3s7vJzV3XwiRCYfxUEjZPmruL2dHZ2fnf4lEIv8jGAw6XC4XTGO7YB4j7o0Kw4ZyejweVFVVobGxMVkoFHgqlRoDcO52nWKMAQBPpVJsamoKhUIB4XD4vvIv3atCRKwGscX2uUlYBgYGgrW1tZGGhoaSAZzJZOjPfvazRxYXFx8OhUL1u3fvRjQahdvtvuUJQ6EQmpub0dHRgXw+X93a2rp/+/btV37+85/rWHU60rffftv53nvv6W+++aY2NTWVYYzFCSFZGKOitc8AIFQdbqN2tzz3mhARALS9vV165ZVXyLPPPouWlhZwzou3egg7d+5sUBTlmwC2i30ej4d8+9vfbr127dp2n8+HaDSKmpqazzz5tm3bcODAAYTDYWiatrO5ufm/AdgBQ3h1AEooFFJ8Ph8LBAI6gFEAHwG4Ym1HkiTIsox8Pk9gzO4Y51y/VwXpXhOiEqFQiGcyGbS2tgKG+ll3uuX3+3cxxp4G8IjYRynlbW1tSlNTU4WiKHC5XMJ2uS1+vx9tbW2IRqPQdT3icrm6ARyAYTdpAJwej4cGAgFEIhEej8c/uX79+g2UCZEYhXCfzI63+uyMACj5ciil3Pq/lWQyWRUIBPxYa2uEAPwrgO/BGDHKYSgzvBljRJznNvvECHIT8/PzOH78OD766CNcvXp13OFw/O+enp6/Pfzww0t//etf8cc//hGXLl0iMAzyJeuxnHNy5MgRcuTIEQ7cOypuq45EFABpb2+Xenp6cPToUfT09MDhcBQJIcjn8zcd4Ha7WwB8DUDQstsDYwSqv815gDIj1zSGrZQbwrd8+SKRCNra2iBJEnbt2lW/bdu2R5944gmSSCSyiqIQr9cLAMsABgGcth7b39/v7Orqov39/ezgwYOMc64DYFtdmLaqEAEAaWpqIktLS9i3bx8AIJ/PU7/fzyRJgtVZmEgkKhRFeRjAtwE0lrVTAcB5m/OsGVE2GA65pRARQrBt2zZUVVUhn887q6ur9wSDwW2pVAper5eEQiFUVFRMLS8v6zDUXFIc6/f76fDwsFJfX6/CEFrhcd/SQrTV1Vk58mOPPRa4evWqvHv3bnb48GEcPHgQ9fX1uwG8DODfAFSuc5wwfFd3MEYZY5QYQGy3g3MOxhg45yCEMEmSSn6e9dQgVn1XSCaT+OCDD/D+++/j/PnziUwm86edO3f+n9/85jcXg8FgjhCSsxyznktiy7JlRiLOOYVxA7nD4dA559C0m/yEro8//nivz+drDQaDEiEELpcLMNTVHqwvQDDbVcvORwBw8+EDwG2FSAiQpmlgjHFJknQAmvkXWH/EcMIUomAwiObmZhw4cADV1dWVfr9/76OPPpqQZfmhxcXFiwCGYAj6PSVAwBYRIkuglAKALMu6eGBlRAB8IxQKPRMMBt0Oh0OoHwdubfcI1qgtSZKIJVTymaOQwCJwor9c+HqwqoKsYQ0O069VU1ODRx55BB0dHfD5fC3V1dUBQsg+RVH+/cSJE1NPPPHE9Q11YouxJYTINBx1rKocBwA3APW5557D4cOH8fzzz6Ozs7ONc94ViUS+2djYKFVVVcE0VGEey4rFIgcAWZZBKRXCCawznd6o4IjvWiL3BKsJaxuyowghvK6uTq2rqxMzu4i5PQRg4fHHH7/IOc9YDzHbzxNCChvu6CbwVQtR6YH29vaSpaUlqaWlBa+88kq5s7ACwF5Kadjr9RJFUbiu63j11Vf3nTt3brfL5ZLa29vR0tICRSkNMBIAjRBiHb6EirSmgXwmIkXE3DillJXZO6LdzwOB8XKU4wOwC8CTAGqwOpoBwA0Ao5zzaUJISc1xziUYweUtofo2YySi7e3tUiwWE6mpBDcHSrcD+JeampquQCAAWZZBCOFPPfVU5Y4dO+oBw8YIhULlbcvm6CNGB2LZNoQQHsYYGGOcEFKaaluvAesLxJ3SCOBbAB61dgXApwDyMLIFrEi4hXN1M/iqhYgD0GOxmB6LxUh7e7tSLBbJD3/4Qw0wbKNQKBRIJpOdVVVVj0ej0W+Gw2F4vV4QQhAOhwvhcFgCAFVV11VHkqlv7hQxS7MIYXlCW8nZeSf+G6sPyrSrwBirJ4REKKUlV4Su65wQIjHGzsuyXO6Rp1jNk9p0vlIhEg9IRLZjsVhRfMY5DwPYMTg42PTGG288NjU1tTMSiWD//v1obW0VAiNuclFRlNLTYIxRABI1nvzndlsIoeCcC6+42NYIj6ZpWFhYwOLiIhhjibq6urHa2to1xrCu65JVUDjnVASIGWNEVVWsrKy4Z2dntyUSiSan00l9Pp/s9/vliooKYXOhUCiQkZGRKIDA7Ozsmmt6++23aT6f3xKqDNgEIZJlGbquozyzUFXVVkVR/nttbe2eF1988aGZmZkqn8+HSCSCioqK8qZE/rJACM8d+710XQdjDJIkcdP/IzpYEqSlpSVcuHABFy9eRLFYnNy1a9d7L7zwwqfWdlKplBvm7ExVVaqqqiRJEnc6nZqu64QxhkuXLkU++eSTby0sLDT4fD7q9/vh9/vh8/mE6kY2m8Xs7CwFQAqFtXb14ODgnV7mXeErFSJd1wkAybyhpf2c85Cu63sBfMvpdHZ1dnbS9vZ22XTqFQkhCtYKyJoR4stIupckSbRDNU0jpgoBDFuN5vN5Ojw8jDNnzmBgYCC3uLh44Z133ul/4YUX/mFt680333R7vV6tsbFRT6fTNBaLKdXV1ay7u7vkp3r99dejCwsLfsbYTpfLtcPj8UBRFCiKAofDAUIIVlZWoOv6FOc8PT8/v0Zt9vf3s0wms2W82HdLiMSosGbIJYT46uvr62tra6t37txJenp68Nxzz3EAYUmSugDsppQ6LEKRhZkTDePhSowxCQCllEJsd4JpPHNCCKeUckIII4RQANLKygpZWloimUwGfr9/rra2djqfzxeXl5fJzMwMmZycXDl//vynxWLx8jrT741Mxy8BOFFRUdHo9/vn3W638H7rTqeTAyCZTIYxxgYAzI+Pj68RmHg8zmKx2P0vRF1dXdLg4GB54LIymUx+LRwOH/D5fAGHw8FNe8QNoA3G1L68f8K4JTDsntLs6/P4ecphjDFd15ksy2LmpZnn8cTjcVy4cAFjY2MIBoPDjz322H9GIpH5bDZLk8mkvLy8rBeLxWEYgdQ7ZSidTntcLtcgpZRLksQJIUzTNK6qKkkmk1zX9WuFQmEaZS+jmVp7fwtRT08PSafTtLe3l1iS0lFZWRnw+/0d4XD4UDgcrvV6vcKQZFg/tcKB1ZtFNpLz8zkgZqgFWHUz0GQyyUdHR8knn3yCWCw2v7KyMnTq1Kn/+MlPfjKWSqXI8vIyZYxJO3fu1IeHh7+IE/D6Qw899IHP53MCQDAY5ACQyWSQTqehqiocDkcxlUrlcXMoZEvlZt81m8jv99+kZ/bt2+f0eDzVVVVVzZ2dnc4dO3agtrbW+hVdMyjNtiRJIl901BEIH5BpAxFJkkqhlmQy6VheXsbo6ChisVh8bGxs8fLly5djsdinAC799re/XdngadxOp7O2pqamyu/3U7fbDY/HU6pGyeVyWF5ezoyNjV0dHR29ZZL/Z13KHR53V7hbIxGmpqbwox/9CH19faX9hw4d4k6nk/l8Pn3Xrl3Ytm1b+aGUEEJFuOLLEh6BruvQNA2SJFk93eCcY2xsDOfOncP58+e169evn1tYWDg5Pz8/CuA8gI0KEAD4CoXCgXQ6fUCWZackSRwwYnWMMWSzWWQymTkA/wHgzJd2cZvIXROiiYmJm96Wl19+Weec5yilSY/H47HEvQBTnUiSJJvlO+TLFCCBGImszM3NYWhoCB988AEuXrw4u7S09Iksy2/KsjwFQ4A+T3pGYyQSeczv9z8fCATcXq8XZqCYMMagKAokSbqqKEpqenp6DJZ8onuVu6bOotHoTfsaGxvXROthDMsaY0w4HV2UUiPC+eUJEGeMgVJKzAR5IdxpAAkA6tzcHB0dHSWmD+icpmnnAJzF2vQRBUa6rXf79u3OZ555hj799NPkpZdeysKSQPbjH/+4a2pq6mu6rrcEAgG43W4oiiLq9pHL5ZBKpSrdbvdEW1vb9Pe+970r9fX1ok/K3//+d1d/fz9OnjzJ5+bmkouLi4vJZDKJtSpMeKy3hFq7m36i9aRAgpGyWiG+wxiTNE37QtP126FpGuecM865bE6jhSNxCkYOz0Qul8PS0pIUj8eJpmlzAMZRln9k/t8A4IDb7W6srKwkjY2NHMaUXszs6AsvvNAxMjLSwhgrFQBYE96KxSIymYy7oqKifefOnYeqqqr2YjWPyN3Q0OBXFMWjaZqjWCyeVVX1I6zmGgEAurq6pObmZnb06NEHL+xhofQGmb4ZBTcXHn4piAAq51xcq6hPmwVwAsA/crkcSyQSUi6XI42NjdzpdCZfe+01x0svvVS0FAXIAFr8fv8z1dXVeysqKuDz+QBDgFTzc6mtrS3Q1NRUZVk1ZE1hgVnDz10u1/ZAIOB3Op05sw0dgLtYLHoZYx5CiMwYc2uaJgokSwLj8/mk8fFx4AGInVlTFwghhBeLxTwhZI4QcpkQ0ihJkosQ4nI4HKWQhWmzWNMhPjeWACkkSRKOPHDOaT6fp9lsVslms2oul7u2a9eui9Zjs9ksGGP4zne+QwHQoaEhd3Nzs/unP/1p4/Hjxx+WZfmR1tbW7U1NTevWqlVVVaGqqgpYW45dfi0Ehmpck4ZgztxQLBbBOc9yzh1Y9c5v2VzruyVEfHBwkHV1dXFRJ88559lsdhnAWUqp5HK5mgC0mpsf5o0mhDDGGP8iwVQRBzMT0xiMB6qsrKzIo6OjGB8fx8jIiPfKlSs3tW+JUzEAqK2t3eHxePZ8//vfb9mxY8eBhYWFhsbGRrS3t5e7J8opwvBzbega5ufnMTo6ik8//RSJROICgGFCyCBjbAEA6+7ulg4ePMj6+vpYJpPRm5ub2VaJod01Ierq6tIJIdySO82Wl5cT9fX1A8VicRxAFMB/BRCGRYgYYyLF4Y4Dqowx6LoOSim3ROXZ3NwcBgYGcOrUKQwPDyujo6Pe27UzMDDgqays3EMI+deGhoZdTz/9dCSdTstutxt+v/92xr/IANhQWsrU1BQ+/PBDnD59GuPj47GVlZWjhJCPACy6XK65TCbD3G63w1w3kj333HO61Ym72dwVIbLUxgsoAN7Q0JCFEQ+bjsfji5IkdUqSRJxOp0uknlJKhX1wx30zjXQuy7LIJlSSyaQ8MjLCzp49Wzx16lTh2rVrM8lkMnf48GHlu9/9Lk+n047f//73juXlZbz44ovs5Zdf5tXV1S0A9gN43OVy1Tc1NcG8riIMO8YBgJqzP3G94qVxmYltRUJIkTHGRc640+nklNIiAD2RSDhOnz5NPvzwQ5w4cSIzPj7+oa7rfwcwYL2mXC5Xkti+vr4tkwYCfEWGdX9/PxkeHl7z2obD4dQvf/lLVFdXB3bs2IFwOIxQKAS/36/CECJRirOhch6Rw2MKj1hcSslms9Ls7CwmJydx9uzZ9PT09MXp6enhZDI5CGDkd7/7nXrq1KmKRCLR4fF4Wvbu3atUVVVpZuppHYCHsbYIgMAQIAZAZ8aJRWYmZYx5xCxTVVUG4DLn/ApjLGsuikUopcTpdBYBFBljUiqVIlNTUxgZGclhndp9wAi6dnR0bJnRx8pXJUSYm5sr3+196623XI2NjcqePXvQ2dmJ3bt3c7/fX1q3R6SpflayovieKWgcq7Xx9Pr163RoaAhDQ0MYGRm5Nj8//0EikfhPGCmn44wxTExMVAYCgUd9Pt+/VFdX+/x+v1CDDlgWgrAgkuMopZSbKpjDEJCSn0KW5QTn/JSu6+9SSpckSYKmaUSSJNn8foFSSnK5HFlZWYHZ51kAmbLzoaOjQz969OiWGoEEd12Ijhw5glOnTvFcLifq6AkA1NTU+C9fvly8cePGvKqq/oqKCtrS0kJhRPQ1mA5JU1VsqLDQRFRJSKlUShkbG+NDQ0Ps9OnTK3Nzcxey2ew/AfSb36UA5Iceemi3x+M5UFNTc7Curs5bUVEBp9MJsx8cAFNVlXDOIcsyIYTI1v6YAkewKlwMAJMk6RKA07Isf0AIWbb0lcLIjr1pii7atZQzcQDYKj6h9fgqhIg/8sgjbHBwUBjZCgDq9XrVmZmZ8Ww2219fXz/lcDiaXS5XCywGqa7rnBmIeNoarLnOkiRpAFAsFuVsNqukUinMzMxgeHh4eXZ2dnRmZmZ8enr6RCqVugQAvb29jl/96le7AoFAdPv27Y/X1tZ2bt++3btv3z7s2rULHo/Hen90ACpjTFJVVaSjMEqpLkmSWBBLfLcIYBjAJIzY2DmrAJn9XXdEEclxnHM8+eSTTjNvyFpKtSW560JkvlnCxim9XalUKqGq6pnFxcWFYDAYjUajz3i93joAq4Vkun7LmZrViWcmlqkAWKFQIHNzc/LVq1dx+fJljI+PTywuLr5348aNk6lUagrAAgD8+te/bnU6nc9UVlZ+o76+Prpv376mzs5ONDU1CT+P9VxiyWJqRv7BOdcIIUUYL4VV3y7CsGv6YQjS2EbvlXBPAEBNTY16/PjxLRPauB1fhU1kvRFCENjS0lIaxhs78tprr9VFIpGgJEltAPaKA81Ie3nNVwmLHVRas1HXdUkklZ07dy4Zj8fPpVKp/oWFhX6YagYACYfDOz0ez+N1dXXPNjY2Ojs6OujevXvh8/mKZluyeKCUUkgG1tMLR2J56dAVAP8EcAxG8FY4Oh0wasXKwylrb5b5Ymxl9VXOZoQ9OMxVXSVJ4owxdHZ2Ti8vLw/BqAZNAqgGEKGUBszvi5GHEUIoIQRmSqS40TIAFwD4fD5QShfT6fS16enpiZmZmY9nZmZiMAobQQjB/v373c3NzcFgMFjb1NTk2bNnD5qbm4UAZWEIh2ym4kqmyiIwEiJh9kHBaiJdHsaoMwdjFDpPCCnlCnHOFbN/3Ay/CHVGJiYmMDExgT/84Q/o7+/H5OSkUF9b0ohej02pOxP/iPhSNpsFjDf4HQAXYKwz9A0YwVoA4JqmgXPOKKVElmVieqKLWF3cU9TxIxwOD1dUVLy3tLQUGx0dHQZwHVh9y5ubm7Xu7u5CRUWFWltbi6amJlRWVgKm8JhtyeZMSwgQN1NqOSGElmVZzgP4EMApGOprsuy6dUtfRVYiB0Ci0SgmJiYQjUYRjUYxOTl5T6iwLYE5mshmSbD4eQMP57yec/4/OefHOecq51zjnDNd17mmaXld1zVuoJnxpZwZpeecc67r+mI+n//l8PDwY+Fw2Id1KlU55zQejz+/tLT0f5eXl/VCocAZKzWR55wXrW2aMF3Xi2Yfyj7i73LO/41zHuScu0TaLedcsVzfmgUk7ic2bUGHV199lf7iF79w6rouA1gxa+izALLHjh1b8vv9ciQSkSsrK1FZWSmm0WLxJ3DOJUJIablXxtgNAAuU0hGn0znQ2tp6MR6PZ4Rv5tlnn3XE43E2ODioA0B1dbVQXWkYqkZWVZUKv5QlNUXkI5X+N/uRhLFc3g0YI9BZQkgpwcy0gZwAVFMN37ejy6auCjIyMsI1TRPGbonXX3/d0dra6m1vb0dHRwcefvhhBAKBIozcHRWGbeEu8x1dBPD/YKjFIZ/PlwYgFvSk8Xic+Xw+64MUqkXYH8xc+MpaRs1Nj7SIhQlVkwVwGcaSecMwSoDK6+VF4PeeXRV2o2yaEPX19fFDhw7puVyuXIhcJ0+edF67dk2dmZlBoVBAdXU1AoGAdSWOcrUwTykdAPA2jGV/V7LZLCRJQj6fdwIgg4OD1oUexUohCszYGgAqyzIRo44FqxCJtI4sgAkYdtDHAHIwvcymKuOEEI0bC0Hc92zqSPTuu+8yGKU7pL+/X2poaJB+8IMfeN57773MysrKFV3XA3V1dYHt27dXhUIhJRQKVYoEe/NBL8Go/boAw7E3ZCkmpG+88YZ7aGhIWVlZ0WHMoEoMDAw4FEXxeTweXygUQjAYhCzLMBPrS8YtpVQHwHVddxSLRSmfz0s3btwIFAqFfDwen3zqqacWRJucc+eNGzecS0tLBQCF+30EEmymEAl1AgCIx+M0nU5LLpdLBTCmadr78/Pzo7Ozs3vOnj37uKIotfv370d9fSkWmoWRNnoGhgo7W1aNyuPxOLtw4QIrr2U/cuQIMpmM4vV6vdu2bUNbWxva29tRWVmpwZhFWXOYiwB4oVBwzM/PS9PT0xgeHlamp6fdJ0+eLBcSfvr0aTxobPZKaQwwPM6HDx/mdXV17K233ioAGAEwl8vl/HNzczcARAqFQq3P57MK0QwMp947AKZhqUYVmZSpVKowMTHBCoXCGi9hX18fOjs79YqKCrWpqQn5fB7BYBCVlZXlNhNgqk5VVfnMzAwGBwdx5swZTE9P0ytXrvjKrkf929/+ZvUDPRBsthCteZMnJiYAQKeUpgghKVVVMTc3F1pZWWlljDW2trYqu3fvRiQS4TBGoAEAn4hYlDmdFk69vLlfb2lpEaco2VIXLlxIBQKBiWQyOREIBAKRSMQfDAaVUCikuFwua+aACzAM9JWVlcLExET68uXLi/Pz8wswVLGo0tUJIaylpUXfv3+/LUSbQSKRYKbDD8BqbGxycnI8GAz+w+fz5dPptGymvnJKaQzApbJgZimL0bKPjY6OAgDEr08fP35cBzCbSqVOzM3NJaampnaeOXPmgKqqDa2trdixY8dNq7B5vV74/f4ph8MxMD8/P76wsHDB5XLFYXFGmufSR0dHHwhbSLBlhOjo0aOst7e3lIhmYSmZTJ5Mp9PDhUKB8tWFqJIwPdECs7KjgLXrGZbq1vv7+6nb7ZZgpHjMADieTCY/nZycfLxQKATn5+cbEokEJEm6SYiKxSIaGhpinZ2d/55IJIZUVU2pqpqAmb+N1VH1gZiR3QtQAFJ3d7f8WVmNf/nLX6w/4Htbj3B3d7fc3d3tAlCeLbkvHA7/r3379k0cPnw48ac//SkxNzeX4JwndF1P6LqeUFX1Iue8l3PeYj3ePK9LeKYfRLbMSFQGaW9vl9xuNyGEaOv81kaJnp4eEcG/adX8cmpqavjFixeFMQ+gFE+7Go/H/wlAisfjvkwmwy2LcImlbBZh/BbHtbJmrfGwB5KtKkQsFovxWCz2mV8khOiccwbc9vdeAayqzL6+vnKVmQbwSSKRmE4mk0o2my0t6mmu7gZFUXIAZt1udxFYk0kJrGZA2mxRCADS29tLe3t7qak+6EbU10bahbESyeeq/e/t7X0gfu9+o9wLP9pGurq65ImJCUc0GnVMTEw4YIQqZBhFkXf6MElXV5fc3d3t+JzrAJBTp04pXV1dImXkgedeuAl8cHBQz2azqtfrVaPRaOm37fHFfguMDw4O6jU1NaoodtzocX6/XzOzAWwVZmPzZWDrdQPa1dUl+Xw+KRqN4uDBgzh48GD5GksiK5PdqlrjQeVeUGdfBdxUT1o0Gi1tWC2CFKXdD/RU3sbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbGxsbmK+T/AxtDWfnfMI91AAAAAElFTkSuQmCC" alt=""><span class="ff-text">Powered by Chutes</span></div>
 <div id="topbar">
  <button id="tog" onclick="toggleSidebar()"><span class="i i18"><svg viewBox="0 0 24 24"><line x1="3" y1="6" x2="21" y2="6"/><line x1="3" y1="12" x2="21" y2="12"/><line x1="3" y1="18" x2="21" y2="18"/></svg></span></button>
  <div class="float-logos"><div class="fl-spacer"></div><img class="fl-t" src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAACQAAAAkCAYAAADhAJiYAAAF6klEQVR4nMWYwYtcWRXGv3PuK7oqXel0JaY6SYNCFr0waTUQhRbHoRaZTQguREE3AREHSVwpGEigFsMgDEMkf4BkpWAHJAsJIoFmWg0ok+kYk1aysG06YUyqU13Vlfeq3333nuOi3muqK1XdVZlIDtzFu1X3vN/9znn3nfMIe5iqGgBCRAoAzWZzJgiCd4wxbxljTqjqMWNMAQBEZAvAf0VkWVX/FIbh7UOHDj3o9kVEfq977gUDAAjD8FvW2t8751o6pDnnImvtH6Mo+m61WuUunzQqCKkqAcCzZ8++Gcfxn7tv5DvmnHPivRcRURFR770458Q55733vnuNtfbjRqPxTq//oWC6VHnPOScphEvHNsBew3sv2boMLIqiX1ar1SC9F+8JQ0S4c+dOod1u3+wSxA0LMWg457yqOlXVra2t2ysrK5O7QqUy8qNHj8biOL6dkljvffeO1Tn3ysN7r957m0L9tVar7R8YPlUNACCKovkMphfkdVjqz6qqxnH8h2q1yqpqMihKYQwR+Var9dNisfihiCRElEt/A3NH1bW1NbRare3rUUxEMD09jYmJCagqACTMnIui6P3x8fErGQOlMdT19fWZUqn0d2NMICJMRKSqICJsbGzo5cuXdXFxkZ1zIKLM6VCWy+XQaDRw48YNnZubI+ccmFmZ2YsIt1qtr01OTt5VVRMAICKSKIo+MMaMee89M++I6aVLl3Dz5k2emppCLpfrDTVEBMwMopdTgZlRq9Vw+vRpOXXqFIkIjDEAQCJCzMz5fP4qgLcBdOJWr9e/fODAgSUASkScSWyMwfLysp45c4YKhQJE5CVlxsbGUCgUEIYhrLVgZmTKorNbzM7O6tWrV3H8+HHy3u8IeXpym42NjbcPHjy4GABAPp//ETOTiHgAnDlSVRw9epTm5+dhjNkBk+XWtWvXtNls4sKFCzo1NcUZVLapYrGI2dlZIqJtJbtNRJSZdd++fe8CWMSTJ0/2WWtXswOn9/zYyy5evKiHDx/Ws2fP+vX19b4L0nNo4OGpqpokSe3x48eHeP/+/aeDIPh8d7h6dgDn3EvDWgvvPZIkQbFYxN27d/nKlSsqIrDWbv/Pe7/jSe016sRWgiD43MTExDeCIAjmUjmFiEyfBVkS9p0nInjvUSqVcP/+fW632xgfH4eI9E3yfpaGDblc7uvMzCeGWrWHqSq897DWvspaAgBjzBeZiI6l86OVBANsWFUGrDnG3vuJ1wHyOoyZC5zL5UbX+P9kqiosIhvpxZvmgYg0WVVXACCrmd+QaQr0H3bOLaWTryWpP4up6iccRdFfvPcxM/ObChsRGe+9WGs/4nK5/Mh7/wk6sr16i4JOHo66KVUVIlIR+efS0tI9BgBr7a8BEDO/skTMjCRJYK2FiIyyVACQtfa3lUrFsapSvV7/jXPuOQBW1aG8ZUqcPHlS2+028vk8ms0m1tbWdNiKUlWVmdl7H7bb7V9lk1kt/bP0jZ8M2eKo914bjYacO3fOl8tlPXLkiFYqFbewsCBRFGl3gzDAR6KqGobhLzKWrOKnhw8fBjMzMx/ncrlZEfH9XrT9VGJmNJtNvX79Oj148ABxHGNubk7Onz9PhUKhU7j3eZ2IiBhj2Dm3Uq/Xv1QulyOkj/92b/T06dOvOOe20tLopdpokFK7dRi7NZDOOanX62+lDDsFyCZardb3Un9uFKgkSbb7ryRJdoVR1URVdXNz88d9YbqgAgAIw/DdbJfOuc/ctXbBOFWVFObn3ffcLS+CVKnvOOc2s0QfVq0BID5LYOfcVrPZ/MFQML1Qz58/PxHH8UeZWqN8cOj3oSGO47/VarWvjgTTBbUd1xcvXvzQWvuv3oTNPst0D1V1vc2BtfbfYRj+pOurx55P8CAozsrLW7dujYVh+O04jueTJHk88NFKLUmST621v2u1Wt+/d+/eeLfP3e451Bt+YWEhqFQqLrteXV0tlUqlWWPMCWb+gqpOotNANAGsOueWNzc3/zE9Pb3eBRIA8HuVOf8Dz+oGAsYI398AAAAASUVORK5CYII=" alt=""><img class="fl-text" src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAAE4AAAAWCAYAAABud6qHAAALrElEQVR4nO1Ye2xb5RU/597rmziJ49hJaBO39JU2WVra0KQDViANezEBGxNTB4itmzYqRjc6VKCdJi3t0BB0QmIgbRTGxooY4JRRha4tdJtLk+bhOA+7eTixGzt183AcO37fa9/H2R+2WQgpsE2apo0jXdn3+373nO+c73zfeQB8Sv8eERF+QigCfGLspwTwQeP+E4b+nyQOAMBsNrMAAENDQ3To0CH1SmBEJADgGhsbARHl/9Aa/3tpZMT3VCymuM6d69kOAEBE7ML5nHe1tdmuHh72OpzOSwP9/aMmIsL/V89jAAA0Gq5Sq2WqOC5PS0Tocrk4IuIsFguXNUwWx9ZUVa26Zu3alVt4XrMOAMDlcvEWiyWH5YiIAch4scWS4bH4WYgDACAiJju2UOb7lPsOAKC5uZmxWK6MXYxZLGshZTeezeHM5g86TE52zpFyOgFARqbLNfVSLEZKR0ffF64ggCEiNJvNbE+P4wGr1b7HbDbzV9qNfybQfIRSHxr/KAN8AgzzUe8LeS21/sV4LvuLLAsMx2mS9fX1mjfeaP1GXh5/dTyeuGy19p1GxCAR4c6dO9Xz520+jYZjdu7cmT5y5IimsfFLtyGSPpEQobx8GT83559AxHet1guf0+uL14uiCAyDnKIoRAQqIqMajQZudtZ/uaHhmncRkbq6HPVGo+FmrbYwP5GIuU+d+tMZRAwTEb7wwgvctm033cMwLCDi0ZMnLSuqq2u+xvN5umg0PH769FtvI6KQUxYR1RMnzqzdsKG6UastWp5Op2J+/2QHIvYRESIiNTcTg4jqgw8+WLRnzyOf12oLa4gUKRQK9iHi2ZxxOjo6tDpd2b2CkAwj4pu9vY56ni+4cWJi+igAALhcU7+PRNLU1tb99eFhb58sE8ViKokikccz4+vu7r0ZAOD22+sLHA634nZPUXNzc/7g4MX3iIgEIYOPRhWy20fv7+623yvLRKlU5gkEkhSPZ/6LIpGqEtlsQ48QETs4ePHo3Fxmfm4uSYJA5PXOXrJaHY0AAK2trWVer5/c7qnprq6BWzwevxoOSxSNKpRIEA0Neaytra1luQDX2Tnwg0uXAnIymZGbTBL5/XEaGHA9mfUcFgDAZhu6bnx8ZlQQiOJxomSSKByWaGTEe/L06Q4jAMCJE39dFQ6nqbd3eOD8+d7dokgUDAp07FhrDQcAwDDIRKNxMBrLnuV5PtTbO7AzmRRiV121/Jsm08rviGLqze7u7prHHnssgoizkiSxxcXFBaoqt4yMeN9KpVJSYWFBgyimqrZsqX7RZnPcNjkZeDgQmJN4nmUR2QOiKL7FcdxwYaGWi0Yj7oaGzX/u73c+U1dX/S2nc6InEgn9IhqNTJaXL2+qqDAdNhrL3jxz5kz1/Py8KIqpKMuyalFR0TFFUV6yWruPFBbml5aWlj+5fv3qbamU8OjWrbX7z5+3bjKZVvxaURSfw2H/aTg8N1pWVrHaYCg9vGlT1X6rdbADEVstlq4Ver3+hNFoLBsfH/99MDhn5nlNUUlJ6cPr16/6iixLRwHgdr2ep7m5uTTHcXqdrviXgUDwj1NTl9tTqZg/63GTv4vHZXI6vY7m5ub8hWf5wgX3K5JENDTk3pt5vxh0Oi/N7969u2AhbnDQfYsgEPX1jexbfD+43dPewUG3deHYmTNtV09MzMou1+WRxfieHscDkkTU3+/cAwDgdE6E5ufTZLePPbcQ19bWeWMwKCgjI96+DM57nySR0t/vfGYhrqOjr1GSyOZwjB0AAOjvdz4hSUR2u+v5hbinn35aOzrquxAKCWSzDW89caLNMDrqE/3+GNlsF/YuxDIAmXuB51kIh8O/PHTokGixePJPnjyZR0RMKBT6TSKRJoZhv5j9hgAAdDodZiNOvsVCHCIVAwAUFBRcR0SMx5MZz0QhkvPztdtOnTpVkY12bHFx8Y0Gg5FNJBLBtjbb3vPne/e3t9sebW+3PirL8tpYLKVqNJqbAYAFAE0oFFQmJyefyhzvQZ6IWFGMjwSDQVWWZSMAgCCkBoPBGGMwGB9wOr2/s9td3+/tHbxVVZMujQYbNm/e8CQAgFarbQwGYzQ3F3qOiJjBwUHe4/Hk79u3T0gmY0f0+nzKy9Pc0tPzlxTLMnnh8Hzk7bePvUhErM1m0wAAZg0HjCgSpFKJCSLCQGC11N3dLSGiOjHhnZ6fj6AsK2U5I+es3tTUJAOA3NSEMhGoRACUuYFVrzcz3tSUSZQlSZJCoZDU1IQyIiocxy4XBJFKSgzbN2/e8kxt7TVPbty4+fCmTXWHq6pqHgVARpZlHQAgAGoURQkAOAKIqLS0tMiIqDidTlAUhQBAAQDm2murB6amLu8homRl5crvbthQ9WJ19cZTJlPVpMs12drZ2b86owJWBINB1eudmEVEdePGjZLX65WJCBmG8ckyIcNwZQAAHKcBWZanAUAEALWhoUECgEyIJQLieYT8/IJSAIDaWmDvuOMOlojYDRuqinW6ImIYNpp1uH814f3Ad0QQKi4uxPn54EuTk956n2/icx7P6Gc9Hs91Pt/FupmZiRui0fCPX3nllQIiVWEYhqusbOQAAA4ePLgUfzp69Gjh1q21z587d3r92NjI591u97e93vHDqVRqbNWqyjt0Ot3zAEBEENbr9bh27cpCImLPnj3LmkwmNuMYjJFlkVRViUejUWIYBiRJFhdXVFxWCQAg0On0P0TEtwAgnQM4HKP3lpRo0e+Hjix2yfyHZRkiIkDEpQxLAAAajeb9OVFMdoZCYdBqC+traq7uWwju6rLvqa5e/b1IJHn4vvvcZqfzZiRS84hIs6TFCDkAoC1bPvtEOCzeMzAwcFdDw6a/5eb379//xP337w2wLFcLAJBIxM+tW7fy2pkZ3S5EPJSFyQAAQ0OeXaKoYCIR77jhhhtYyhjnQzrl8jgmEAjLqqpuHR72HJufD/yC4/gEx7F3Go3l+3w+f1qSUn8AAAaA0ojELWYkSTLHMKgSqcKHFSNOVVV5dnZWyL6ziDja2zt8vK7uM3cODV08FY2GnmIYzZSq0naTacWvIpEkC6D2Hjy4mr/7bpAAUM4ey8W8ZYYhGQAAkbHqdHkPmUwrnn/vva4fCULKVVRUyOfl5d2zbNkyzfj4+EDGcNEXfL6ZBysqTD/r7x9KRiLR4xzHFBoMVz28Zs3qxosXxy8Iwvw5o9G4TFUVmUhVljScqqpFK1YYuJmZiQN6ffkTW7ZsuysWi4LBUAyBQEiJRiO76+pqPPX19QUAWIrIagRByO0CAQAkk4nxVEpiOI7f2t7ertu+HRJEmUQznU7PrlmzZnUodF0DEb139uxZbG4mJp2+8IDb7Stft27trYKw8lZBSEJJiR6CwZA6Nub83o4d17vM5pPlPM8XL+VpPM9jUVGRNhIJGwCAefzxn75+4MDPb6uurrmntLT8r+FwCHg+D8rKDDA5OX1JFBOPZJPg4a4u+66Kisrf1tbWHg4Go4c1Gg6KiwvA55sei8fD32xqapItFoumpMTARSIR45KGk2XpVVGEEVlmj1+65DvOMPiYLMtrotGId3Z27g833bTVTkSYuRjVg4qiQHl5eTqzy6hmF2Pv7LR/t7Jy+VeJ2OsR8Uy2vlRFMf6dQMD/E57XfPHZZ5/tfOihh9I7dgAgbvYDQOPIiPdulmW/kE6nSwQhOTY97X9tx47rB4gIX3755ZgkpX6GyInj4+PxhZsFAMl4PH4QEWMAAC0tx5SWlpZ77XbXG0VFBV9OJAQTz6eT0eh8e09P++u7du0KZksqBhFfe/fdszZF2XB3MinUsSybDgbx3DvvHH917969UQAAQRBCoVDocSKYXWrjPpauVNf9u/RRNe2/KvNjeH5sTQuQaRJ8ImFmsznXIcDm5uZcp4LN5lwfYJLpVHz4jssuhs0+iyLo+92PpQp3XNA1YZfC/SMfXFLm+52TRevI8WOv1EVZsK4r4nLzS8n+lD6l/xz9HUoIYo9JHffMAAAAAElFTkSuQmCC" alt=""></div>
  <div class="tc">
   <button id="stopbtn" onclick="stopGeneration()" title="Stop generation" style="display:none"><span class="i i16"><svg viewBox="0 0 24 24" fill="none"><rect x="8" y="8" width="8" height="8" rx="1.5" fill="#ff4466"/></svg></span></button><div class="ms"><button class="mb" id="mbtn" onclick="toggleMD()"><span id="mname">Loading...</span><span class="arr"><span class="i i12"><svg viewBox="0 0 24 24"><polyline points="6 9 12 15 18 9"/></svg></span></span></button><div class="md" id="mdd"></div></div>
   <div class="tkb" id="tkb" style="display:none"><span class="tkb-total" id="tkb-total"></span><span class="tkb-inc" id="tkb-inc"></span></div>
  </div>
  <div class="tr">
    <div class="search-container">
     <input type="text" class="search-input" id="chat-search" placeholder="Search a word" autocomplete="off">
     <div class="search-nav" id="search-nav" style="display:none">
      <span class="search-count" id="search-count"></span>
      <button class="search-nav-btn" id="search-prev" title="Previous">&#9664;</button>
      <button class="search-nav-btn" id="search-next" title="Next">&#9654;</button>
      <button class="search-close" id="search-close" title="Close">&#10005;</button>
     </div>
    </div>
   <button class="tb-solid" onclick="summarizeChat()" title="Summarize chat"><span class="i i16"><svg viewBox="0 0 24 24"><line x1="8" y1="6" x2="21" y2="6"/><line x1="8" y1="12" x2="21" y2="12"/><line x1="8" y1="18" x2="16" y2="18"/><line x1="3" y1="6" x2="3.01" y2="6"/><line x1="3" y1="12" x2="3.01" y2="12"/><line x1="3" y1="18" x2="3.01" y2="18"/></svg></span></button>
   <button class="tb" onclick="openSet()" title="Settings"><span class="i i16"><svg viewBox="0 0 24 24"><circle cx="12" cy="12" r="3"/><path d="M19.4 15a1.65 1.65 0 00.33 1.82l.06.06a2 2 0 010 2.83 2 2 0 01-2.83 0l-.06-.06a1.65 1.65 0 00-1.82-.33 1.65 1.65 0 00-1 1.51V21a2 2 0 01-4 0v-.09A1.65 1.65 0 009 19.4a1.65 1.65 0 00-1.82.33l-.06.06a2 2 0 01-2.83-2.83l.06-.06A1.65 1.65 0 004.68 15a1.65 1.65 0 00-1.51-1H3a2 2 0 010-4h.09A1.65 1.65 0 004.6 9a1.65 1.65 0 00-.33-1.82l-.06-.06a2 2 0 012.83-2.83l.06.06A1.65 1.65 0 009 4.68a1.65 1.65 0 001-1.51V3a2 2 0 014 0v.09a1.65 1.65 0 001 1.51 1.65 1.65 0 001.82-.33l.06-.06a2 2 0 012.83 2.83l-.06.06A1.65 1.65 0 0019.4 9a1.65 1.65 0 001.51 1H21a2 2 0 010 4h-.09a1.65 1.65 0 00-1.51 1z"/></svg></span></button>
   <div class="prof-btn-stack">
    <button class="tb" id="profile-tb-btn" onclick="openProfileModal()" title="Profiles"><span class="i i16"><svg viewBox="0 0 24 24"><path d="M17 21v-2a4 4 0 00-4-4H5a4 4 0 00-4 4v2"/><circle cx="9" cy="7" r="4"/><path d="M23 21v-2a4 4 0 00-3-3.87"/><path d="M16 3.13a4 4 0 010 7.75"/></svg></span></button>
    <button class="tb" id="prof-logout-btn" onclick="logoutProfile()" title="Log out"><span class="i i12"><svg viewBox="0 0 24 24"><path d="M9 21H5a2 2 0 01-2-2V5a2 2 0 012-2h4"/><polyline points="16 17 21 12 16 7"/><line x1="21" y1="12" x2="9" y2="12"/></svg></span></button>
   </div>
  </div>
 </div>
 <div id="msgs"><div class="es"><h2>What do you want to know?</h2><p>Start a conversation or pick a chat from the sidebar.</p></div></div>
 <div id="inarea">
  <div id="fba"><div class="fb-list" id="fb-list"></div></div>
  <div class="iw">
   <div class="ib">
    <div class="plus-wrap">
     <button id="plusbtn" onclick="togglePlusMenu(event)" title="More options">+</button>
     <div class="plus-menu" id="plusmenu">
      <div class="plus-opt" id="wsopt" onclick="toggleSearch();updatePlusState()" data-ws><span class="i i14"><svg viewBox="0 0 24 24"><circle cx="11" cy="11" r="8"/><line x1="21" y1="21" x2="16.65" y2="16.65"/></svg></span>Web Search</div>
      <div class="plus-opt" onclick="uploadFile();closePlusMenu()"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M21.44 11.05l-9.19 9.19a6 6 0 01-8.49-8.49l9.19-9.19a4 4 0 015.66 5.66l-9.2 9.19a2 2 0 01-2.83-2.83l8.49-8.48"/></svg></span>Attach Files</div>
      <div class="plus-opt" onclick="openTemplates();closePlusMenu()"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="16" y1="13" x2="8" y2="13"/><line x1="16" y1="17" x2="8" y2="17"/><polyline points="10 9 9 9 8 9"/></svg></span>Templates</div>
      <div class="plus-opt" onclick="openCompare();closePlusMenu()"><span class="i i14"><svg viewBox="0 0 24 24"><rect x="2" y="3" width="9" height="18" rx="2"/><rect x="13" y="3" width="9" height="18" rx="2"/></svg></span>Compare Models</div>
     </div>
    </div>
    <textarea id="uinp" placeholder="Ask anything… or paste a URL to fetch" rows="1"></textarea>
    <button id="vbtn" onclick="toggleVoice()" title="Voice"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M12 1a3 3 0 00-3 3v8a3 3 0 006 0V4a3 3 0 00-3-3z"/><path d="M19 10v2a7 7 0 01-14 0v-2"/><line x1="12" y1="19" x2="12" y2="23"/><line x1="8" y1="23" x2="16" y2="23"/></svg></span></button>
   </div>
   <button id="sbtn" onclick="sendMessage()"><img src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAABwAAAAcCAYAAAByDd+UAAABw0lEQVR4nL2WPU5bQRDHf7M2DZIbOAISgoY7pM4Vkip1rpAul+ACHIQDUCNRQRFFstyli/2nYMZer997u2shT2O/nY//fO3sGBWSZIDlR2amghf8TfDGyMYYmTEB1Ay5fNjUmPwgoKQUf2tAI/oz1900CWfe9oBYruff07YkpUIpubdH0aS+M1P23R3lkK5HurU7zwUi55KSmW0kXQG/gRsgsevIaKa8vgb8B86BH8CjpJmZrSPd237Iww6mpAtJf9VPf1y3rGkCmPvhnqce3TWwBF6BWRbVJXAGvACLLNIN8Az8MrNVZGk/07KDRqnVRdI3SStJP2uyA+cpv3M1wKj39yx9X/xsFr9T9iSlJrCC8ob56mcJwMzWtcs+Z79+PaChT4cNRaufhMysrX6fSYnjUnoURdOcLKXg3XVywM5BnU8mddxjC0DRl9Z/Lr8GFn7vWvpg52jLm5cN9UtJbz5plpLuJM1rWdrjl2/WhFKU4FbSg6QnSfdDr0Opd8DrHeKtVAZjBTPRsDi5nPHxJNnY/Aw5M1sPArpQ+8Y17ZTx0ZTVXfVgvzkCbFS/tgjHtWlJc76B9y3CI8BwuDSVRquOvQOC88GWdq5oUQAAAABJRU5ErkJggg==" alt=""></button>
  </div>
 </div>
</div>
<!-- Floating KB shortcut btn -->
<button id="kb-float-btn" onclick="$('kbm').classList.add('show')" title="Keyboard shortcuts">?</button>
<!-- Image Folder Hamburger Button -->
<button id="img-folder-btn" onclick="toggleImgFolderPanel()" title="Image folder viewer">
 <svg viewBox="0 0 24 24" width="18" height="18" fill="none" stroke="currentColor" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"><rect x="3" y="3" width="18" height="18" rx="2"/><line x1="3" y1="9" x2="21" y2="9"/><line x1="9" y1="21" x2="9" y2="9"/></svg>
</button>
<!-- Image Folder Panel -->
<div id="img-folder-panel">
 <div class="ifp-header">
  <svg viewBox="0 0 24 24" width="14" height="14" fill="none" stroke="var(--ac)" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round" style="flex-shrink:0"><rect x="3" y="3" width="18" height="18" rx="2"/><line x1="3" y1="9" x2="21" y2="9"/><line x1="9" y1="21" x2="9" y2="9"/></svg>
  <span class="ifp-title" id="ifp-title">Image Folder</span>
  <button class="ifp-close" onclick="toggleImgFolderPanel()" title="Close">✕</button>
 </div>
 <button class="ifp-load-btn" onclick="loadImageFolder()">
  <svg viewBox="0 0 24 24" width="13" height="13" fill="none" stroke="currentColor" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg>
  Upload Folder…
 </button>
 <div class="ifp-count" id="ifp-count"></div>
 <div class="ifp-list" id="ifp-list"></div>
</div>
<!-- Image Folder Hover Preview -->
<div class="ifp-preview" id="ifp-preview">
 <img id="ifp-preview-img" src="" alt="">
 <embed id="ifp-preview-pdf" src="" type="application/pdf" style="display:none;width:100%;height:460px;border-radius:6px">
 <div class="ifp-preview-name" id="ifp-preview-name"></div>
 <div class="ifp-preview-meta" id="ifp-preview-meta"></div>
</div>
<!-- Modals -->
<div class="mo" id="setm"><div class="ml"><h3><span class="i i16"><svg viewBox="0 0 24 24"><circle cx="12" cy="12" r="3"/><path d="M19.4 15a1.65 1.65 0 00.33 1.82l.06.06a2 2 0 010 2.83 2 2 0 01-2.83 0l-.06-.06a1.65 1.65 0 00-1.82-.33 1.65 1.65 0 00-1 1.51V21a2 2 0 01-4 0v-.09A1.65 1.65 0 009 19.4a1.65 1.65 0 00-1.82.33l-.06.06a2 2 0 01-2.83-2.83l.06-.06A1.65 1.65 0 004.68 15a1.65 1.65 0 00-1.51-1H3a2 2 0 010-4h.09A1.65 1.65 0 004.6 9a1.65 1.65 0 00-.33-1.82l-.06-.06a2 2 0 012.83-2.83l.06.06A1.65 1.65 0 009 4.68a1.65 1.65 0 001-1.51V3a2 2 0 014 0v.09a1.65 1.65 0 001 1.51 1.65 1.65 0 001.82-.33l.06-.06a2 2 0 012.83 2.83l-.06.06A1.65 1.65 0 0019.4 9a1.65 1.65 0 001.51 1H21a2 2 0 010 4h-.09a1.65 1.65 0 00-1.51 1z"/></svg></span> Settings</h3>
<label>API Key (Chutes.ai)</label><input type="password" id="sak" placeholder="API key">
<label>Brave Search API Key</label><input type="password" id="sbk" placeholder="Optional">
<label>Base URL</label><input type="text" id="sbu">
<label>Max Tokens per Request</label><div class="sr"><input type="range" id="smt" min="1000" max="30000" step="500"><span class="sv" id="smtv">10000</span></div>
<label>Max Messages in Memory</label><div class="sr"><input type="range" id="smm" min="10" max="100" step="5"><span class="sv" id="smmv">50</span></div>

<label>Theme</label><div class="thr"><div class="thc" data-theme="void" onclick="setTh('void')">Void</div><div class="thc" data-theme="nebula" onclick="setTh('nebula')">Nebula</div><div class="thc" data-theme="abyss" onclick="setTh('abyss')">Abyss</div><div class="thc" data-theme="matrix" onclick="setTh('matrix')">Matrix</div></div>

<div style="margin:10px 0 4px"><button class="bc" onclick="closeM('setm');openSP()" style="width:100%;display:flex;align-items:center;justify-content:center;gap:8px"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M12 20h9"/><path d="M16.5 3.5a2.121 2.121 0 013 3L7 19l-4 1 1-4L16.5 3.5z"/></svg></span> System Prompt</button></div>
<div style="margin:0 0 4px"><button class="bc" onclick="openThemeEditor()" style="width:100%;display:flex;align-items:center;justify-content:center;gap:8px"><span class="i i14"><svg viewBox="0 0 24 24"><circle cx="13.5" cy="6.5" r="1.5"/><circle cx="17.5" cy="10.5" r="1.5"/><circle cx="8.5" cy="7.5" r="1.5"/><circle cx="6.5" cy="12.5" r="1.5"/><path d="M12 2C6.5 2 2 6.5 2 12s4.5 10 10 10c.926 0 1.648-.746 1.648-1.688 0-.437-.18-.835-.437-1.125-.29-.289-.438-.652-.438-1.125a1.64 1.64 0 011.668-1.668h1.996c3.051 0 5.555-2.503 5.555-5.554C21.965 6.012 17.461 2 12 2z"/></svg></span> Open Theme Editor</button></div>
<div class="mbs"><button class="bc" onclick="closeM('setm');openExp()" style="margin-right:auto;display:flex;align-items:center;gap:6px"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M21 15v4a2 2 0 01-2 2H5a2 2 0 01-2-2v-4"/><polyline points="17 8 12 3 7 8"/><line x1="12" y1="3" x2="12" y2="15"/></svg></span> Export</button><button class="bc" onclick="closeM('setm')">Cancel</button><button class="bs" onclick="saveSet()">Save</button></div></div></div>
<div class="mo" id="spm"><div class="ml spm"><h3><span class="i i16"><svg viewBox="0 0 24 24"><path d="M12 20h9"/><path d="M16.5 3.5a2.121 2.121 0 013 3L7 19l-4 1 1-4L16.5 3.5z"/></svg></span> System Prompt</h3>
<p style="font-size:12px;color:var(--t3);margin-bottom:10px">Customize the AI behavior for this chat.</p>
<textarea id="spi" placeholder="You are a helpful AI assistant..."></textarea>
<div class="mbs"><button class="bc" onclick="closeM('spm')">Cancel</button><button class="bs" onclick="saveSP()">Save</button></div></div></div>
<div class="mo" id="expm"><div class="ml"><h3><span class="i i16"><svg viewBox="0 0 24 24"><path d="M21 15v4a2 2 0 01-2 2H5a2 2 0 01-2-2v-4"/><polyline points="17 8 12 3 7 8"/><line x1="12" y1="3" x2="12" y2="15"/></svg></span> Export Chat</h3>
<p style="font-size:12px;color:var(--t3);margin-bottom:16px">Choose a format to save this conversation.</p>
<div class="exp-grid">
 <button class="exp-btn" onclick="doExp('md')">
  <svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="16" y1="13" x2="8" y2="13"/></svg>
  <span>Markdown</span><small>.md file</small>
 </button>
 <button class="exp-btn" onclick="doExp('txt')">
  <svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="9" y1="13" x2="15" y2="13"/><line x1="9" y1="17" x2="13" y2="17"/></svg>
  <span>Plain Text</span><small>.txt file</small>
 </button>
 <button class="exp-btn" onclick="doExp('html')">
  <svg viewBox="0 0 24 24"><polyline points="16 18 22 12 16 6"/><polyline points="8 6 2 12 8 18"/></svg>
  <span>HTML</span><small>Styled page</small>
 </button>
</div>
<div class="mbs" style="margin-top:16px"><button class="bc" onclick="closeM('expm')">Cancel</button></div></div></div>
<div id="ctx"></div>
<div id="cdlg"><div class="cdlg-box"><div class="cdlg-msg" id="cdlg-msg"></div><div class="cdlg-btns" id="cdlg-btns"></div></div></div>
<div id="rnml" class="mo"><div class="ml"><h3><span class="i i16"><svg viewBox="0 0 24 24"><path d="M12 20h9"/><path d="M16.5 3.5a2.121 2.121 0 013 3L7 19l-4 1 1-4L16.5 3.5z"/></svg></span> Rename Chat</h3>
<p id="rnml-prompt" style="font-size:12px;color:var(--t3);margin-bottom:14px">Enter a new title for this chat.</p>
<input type="text" id="rnml-input" placeholder="Chat title" style="width:100%;background:var(--bg2);border:1px solid rgba(255,255,255,0.06);color:var(--t1);padding:10px 14px;border-radius:10px;font-family:inherit;font-size:14px;outline:none">
<div class="mbs"><button class="bc" onclick="closeM('rnml')">Cancel</button><button class="bs" onclick="saveRename()">Rename</button></div></div></div>

<div class="mo" id="tplm"><div class="ml">
<h3><span class="i i16"><svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="16" y1="13" x2="8" y2="13"/><line x1="16" y1="17" x2="8" y2="17"/><polyline points="10 9 9 9 8 9"/></svg></span> Prompt Templates</h3>
<p style="font-size:13px;color:var(--t3);margin-bottom:14px">Click a template to paste it into the input. Save your own below.</p>
<div class="tpl-search-wrap">
 <input type="text" id="tpl-search" placeholder="Search templates…" oninput="filterTemplates(this.value)">
 <button class="tpl-search-clear" onclick="clearTplSearch()" title="Clear">×</button>
</div>
<div class="tpl-list" id="tpl-list"></div>
<div class="tpl-add">
 <input type="text" id="tpl-name-inp" placeholder="Template name (e.g. Code Review)">
 <textarea id="tpl-text-inp" placeholder="Prompt text…"></textarea>
 <div class="mbs"><button class="bc" onclick="closeM('tplm')">Close</button><button class="bs" onclick="saveNewTemplate()">Save Template</button></div>
</div></div></div>

<div class="mo" id="foldrm"><div class="ml" style="min-width:340px;max-width:420px">
<h3><span class="i i16"><svg viewBox="0 0 24 24"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg></span> Move to Folder</h3>
<p style="font-size:12px;color:var(--t3);margin-bottom:14px">Choose a folder for this chat.</p>
<div class="foldr-list" id="foldr-list"></div>
<div class="foldr-add-row">
 <div style="display:flex;gap:6px;align-items:center"><input type="color" id="foldr-color-inp" class="fd-color-btn" value="#6366f1" title="Pick folder color"><input type="text" id="foldr-new-inp" placeholder="New folder name…" onkeydown="if(event.key==='Enter')addFolderFromModal()" style="flex:1"></div>
 <button onclick="addFolderFromModal()">+ Add</button>
</div>
<div class="mbs" style="margin-top:12px"><button class="bc" onclick="closeM('foldrm')">Cancel</button></div>
</div></div>

<div class="mo" id="summm"><div class="ml">
<h3><span class="i i16"><svg viewBox="0 0 24 24"><line x1="8" y1="6" x2="21" y2="6"/><line x1="8" y1="12" x2="21" y2="12"/><line x1="8" y1="18" x2="16" y2="18"/><line x1="3" y1="6" x2="3.01" y2="6"/><line x1="3" y1="12" x2="3.01" y2="12"/><line x1="3" y1="18" x2="3.01" y2="18"/></svg></span> Chat Summary</h3>
<div id="summ-content"><div class="summ-loading">⏳ Generating summary…</div></div>
<div class="mbs" style="margin-top:14px"><button class="bc" onclick="closeM('summm')">Close</button><button class="bs" id="summ-copy-btn" onclick="copySummary()" style="display:none">Copy</button></div>
</div></div>

<!-- Keyboard Shortcut Modal -->
<div class="mo" id="profm"><div class="ml" style="min-width:760px;max-width:880px">
<h3><span class="i i16"><svg viewBox="0 0 24 24"><path d="M17 21v-2a4 4 0 00-4-4H5a4 4 0 00-4 4v2"/><circle cx="9" cy="7" r="4"/><path d="M23 21v-2a4 4 0 00-3-3.87"/><path d="M16 3.13a4 4 0 010 7.75"/></svg></span> Profiles</h3>
<p style="font-size:12px;color:var(--t3);margin:6px 0 14px">Each profile saves its own chats, folders, theme and settings. Auto-saves when the app closes.</p>
<div id="prof-list"></div>
<div class="prof-sep"></div>
<div class="prof-create">
 <p style="font-size:13px;font-weight:600;color:var(--t2);margin-bottom:2px">Create new profile</p>
 <input type="text" id="prof-new-name" placeholder="Profile name (e.g. Your name)" maxlength="32" autocomplete="off">
 <input type="password" id="prof-new-pw" placeholder="Password" autocomplete="new-password">
 <input type="password" id="prof-new-pw2" placeholder="Confirm password" autocomplete="new-password">
 <div style="display:flex;gap:8px;align-items:center">
  <button class="prof-create-btn" onclick="createProfile()"><svg viewBox="0 0 24 24" width="15" height="15" fill="none" stroke="currentColor" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"><path d="M16 21v-2a4 4 0 00-4-4H6a4 4 0 00-4 4v2"/><circle cx="9" cy="7" r="4"/><line x1="19" y1="8" x2="19" y2="14"/><line x1="22" y1="11" x2="16" y2="11"/></svg>Create &amp; Switch</button>
  <button class="bc" onclick="closeM('profm')" style="padding:10px 22px;border-radius:10px;cursor:pointer;font-family:inherit;font-size:14px;font-weight:500;transition:all .15s">Close</button>
 </div>
</div>
</div></div>

<!-- Delete Profile Modal -->
<div class="mo" id="del-prof-m"><div class="ml" style="min-width:360px;max-width:430px">
<h3><span class="i i16"><svg viewBox="0 0 24 24"><polyline points="3 6 5 6 21 6"/><path d="M19 6v14a2 2 0 01-2 2H7a2 2 0 01-2-2V6m3 0V4a2 2 0 012-2h4a2 2 0 012 2v2"/></svg></span> Delete Profile</h3>
<p id="del-prof-msg" style="font-size:13px;color:var(--t2);margin:10px 0 16px;line-height:1.5"></p>
<input type="password" id="del-prof-pw" placeholder="Profile password" autocomplete="off"
 style="width:100%;background:var(--bg2);border:1px solid rgba(255,255,255,.08);color:var(--t1);padding:10px 14px;border-radius:9px;font-family:inherit;font-size:14px;outline:none;box-sizing:border-box"
 onkeydown="if(event.key==='Enter')confirmDeleteProfile();if(event.key==='Escape')closeM('del-prof-m')">
<div class="mbs" style="margin-top:16px">
 <button class="bc" onclick="closeM('del-prof-m')">Cancel</button>
 <button onclick="confirmDeleteProfile()" style="padding:10px 22px;border-radius:10px;cursor:pointer;font-family:inherit;font-size:14px;font-weight:600;background:#ef4444;border:none;color:#fff;transition:all .15s" onmouseover="this.style.background='#dc2626'" onmouseout="this.style.background='#ef4444'">Delete</button>
</div>
</div></div>

<div class="mo" id="kbm"><div class="ml" style="min-width:560px;max-width:640px"><h3><span class="i i16"><svg viewBox="0 0 24 24"><rect x="2" y="7" width="20" height="14" rx="2"/><path d="M16 7V5a2 2 0 00-2-2h-4a2 2 0 00-2 2v2"/><line x1="12" y1="12" x2="12" y2="16"/><line x1="10" y1="14" x2="14" y2="14"/></svg></span> Keyboard Shortcuts</h3>
<div class="kb-grid" style="margin-top:16px">
 <div class="kb-row"><span class="kb-desc">New chat</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">N</span></span></div>
 <div class="kb-row"><span class="kb-desc">Send message</span><span class="kb-keys"><span class="kb-key">Enter</span></span></div>
 <div class="kb-row"><span class="kb-desc">New line in input</span><span class="kb-keys"><span class="kb-key">Shift</span><span class="kb-key">Enter</span></span></div>
 <div class="kb-row"><span class="kb-desc">Open Settings</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">/</span></span></div>
 <div class="kb-row"><span class="kb-desc">Search chats</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">F</span></span></div>
 <div class="kb-row"><span class="kb-desc">This help modal</span><span class="kb-keys"><span class="kb-key">?</span></span></div>
 <div class="kb-row"><span class="kb-desc">Close any modal</span><span class="kb-keys"><span class="kb-key">Esc</span></span></div>
 <div class="kb-row"><span class="kb-desc">Toggle sidebar</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">B</span></span></div>
 <div class="kb-row"><span class="kb-desc">Export chat</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">E</span></span></div>
 <div class="kb-row"><span class="kb-desc">Web Search toggle</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">W</span></span></div>
 <div class="kb-row"><span class="kb-desc">Attach file(s)</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">U</span></span></div>
 <div class="kb-row"><span class="kb-desc">Compare models</span><span class="kb-keys"><span class="kb-key">Ctrl</span><span class="kb-key">M</span></span></div>
</div>
<div class="mbs" style="margin-top:18px"><button class="bs" onclick="closeM('kbm')">Got it</button></div>
</div></div>

<!-- Theme Editor Modal -->
<div class="mo" id="them"><div class="ml" style="min-width:420px;max-width:500px">
<h3><span class="i i16"><svg viewBox="0 0 24 24"><circle cx="13.5" cy="6.5" r="1.5"/><circle cx="17.5" cy="10.5" r="1.5"/><circle cx="8.5" cy="7.5" r="1.5"/><circle cx="6.5" cy="12.5" r="1.5"/><path d="M12 2C6.5 2 2 6.5 2 12s4.5 10 10 10c.926 0 1.648-.746 1.648-1.688 0-.437-.18-.835-.437-1.125-.29-.289-.438-.652-.438-1.125a1.64 1.64 0 011.668-1.668h1.996c3.051 0 5.555-2.503 5.555-5.554C21.965 6.012 17.461 2 12 2z"/></svg></span> Theme Editor</h3>
<p style="font-size:12px;color:var(--t3);margin-bottom:16px">Customize accent color, font size, and assistant text width.</p>
<div class="te-row">
 <span class="te-label">Accent Color</span>
 <input type="color" id="te-accent" class="te-swatch" value="#e88a2a" oninput="applyTeAccent(this.value)">
 <span class="te-val" id="te-accent-val">#e88a2a</span>
</div>
<div class="te-row">
 <span class="te-label">Font Size</span>
 <input type="range" class="te-slider" id="te-font" min="12" max="20" step="1" value="17" oninput="applyTeFont(this.value)">
 <span class="te-val" id="te-font-val">17px</span>
</div>
<div class="te-row">
 <span class="te-label">Text Width</span>
 <input type="range" class="te-slider" id="te-radius" min="40" max="90" step="5" value="80" oninput="applyTeRadius(this.value)">
 <span class="te-val" id="te-radius-val">80%</span>
</div>
<div class="te-preview" id="te-preview">
 <div style="display:flex;justify-content:flex-end;margin-bottom:10px"><div style="padding:10px 16px;background:var(--ubg);border-radius:20px 20px 5px 20px;font-size:var(--chat-font-size,16px);color:var(--t1);width:fit-content;max-width:72%">This is your prompt bubble.</div></div>
 <div style="padding:4px 0;font-size:var(--chat-font-size,16px);color:#cccce0;line-height:1.7">And this is an AI response — no background, full width text.</div>
</div>
<div class="mbs" style="margin-top:16px"><button class="bc" onclick="resetThemeEditor()">Reset</button><button class="bc" onclick="closeM('them')">Close</button><button class="bs" onclick="saveThemeEditor()">Apply</button></div>
</div></div>

<!-- Compare overlay -->
<div id="compare-overlay">
 <div class="cmp-header">
  <h3>Model Comparison</h3>
  <div class="cmp-selectors">
   <select class="cmp-sel" id="cmp-m1"></select>
   <span style="color:var(--t3);font-size:13px">vs</span>
   <select class="cmp-sel" id="cmp-m2"></select>
  </div>
  <button class="cmp-close" onclick="closeCompare()" title="Close">✕</button>
 </div>
 <div class="cmp-body">
  <div class="cmp-pane" id="cmp-pane1">
   <div class="cmp-pane-head"><span>Model A</span><span class="cmp-pane-model" id="cmp-label1">—</span></div>
   <div class="cmp-pane-body" id="cmp-body1"><div style="color:var(--t3);text-align:center;padding:40px 20px;font-size:13px">Send a message to start comparing</div></div>
  </div>
  <div class="cmp-pane" id="cmp-pane2">
   <div class="cmp-pane-head"><span>Model B</span><span class="cmp-pane-model" id="cmp-label2">—</span></div>
   <div class="cmp-pane-body" id="cmp-body2"><div style="color:var(--t3);text-align:center;padding:40px 20px;font-size:13px">Send a message to start comparing</div></div>
  </div>
 </div>
 <div class="cmp-input-row">
  <div class="cmp-input-wrap" style="flex-direction:column;align-items:stretch;gap:6px">
   <div id="cmp-fba"><div class="fb-list" id="cmp-fb-list"></div></div>
   <div style="display:flex;align-items:center;gap:8px">
   <div class="cmp-ib">
    <div class="plus-wrap">
     <button id="cmp-plusbtn" onclick="toggleCmpPlusMenu(event)" title="More options">+</button>
     <div class="plus-menu" id="cmp-plusmenu">
      <div class="plus-opt" id="cmp-wsopt" onclick="toggleSearch();updatePlusState();closeCmpPlusMenu()" data-ws><span class="i i14"><svg viewBox="0 0 24 24"><circle cx="11" cy="11" r="8"/><line x1="21" y1="21" x2="16.65" y2="16.65"/></svg></span>Web Search</div>
      <div class="plus-opt" onclick="uploadFileForCompare();closeCmpPlusMenu()"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M21.44 11.05l-9.19 9.19a6 6 0 01-8.49-8.49l9.19-9.19a4 4 0 015.66 5.66l-9.2 9.19a2 2 0 01-2.83-2.83l8.49-8.48"/></svg></span>Attach Files</div>
      <div class="plus-opt" onclick="openTemplatesForCompare();closeCmpPlusMenu()"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="16" y1="13" x2="8" y2="13"/><line x1="16" y1="17" x2="8" y2="17"/><polyline points="10 9 9 9 8 9"/></svg></span>Templates</div>
     </div>
    </div>
    <textarea id="cmp-inp" placeholder="Enter prompt to compare both models…" rows="1"
     oninput="autoR(this)"
     onkeydown="if(event.key==='Enter'&&!event.shiftKey){event.preventDefault();runCompare()}"></textarea>
    <button id="cmp-vbtn" onclick="toggleVoice()" title="Voice"><span class="i i14"><svg viewBox="0 0 24 24"><path d="M12 1a3 3 0 00-3 3v8a3 3 0 006 0V4a3 3 0 00-3-3z"/><path d="M19 10v2a7 7 0 01-14 0v-2"/><line x1="12" y1="19" x2="12" y2="23"/><line x1="8" y1="23" x2="16" y2="23"/></svg></span></button>
   </div>
   <button class="cmp-run-btn" onclick="runCompare()" title="Compare">▶</button>
   </div>
  </div>
 </div>
</div>
<script>
const I={
 copy:'<span class="i i14"><svg viewBox="0 0 24 24"><rect x="9" y="9" width="13" height="13" rx="2"/><path d="M5 15H4a2 2 0 01-2-2V4a2 2 0 012-2h9a2 2 0 012 2v1"/></svg></span>',
 spk:'<span class="i i14"><svg viewBox="0 0 24 24"><polygon points="11 5 6 9 2 9 2 15 6 15 11 19 11 5"/><path d="M15.54 8.46a5 5 0 010 7.07"/></svg></span>',
 edit:'<span class="i i14"><svg viewBox="0 0 24 24"><path d="M12 20h9"/><path d="M16.5 3.5a2.121 2.121 0 013 3L7 19l-4 1 1-4L16.5 3.5z"/></svg></span>',
 branch:'<span class="i i14"><svg viewBox="0 0 24 24"><line x1="6" y1="3" x2="6" y2="15"/><circle cx="18" cy="6" r="3"/><circle cx="6" cy="18" r="3"/><path d="M18 9a9 9 0 01-9 9"/></svg></span>',
 regen:'<span class="i i14"><svg viewBox="0 0 24 24"><polyline points="23 4 23 10 17 10"/><path d="M20.49 15a9 9 0 11-2.12-9.36L23 10"/></svg></span>',
 pin:'<span class="i i12"><svg viewBox="0 0 24 24"><path d="M12 2l2.09 6.26L21 9.27l-5 4.87L17.18 21 12 17.27 6.82 21 8 14.14 3 9.27l6.91-1.01L12 2z" fill="currentColor" stroke="none"/></svg></span>',
 pinO:'<span class="i i12"><svg viewBox="0 0 24 24"><path d="M12 2l2.09 6.26L21 9.27l-5 4.87L17.18 21 12 17.27 6.82 21 8 14.14 3 9.27l6.91-1.01L12 2z"/></svg></span>',
 del:'<span class="i i12"><svg viewBox="0 0 24 24"><polyline points="3 6 5 6 21 6"/><path d="M19 6v14a2 2 0 01-2 2H7a2 2 0 01-2-2V6m3 0V4a2 2 0 012-2h4a2 2 0 012 2v2"/></svg></span>',
 file:'<span class="i i14"><svg viewBox="0 0 24 24"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/></svg></span>',
 img:'<span class="i i14"><svg viewBox="0 0 24 24"><rect x="3" y="3" width="18" height="18" rx="2"/><circle cx="8.5" cy="8.5" r="1.5"/><polyline points="21 15 16 10 5 21"/></svg></span>',
 chev:'<span class="i i12"><svg viewBox="0 0 24 24"><polyline points="9 18 15 12 9 6"/></svg></span>',
 menu:'<span class="i i12"><svg viewBox="0 0 24 24"><line x1="3" y1="6" x2="21" y2="6"/><line x1="3" y1="12" x2="21" y2="12"/><line x1="3" y1="18" x2="21" y2="18"/></svg></span>',
};
marked.setOptions({highlight:(c,l)=>{if(l&&hljs.getLanguage(l))return hljs.highlight(c,{language:l}).value;return hljs.highlightAuto(c).value},breaks:true,gfm:true});
// Custom themed dialogs
function showDialog(msg,onOk,showCancel){
 $('cdlg-msg').textContent=msg;
 const btns=$('cdlg-btns');btns.innerHTML='';
 if(showCancel){const cb=document.createElement('button');cb.className='cdlg-cancel';cb.textContent='Cancel';cb.onclick=()=>{$('cdlg').classList.remove('show')};btns.appendChild(cb)}
 const ob=document.createElement('button');ob.className='cdlg-ok';ob.textContent='OK';ob.onclick=()=>{$('cdlg').classList.remove('show');if(onOk)onOk()};btns.appendChild(ob);
 $('cdlg').classList.add('show');
}
function showAlert(msg){showDialog(msg,null,false)}
function showConfirm(msg,onOk){showDialog(msg,onOk,true)}
let cid=null,pf=[],stm=false,rec=null,isRec=false,webS=false,models=[],curM='',selTh='void',lastU=null,maxTk=10000,userScrolled=false,activeStreams=new Set(),streamBuffer={},sessionTokens=0,lastTotalTokens=0,titleOverrides={},chatTokens={};
let chatDrafts={},activeFolderTab='All',folderColors={},folderList=[],_ctxSel='';
async function init(){
 const c=await pywebview.api.get_config();window._cfg=c;curM=c.current_model||'';models=c.models||[];maxTk=c.max_tokens||10000;selTh=c.theme||'void';applyTh(selTh);
 $('mname').textContent=short(curM);$('sak').value=c.api_key||'';$('sbk').value=c.brave_api_key||'';$('sbu').value=c.base_url||'';
 $('smt').value=c.max_tokens||10000;$('smtv').textContent=c.max_tokens||10000;$('smm').value=c.max_memory||50;$('smmv').textContent=c.max_memory||50;
 folderColors=c.folder_colors||{};folderList=c.folders||['General'];
 upTh();restoreThemeOverrides();renderFolderTabs();refreshCL();setupRangeListeners();
 setupRenameModal();
 // Sync topbar logout button visibility
 const logoutBtn=$('prof-logout-btn');
 if(logoutBtn){logoutBtn.classList.toggle('visible',!!c.current_profile_id);}
 const msgsEl=$('msgs');if(msgsEl){msgsEl.addEventListener('scroll',()=>{const nearBottom=isNearBottom();if(!nearBottom)userScrolled=true;else if(!stm)userScrolled=false})}
 const fbListEl=$('fb-list');if(fbListEl){fbListEl.addEventListener('click',function(e){const btn=e.target.closest('.pf-rm-btn');if(btn){e.stopPropagation();const idx=parseInt(btn.getAttribute('data-pf-idx'));if(!isNaN(idx))removePF(idx);}});}
}
function setupRenameModal(){
 const inp=$('rnml-input');
 if(inp){
  inp.addEventListener('keydown',e=>{
   if(e.key==='Enter'){e.preventDefault();saveRename()}
   else if(e.key==='Escape'){e.preventDefault();closeM('rnml');renameId=null}
  });
 }
}
function $(id){return document.getElementById(id)}
function short(m){return m?m.split('/').pop():'Select Model'}
function esc(s){const d=document.createElement('div');d.textContent=s||'';return d.innerHTML}
function rmd(t){if(typeof t!=='string')t=String(t);try{return marked.parse(t)}catch(e){return '<p>'+esc(t)+'</p>'}}
function applyTh(t){document.documentElement.setAttribute('data-theme',t);selTh=t}
function setTh(t){applyTh(t);upTh()}
function upTh(){document.querySelectorAll('.thc').forEach(c=>c.classList.toggle('active',c.getAttribute('data-theme')===selTh))}
function upTkB(u){const el=$('tkb'),totalEl=$('tkb-total'),incEl=$('tkb-inc');if(!u||!u.total_tokens){el.style.display='none';return}el.style.display='flex';const latest=u.completion_tokens||0;if(cid){chatTokens[cid]=(chatTokens[cid]||0)+latest}const sessionTotal=cid?(chatTokens[cid]||0):0;totalEl.textContent=sessionTotal.toLocaleString();if(latest>0){incEl.textContent=`+${latest.toLocaleString()}`;incEl.style.display='inline'}else{incEl.style.display='none'}}
function toggleSidebar(){$('sidebar').classList.toggle('collapsed');const fp=$('footer-powered'),fi=$('footer-icon');if($('sidebar').classList.contains('collapsed')){if(fp){fp.style.opacity='0';fp.style.transform='translateY(10px)';fp.style.filter='blur(4px)'}if(fi){fi.style.opacity='0';fi.style.transform='scale(0.8) rotate(-10deg)'}}else{if(fp){fp.style.opacity='';fp.style.transform='';fp.style.filter=''}if(fi){fi.style.opacity='';fi.style.transform=''}}}
async function refreshCL(){
 const chats=await pywebview.api.list_chats();const el=$('chatlist');
 const filtered=activeFolderTab==='All'?chats:chats.filter(c=>c.folder===activeFolderTab);
 if(!filtered.length){el.innerHTML=`<div style="text-align:center;padding:20px;color:var(--t3);font-size:12px">${activeFolderTab==='All'?'No chats yet':'No chats in this folder'}</div>`;return}
 el.innerHTML=filtered.map(c=>{
  const title=titleOverrides[c.id]||c.title;
  const fc=folderColors[c.folder]||'#555566';
  const fdot=c.folder&&c.folder!=='General'?`<span class="ci-fdot" style="background:${fc}"></span>`:'';
  const draftDot=chatDrafts[c.id]?`<span class="ci-draft" title="Draft saved"></span>`:'';
  return`<div class="ci ${c.id===cid?'active':''}" onclick="loadChat('${c.id}')" data-tip="${esc(title)}" data-id="${c.id}" draggable="true" ondragstart="onChatDragStart(event,'${c.id}')" ondragend="onChatDragEnd(event)">${c.pinned?`<button class="cip" onclick="event.stopPropagation();togglePin('${c.id}')" title="Unpin">`+I.pin+'</button>':''}${c.branched?'<span class="ci-branch">'+I.branch+'</span>':''}${fdot}<span class="cit">${esc(title)}</span>${draftDot}<span class="cia"><div class="chat-menu"><button onclick="event.stopPropagation();toggleChatMenu('${c.id}')" title="Actions">${I.menu}</button><div class="chat-menu-dropdown" id="cm-${c.id}"><div class="chat-menu-item" onclick="event.stopPropagation();togglePin('${c.id}');hideChatMenu('${c.id}')">${c.pinned?I.pin:I.pinO}<span>${c.pinned?'Unpin':'Pin'}</span></div><div class="chat-menu-item" onclick="event.stopPropagation();renameChat('${c.id}');hideChatMenu('${c.id}')">${I.edit}<span>Rename</span></div><div class="chat-menu-item" onclick="event.stopPropagation();moveChatToFolder('${c.id}');hideChatMenu('${c.id}')"><span class="i i12"><svg viewBox="0 0 24 24"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg></span><span>Move to folder</span></div><div class="chat-menu-item delete" onclick="event.stopPropagation();deleteChat('${c.id}');hideChatMenu('${c.id}')">${I.del}<span>Delete</span></div></div></div></span></div>`
 }).join('')
}
async function createNewChat(){const c=await pywebview.api.create_chat();cid=c.id;renderMsgs(c.messages);refreshCL();$('tkb').style.display='none';sessionTokens=0;$('tkb-total').textContent='';$('tkb-inc').textContent=''}
async function loadChat(id){
 if(!id)return;
 // Save draft for current chat before switching
 const inp=$('uinp');
 if(cid&&inp&&inp.value.trim())chatDrafts[cid]=inp.value;
 else if(cid&&inp&&!inp.value.trim()&&chatDrafts[cid])delete chatDrafts[cid];
 const c=await pywebview.api.get_chat(id);if(!c)return;cid=c.id;renderMsgs(c.messages);refreshCL();
 const prevTokens=chatTokens[cid]||0;if(prevTokens>0){$('tkb').style.display='flex';$('tkb-total').textContent=prevTokens.toLocaleString();$('tkb-inc').textContent=''}else{$('tkb').style.display='none';$('tkb-total').textContent='';$('tkb-inc').textContent=''}
 if(inp){
  const draft=chatDrafts[cid]||'';
  inp.value=draft;autoR(inp);
  if(draft){inp.focus();inp.setSelectionRange(draft.length,draft.length)}
 }
 userScrolled=false;
 if(activeStreams.has(cid)){stm=true;$('sbtn').disabled=true;if(!streamBuffer[cid])streamBuffer[cid]={};if(streamBuffer[cid].content||streamBuffer[cid].reasoning){restoreStreamingUI()}}
 else{stm=false;$('sbtn').disabled=false;const smsg=$('smsg');if(smsg)smsg.remove()}
}
async function deleteChat(id){await pywebview.api.delete_chat(id);if(cid===id){cid=null;showE()}refreshCL()}
async function deleteAllChats(){showConfirm('Delete ALL chats?',async()=>{await pywebview.api.delete_all_chats();cid=null;showE();refreshCL()})}
async function togglePin(id){await pywebview.api.toggle_pin(id);refreshCL()}
function toggleChatMenu(id){const menu=$(`cm-${id}`);if(!menu)return;const isOpen=menu.classList.contains('show');document.querySelectorAll('.chat-menu-dropdown').forEach(m=>m.classList.remove('show'));document.querySelectorAll('.ci.menu-open').forEach(c=>c.classList.remove('menu-open'));if(!isOpen){const ci=document.querySelector(`.ci[data-id="${id}"]`);if(ci)ci.classList.add('menu-open');const btn=document.querySelector(`.ci[data-id="${id}"] .chat-menu button`);if(btn){const rect=btn.getBoundingClientRect();menu.style.visibility='hidden';menu.style.display='block';const mw=menu.offsetWidth||130;menu.style.display='';menu.style.visibility='';menu.classList.add('show');menu.style.left=(rect.left-mw-4)+'px';menu.style.top=Math.min(rect.top,window.innerHeight-menu.offsetHeight-8)+'px'}}}
function hideChatMenu(id){const menu=$(`cm-${id}`);if(menu)menu.classList.remove('show');const ci=document.querySelector(`.ci[data-id="${id}"]`);if(ci)ci.classList.remove('menu-open')}
document.addEventListener('click',function(e){if(!e.target.closest('.chat-menu')){document.querySelectorAll('.chat-menu-dropdown').forEach(m=>m.classList.remove('show'));document.querySelectorAll('.ci.menu-open').forEach(c=>c.classList.remove('menu-open'))}});
let renameId=null;async function renameChat(id){const c=await pywebview.api.get_chat(id);if(!c)return;renameId=id;$("rnml-input").value=c.title;$("rnml").classList.add("show");$("rnml-input").focus();$("rnml-input").select()}async function saveRename(){const nt=$("rnml-input").value.trim();if(!renameId||!nt)return;const c=await pywebview.api.get_chat(renameId);if(!c||nt===c.title){closeM("rnml");renameId=null;return}await pywebview.api.rename_chat(renameId,nt);refreshCL();if(cid===renameId){const el=document.querySelector(`.ci[data-id="${renameId}"] .cit`);if(el)el.textContent=nt}closeM("rnml");renameId=null}
function showE(){$('msgs').innerHTML='<div class="es"><h2>What do you want to know?</h2><p>Test my knowledge or pick a chat.</p></div>';$('tkb').style.display='none';sessionTokens=0;$('tkb-total').textContent='';$('tkb-inc').textContent=''}
function renderMsgs(msgs){
 const a=$('msgs');const f=msgs.filter(m=>m.role!=='system');
 if(!f.length){a.innerHTML='<div class="es"><h2>What do you want to know?</h2><p>Test my knowledge or pick a chat.</p></div>';return}
 a.innerHTML=f.map((m,i)=>{const ri=i+1;const iu=m.role==='user';const il=i===f.length-1;
  let cs=typeof m.content==='string'?m.content:(Array.isArray(m.content)?m.content.filter(c=>c.type==='text').map(c=>c.text).join('\n'):String(m.content));
  let wsBlock='';
  if(iu){const wsM=cs.match(/^\[Web Search Results\]\n([\s\S]*?)\n\n\[User Question\]\n/);if(wsM){wsBlock=`<div class="ws-toggle" onclick="this.nextElementSibling.classList.toggle('show')"><span class="i i12"><svg viewBox="0 0 24 24"><line x1="3" y1="6" x2="21" y2="6"/><line x1="3" y1="12" x2="21" y2="12"/><line x1="3" y1="18" x2="21" y2="18"/></svg></span> Web results</div><div class="ws-content">${esc(wsM[1])}</div>`;cs=cs.replace(wsM[0],'')}}
  const usageBadge=m.usage?`<div class="mm"><span>Tokens: ${m.usage.total_tokens||'—'} (in: ${m.usage.prompt_tokens||0}, out: ${m.usage.completion_tokens||0})</span></div>`:'';
  if(iu){
   return `<div class="user-wrap"><div class="msg user" data-idx="${ri}"><div class="mh"><span class="mr">You</span></div>${wsBlock}<div class="mc" id="mc-${ri}">${rmd(cs)}</div>${usageBadge}</div>
  <div class="mas">
   <button class="ma2" onclick="copyM(${ri})" title="Copy">${I.copy}</button>
   <button class="ma2" onclick="editM(${ri})" title="Edit">${I.edit}</button>
  </div></div>`;
  }
  return `<div class="msg ${m.role}" data-idx="${ri}"><div class="mh"><span class="mr">${short(curM)}</span></div>${wsBlock}${m.reasoning?`<div class="rb"><div class="rt" onclick="toggleRB(this)"><span class="rt-icon"><svg viewBox="0 0 24 24" width="14" height="14" fill="none" stroke="var(--ac)" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M12 2a7 7 0 0 1 5 11.9V17a2 2 0 0 1-2 2h-6a2 2 0 0 1-2-2v-3.1A7 7 0 0 1 12 2z"/><line x1="9" y1="21" x2="15" y2="21"/></svg></span><span class="rblk-label-static">${m.reasoning_duration?'Thought for '+m.reasoning_duration+'s':'Thought'}</span><span class="ra">${I.chev}</span></div><div class="rc">${esc(m.reasoning)}</div></div>`:''}
<div class="mc" id="mc-${ri}">${rmd(cs)}</div>${usageBadge}
  <div class="mas">
   <button class="ma2" onclick="copyM(${ri})" title="Copy">${I.copy}</button>
   <button class="ma2" onclick="branchC(${ri})" title="Branch">${I.branch}</button>
   ${il?`<button class="ma2" onclick="regen()" title="Regenerate">${I.regen}</button>`:''}
   <div class="tts-wrap"><button class="ma2 tts-btn" id="ttsb-${ri}" onclick="speakMsg(${ri},this)" title="Speak">${I.spk}</button><button class="tts-chevron" onclick="openVoicePicker(event,${ri})" title="Voice">&#9660;</button><div class="voice-picker" id="vp-${ri}"></div></div>
  </div></div>`;
 }).join('');
 wrapTables(a);
 a.scrollTop=a.scrollHeight}
function stPH(){return `<div class="msg assistant" id="smsg"><div class="mh"><span class="mr">${short(curM)}</span></div><div id="rblk" class="rb" style="display:none"><div class="rt" onclick="toggleR()"><span class="rt-icon"><svg viewBox="0 0 24 24" width="14" height="14" fill="none" stroke="var(--ac)" stroke-width="2" stroke-linecap="round" stroke-linejoin="round"><path d="M12 2a7 7 0 0 1 5 11.9V17a2 2 0 0 1-2 2h-6a2 2 0 0 1-2-2v-3.1A7 7 0 0 1 12 2z"/><line x1="9" y1="21" x2="15" y2="21"/></svg></span><span id="rblk-label">Thinking…</span><span class="ra" id="rarr">${I.chev}</span></div><div class="rc" id="rcnt"></div></div><div class="mc" id="scnt"><span class="sd"><img src="data:image/png;base64,iVBORw0KGgoAAAANSUhEUgAAABwAAAAcCAYAAAByDd+UAAABw0lEQVR4nL2WPU5bQRDHf7M2DZIbOAISgoY7pM4Vkip1rpAul+ACHIQDUCNRQRFFstyli/2nYMZer997u2shT2O/nY//fO3sGBWSZIDlR2amghf8TfDGyMYYmTEB1Ay5fNjUmPwgoKQUf2tAI/oz1900CWfe9oBYruff07YkpUIpubdH0aS+M1P23R3lkK5HurU7zwUi55KSmW0kXQG/gRsgsevIaKa8vgb8B86BH8CjpJmZrSPd237Iww6mpAtJf9VPf1y3rGkCmPvhnqce3TWwBF6BWRbVJXAGvACLLNIN8Az8MrNVZGk/07KDRqnVRdI3SStJP2uyA+cpv3M1wKj39yx9X/xsFr9T9iSlJrCC8ob56mcJwMzWtcs+Z79+PaChT4cNRaufhMysrX6fSYnjUnoURdOcLKXg3XVywM5BnU8mddxjC0DRl9Z/Lr8GFn7vWvpg52jLm5cN9UtJbz5plpLuJM1rWdrjl2/WhFKU4FbSg6QnSfdDr0Opd8DrHeKtVAZjBTPRsDi5nPHxJNnY/Aw5M1sPArpQ+8Y17ZTx0ZTVXfVgvzkCbFS/tgjHtWlJc76B9y3CI8BwuDSVRquOvQOC88GWdq5oUQAAAABJRU5ErkJggg==" alt=""></span></div></div>`}
async function sendMessage(){
 const inp=$('uinp');let t=inp.value.trim();if(!t||stm)return;if(!cid)await createNewChat();
 if(webS&&t){const sr=await pywebview.api.web_search(t);if(sr&&!sr.startsWith('Brave API'))t=`[Web Search Results]\n${sr}\n\n[User Question]\n${t}`}
 // Inline URL fetch: detect URL at start of message
 const urlRx=/^(https?:\/\/[^\s]+)(\s|$)/;
 const urlMatch=t.match(urlRx);
 if(urlMatch&&!webS){
  const url=urlMatch[1];const rest=t.slice(url.length).trim();
  const fetched=await pywebview.api.fetch_url(url);
  if(fetched&&fetched.content){t=`[Fetched URL: ${url}]\n\n${fetched.content}\n\n${rest||'Summarize the above content.'}`;showCopyOk('URL fetched ✓')}
  else if(fetched&&fetched.error){showAlert('URL fetch failed: '+fetched.error)}
 }
 // Build file content: support multiple files
 let fc = null;
 if(Array.isArray(pf)&&pf.length>0){
  if(pf.length===1){fc=pf[0]}
  else{// Multiple files: combine as text blocks
   const parts=pf.map(f=>{
    if(f.type==='text')return `[File: ${f.filename}]\n${f.content}`;
    return `[Image: ${f.filename}]`;
   });
   fc={type:'text',filename:`${pf.length} files`,content:parts.join('\n\n---\n\n')};}
 }
 const fcj=fc?JSON.stringify(fc):null;
 clearPF();inp.value='';autoR(inp);
 const a=$('msgs');if(a.querySelector('.es'))a.innerHTML='';
 const wsMatch=t.match(/^\[Web Search Results\]\n([\s\S]*?)\n\n\[User Question\]\n/);
 const dt=wsMatch?t.replace(wsMatch[0],''):t;
 const wsData=wsMatch?wsMatch[1]:'';
 let umHtml=`<div class="msg user"><div class="mh"><span class="mr">You</span></div>`;
 if(wsData)umHtml+=`<div class="ws-toggle" onclick="this.nextElementSibling.classList.toggle('show')"><span class="i i12"><svg viewBox="0 0 24 24"><line x1="3" y1="6" x2="21" y2="6"/><line x1="3" y1="12" x2="21" y2="12"/><line x1="3" y1="18" x2="21" y2="18"/></svg></span> Web results</div><div class="ws-content">${esc(wsData)}</div>`;
 umHtml+=`<div class="mc">${rmd(dt)}</div></div>`;
 a.innerHTML+=umHtml;
  a.innerHTML+=stPH();a.scrollTop=a.scrollHeight;stm=true;activeStreams.add(cid);$('sbtn').disabled=true;$('stopbtn').style.display='flex';$('stopbtn').disabled=false;userScrolled=false;
 await pywebview.api.stream_message(cid,t,fcj)}
function isNearBottom(){const el=$('msgs');if(!el)return true;return el.scrollHeight-el.scrollTop-el.clientHeight<100}
function scrollToBottomIfNear(){if(!userScrolled){const el=$('msgs');if(el)el.scrollTop=el.scrollHeight}}
let _rblkStart=null;
function onReasoningToken(tk){const b=$('rblk'),el=$('rcnt');if(!b||!el)return;if(b.style.display==='none'||!_rblkStart){b.style.display='block';_rblkStart=Date.now();const lbl=$('rblk-label');if(lbl)lbl.textContent='Thinking…'}el.textContent+=tk;const lbl=$('rblk-label');if(lbl&&_rblkStart){const secs=Math.round((Date.now()-_rblkStart)/1000);lbl.textContent='Thought for '+secs+'s';}if(cid)streamBuffer[cid]={...streamBuffer[cid],reasoning:(streamBuffer[cid]?.reasoning||'')+tk};scrollToBottomIfNear()}
function toggleR(){const el=$('rcnt'),ar=$('rarr');if(!el)return;el.classList.toggle('show');if(ar)ar.classList.toggle('open')}
function toggleRB(rt){const rc=rt.nextElementSibling;const ra=rt.querySelector('.ra');if(!rc)return;rc.classList.toggle('show');if(ra)ra.classList.toggle('open')}
function onStreamToken(tk){const el=$('scnt');if(!el)return;if(el.querySelector('.sd')){el.innerHTML='';el._r=''}if(!el._r)el._r='';el._r+=tk;el.innerHTML=rmd(el._r);if(cid)streamBuffer[cid]={...streamBuffer[cid],content:(streamBuffer[cid]?.content||'')+tk};scrollToBottomIfNear()}
function onStreamDone(r){const doneCid=cid;stm=false;activeStreams.delete(doneCid);userScrolled=false;if(streamBuffer[doneCid])delete streamBuffer[doneCid];$('sbtn').disabled=false;$('stopbtn').style.display='none';lastU=r.usage;if(doneCid===cid){upTkB(r.usage);const dur=_rblkStart?Math.round((Date.now()-_rblkStart)/1000):0;_rblkStart=null;if(dur>0)try{pywebview.api.save_reasoning_duration(doneCid,dur);}catch(e){}pywebview.api.get_chat(doneCid).then(c=>{if(c&&cid===doneCid){renderMsgs(c.messages)}})}refreshCL()}
function onChatTitleUpdated(r){titleOverrides[r.cid]=r.title;const el=document.querySelector(`.ci[data-id="${r.cid}"] .cit`);if(el)el.textContent=r.title;}
function onStreamError(e){const errCid=cid;stm=false;activeStreams.delete(errCid);userScrolled=false;if(streamBuffer[errCid])delete streamBuffer[errCid];$('sbtn').disabled=false;$('stopbtn').style.display='none';if(cid===errCid){const el=$('scnt');if(el)el.innerHTML=`<span style="color:var(--dg)">Error: ${esc(e)}</span>`}}function onStreamStopped(){const stoppedCid=cid;stm=false;activeStreams.delete(stoppedCid);userScrolled=false;if(streamBuffer[stoppedCid])delete streamBuffer[stoppedCid];$('sbtn').disabled=false;$('stopbtn').style.display='none';if(cid===stoppedCid){const el=$('scnt');if(el)el.innerHTML+='<div style="color:var(--t3);font-size:12px;margin-top:8px;font-style:italic">Generation stopped by user</div>';pywebview.api.get_chat(stoppedCid).then(c=>{if(c&&cid===stoppedCid){const smsg=$('smsg');if(smsg)smsg.remove();renderMsgs(c.messages)}})}}async function stopGeneration(){if(!stm||!cid)return;$('stopbtn').disabled=true;await pywebview.api.stop_stream(cid)}
async function copyM(i){
 const el=$('mc-'+i);if(!el)return;
 const html=el.innerHTML;
 const styled=`<div style="font-family:Calibri,sans-serif;font-size:11pt;color:#000">${html.replace(/color:var\([^)]+\)/g,'color:#000')}</div>`;
 try{
  const blob=new Blob([styled],{type:'text/html'});
  const txt=el.innerText;
  await navigator.clipboard.write([new ClipboardItem({'text/html':blob,'text/plain':new Blob([txt],{type:'text/plain'})})]);
 }catch(e){const ta=document.createElement('textarea');ta.value=el.innerText;document.body.appendChild(ta);ta.select();document.execCommand('copy');document.body.removeChild(ta)}
 showCopyOk()}
function copyTable(btn){
 const tbl=btn.closest('.tbl-wrap').querySelector('table');if(!tbl)return;
 // Build HTML version with Calibri bold headers (always, regardless of theme)
 const rows=Array.from(tbl.querySelectorAll('tr'));
 let htmlRows='';
 rows.forEach(r=>{
  const cells=Array.from(r.querySelectorAll('th,td'));
  const isHeader=r.querySelectorAll('th').length>0;
  let htmlCells='';
  cells.forEach(c=>{
   const txt=c.innerText.replace(/</g,'&lt;').replace(/>/g,'&gt;');
   if(isHeader){
    htmlCells+=`<th style="font-family:Calibri,sans-serif;font-weight:bold;font-size:11pt;border:1px solid #ccc;padding:6px 10px;background:#f2f2f2;color:#000">${txt}</th>`;
   }else{
    htmlCells+=`<td style="font-family:Calibri,sans-serif;font-size:11pt;border:1px solid #ccc;padding:6px 10px;color:#000">${txt}</td>`;
   }
  });
  htmlRows+=`<tr>${htmlCells}</tr>`;
 });
 const htmlTable=`<table style="border-collapse:collapse;font-family:Calibri,sans-serif">${htmlRows}</table>`;
 const tsv=rows.map(r=>Array.from(r.querySelectorAll('th,td')).map(c=>c.innerText.replace(/\t/g,' ')).join('\t')).join('\n');
 if(navigator.clipboard&&navigator.clipboard.write){
  const htmlBlob=new Blob([htmlTable],{type:'text/html'});
  const txtBlob=new Blob([tsv],{type:'text/plain'});
  navigator.clipboard.write([new ClipboardItem({'text/html':htmlBlob,'text/plain':txtBlob})]).then(()=>showCopyOk('Table copied')).catch(()=>{fallbackCopy(tsv);showCopyOk('Table copied')});
 }else if(navigator.clipboard&&navigator.clipboard.writeText){
  navigator.clipboard.writeText(tsv).then(()=>showCopyOk('Table copied')).catch(()=>{fallbackCopy(tsv);showCopyOk('Table copied')});
 }else{fallbackCopy(tsv);showCopyOk('Table copied')}
}
function showCopyOk(msg){const d=document.createElement('div');d.className='copy-ok';d.textContent=msg||'Copied to clipboard';document.body.appendChild(d);setTimeout(()=>d.remove(),1600)}

/* ── IMAGE FOLDER PANEL ── */
var _ifpFiles=[];
function toggleImgFolderPanel(){
 const panel=$('img-folder-panel'),btn=$('img-folder-btn');
 const open=panel.classList.toggle('open');
 if(btn)btn.classList.toggle('active',open);
}
async function loadImageFolder(){
 const result=await pywebview.api.upload_folder();
 if(!result)return;
 if(result.error){showAlert(result.error);return;}
 _ifpFiles=result.files||[];
 const folderName=result.folder||'Images';
 $('ifp-title').textContent=folderName;
 $('ifp-count').textContent=_ifpFiles.length?`${_ifpFiles.length} file${_ifpFiles.length!==1?'s':''}`:' No files found';
 renderIFPList();
}
function renderIFPList(){
 const list=$('ifp-list');if(!list)return;
 if(!_ifpFiles.length){list.innerHTML='<div style="text-align:center;padding:20px;color:var(--t3);font-size:12px">No images or PDFs in this folder</div>';return;}
 const pdfIcon=`<svg viewBox="0 0 24 24" width="16" height="16" fill="none" stroke="#e05555" stroke-width="1.8" stroke-linecap="round" stroke-linejoin="round"><path d="M14 2H6a2 2 0 00-2 2v16a2 2 0 002 2h12a2 2 0 002-2V8z"/><polyline points="14 2 14 8 20 8"/><line x1="9" y1="13" x2="15" y2="13"/><line x1="9" y1="17" x2="12" y2="17"/></svg>`;
 list.innerHTML=_ifpFiles.map((f,i)=>`<div class="ifp-item" data-ifp-idx="${i}" onmouseenter="showIFPPreview(event,${i})" onmouseleave="hideIFPPreview()">${f.type==='pdf'?`<div class="ifp-pdf-thumb">${pdfIcon}</div>`:`<img class="ifp-thumb" src="" onerror="this.style.display='none'">`}<span class="ifp-name" title="${esc(f.filename)}">${esc(f.filename)}</span></div>`).join('');
 list.querySelectorAll('.ifp-thumb').forEach(img=>{
  const idx=parseInt(img.closest('.ifp-item').getAttribute('data-ifp-idx'));
  const f=_ifpFiles[idx];
  if(f&&f.type==='image')img.src=`data:${f.media_type};base64,${f.base64}`;
 });
}
function showIFPPreview(e,idx){
 const f=_ifpFiles[idx];if(!f)return;
 const prev=$('ifp-preview');
 const imgEl=$('ifp-preview-img');
 const embedEl=$('ifp-preview-pdf');
 if(f.type==='pdf'){
  imgEl.style.display='none';
  embedEl.src=`data:application/pdf;base64,${f.base64}`;
  embedEl.style.display='block';
 } else {
  embedEl.style.display='none';
  embedEl.src='';
  imgEl.src=`data:${f.media_type};base64,${f.base64}`;
  imgEl.style.display='block';
 }
 $('ifp-preview-name').textContent=f.filename;
 const kb=f.size?Math.round(f.size/1024)+'KB':'';
 const pages=f.page_count?` · ${f.page_count} page${f.page_count!==1?'s':''}`:''
 $('ifp-preview-meta').innerHTML=`Type: ${f.media_type}${pages}<br>${kb?'Size: '+kb:''}`;
 const panel=$('img-folder-panel');
 const panelRect=panel.getBoundingClientRect();
 const itemRect=e.currentTarget.getBoundingClientRect();
 const previewW=Math.min(728,panelRect.left-16);
 prev.style.width=previewW+'px';
 prev.style.right=(window.innerWidth-panelRect.left+8)+'px';
 const top=Math.min(itemRect.top,window.innerHeight-prev.offsetHeight-16);
 prev.style.top=Math.max(8,top)+'px';
 prev.classList.add('show');
}
function hideIFPPreview(){$('ifp-preview').classList.remove('show')}

// Post-process: wrap tables with copy button
function wrapTables(container){
 container.querySelectorAll('table').forEach(t=>{
  if(t.parentElement.classList.contains('tbl-wrap'))return;
  const w=document.createElement('div');w.className='tbl-wrap';
  t.parentElement.insertBefore(w,t);w.appendChild(t);
  const btn=document.createElement('button');btn.className='cp-tbl';btn.textContent='Copy table';btn.onclick=function(){copyTable(this)};
  w.appendChild(btn);
 })}
// Context menu
(function(){
 const ctx=$('ctx');
 document.addEventListener('contextmenu',function(e){
  const msg=e.target.closest('.msg');
  const inp=e.target.closest('#uinp');
  const ci=e.target.closest('.ci');
  if(!msg && !inp && !ci){ctx.classList.remove('show');return}
  e.preventDefault();
  let items='';
  if(ci){
   const id=ci.getAttribute('data-id');
   const isPinned=ci.querySelector('.cip')!==null;
   items+=`<div class="ctx-item" onclick="togglePin('${id}');closeCtx()">${isPinned?I.pin:I.pinO}<span style="margin-left:8px">${isPinned?'Unpin':'Pin'}</span></div>`;
   items+=`<div class="ctx-item" onclick="renameChat('${id}');closeCtx()">${I.edit}<span style="margin-left:8px">Rename</span></div>`;
   items+=`<div class="ctx-item" onclick="moveChatToFolder('${id}');closeCtx()"><span class="i i12"><svg viewBox="0 0 24 24"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg></span><span style="margin-left:8px">Move to folder</span></div>`;
   items+=`<div class="ctx-item" style="color:var(--dg)" onclick="deleteChat('${id}');closeCtx()">${I.del}<span style="margin-left:8px">Delete</span></div>`;
  }else if(msg){
   const idx=parseInt(msg.getAttribute('data-idx'));const iu=msg.classList.contains('user');
   const sel=window.getSelection().toString().trim();
   _ctxSel=sel;
   if(sel){
    items+=`<div class="ctx-item" onclick="copySelection();closeCtx()">Copy selection</div>`;
   }else{
    items+=`<div class="ctx-item" onclick="copyM(${idx});closeCtx()">Copy message</div>`;
    if(iu)items+=`<div class="ctx-item" onclick="editM(${idx});closeCtx()">Edit & resend</div>`;
    if(!iu)items+=`<div class="ctx-item" onclick="branchC(${idx});closeCtx()">Branch chat</div>`;
    const tbl=e.target.closest('table')||msg.querySelector('table');
    if(tbl)items+=`<div class="ctx-item" onclick="copyTableFromCtx();closeCtx()">Copy table</div>`;
    ctx._tbl=tbl;
   }
  }else if(inp){
   const ta=$('uinp');const hasSel=ta.selectionStart!==ta.selectionEnd;
   if(hasSel)items+=`<div class="ctx-item" onclick="inputCtx('cut');closeCtx()">Cut</div>`;
   if(hasSel)items+=`<div class="ctx-item" onclick="inputCtx('copy');closeCtx()">Copy</div>`;
   items+=`<div class="ctx-item" onclick="inputCtx('paste');closeCtx()">Paste</div>`;
   if(hasSel)items+=`<div class="ctx-item" onclick="inputCtx('delete');closeCtx()">Delete</div>`;
   items+=`<div class="ctx-sep"></div>`;
   items+=`<div class="ctx-item" onclick="inputCtx('selectAll');closeCtx()">Select All</div>`;
   ctx._inp=ta;
  }
  ctx.innerHTML=items;
  let x=e.clientX,y=e.clientY;
  ctx.style.left=x+'px';ctx.style.top=y+'px';ctx.classList.add('show');
  requestAnimationFrame(()=>{
   if(x+ctx.offsetWidth>window.innerWidth)ctx.style.left=(window.innerWidth-ctx.offsetWidth-4)+'px';
   if(y+ctx.offsetHeight>window.innerHeight)ctx.style.top=(window.innerHeight-ctx.offsetHeight-4)+'px';
  });
 });
 document.addEventListener('click',()=>ctx.classList.remove('show'));
})();
function inputCtx(action){
 const ta=$('uinp');if(!ta)return;
 const start=ta.selectionStart,end=ta.selectionEnd;
 if(action==='cut'&&start!==end){
  const text=ta.value.substring(start,end);
  const doCut=()=>{ta.value=ta.value.substring(0,start)+ta.value.substring(end);ta.selectionStart=ta.selectionEnd=start;ta.focus();autoR(ta)};
  if(navigator.clipboard&&navigator.clipboard.writeText){navigator.clipboard.writeText(text).then(doCut).catch(()=>{fallbackCopy(text);doCut()})}
  else{fallbackCopy(text);doCut()}
 }
 else if(action==='copy'&&start!==end){
  const text=ta.value.substring(start,end);
  if(navigator.clipboard&&navigator.clipboard.writeText){navigator.clipboard.writeText(text).catch(()=>fallbackCopy(text))}
  else{fallbackCopy(text)}
 }
 else if(action==='paste'){
  const curStart=ta.selectionStart,curEnd=ta.selectionEnd;
  const doPaste=(t)=>{ta.value=ta.value.substring(0,curStart)+t+ta.value.substring(curEnd);const newPos=curStart+t.length;ta.selectionStart=ta.selectionEnd=newPos;ta.focus();autoR(ta)};
  if(navigator.clipboard&&navigator.clipboard.readText){navigator.clipboard.readText().then(doPaste).catch(()=>fallbackPaste(doPaste))}
  else{fallbackPaste(doPaste)}
 }
 else if(action==='delete'&&start!==end){ta.value=ta.value.substring(0,start)+ta.value.substring(end);ta.selectionStart=ta.selectionEnd=start;ta.focus();autoR(ta)}
 else if(action==='selectAll'){ta.selectionStart=0;ta.selectionEnd=ta.value.length;ta.focus();}
}
function fallbackCopy(text){
 const tmp=document.createElement('textarea');tmp.value=text;tmp.style.position='fixed';tmp.style.opacity='0';document.body.appendChild(tmp);tmp.select();
 try{document.execCommand('copy')}catch(e){}document.body.removeChild(tmp);
}
function fallbackPaste(callback){
 const tmp=document.createElement('textarea');tmp.style.position='fixed';tmp.style.opacity='0';document.body.appendChild(tmp);tmp.focus();
 try{if(document.execCommand('paste')){callback(tmp.value)}else{const saved=tmp.value;if(saved)callback(saved)}}catch(e){}document.body.removeChild(tmp);
}
function closeCtx(){$('ctx').classList.remove('show')}
function copyTableFromCtx(){const tbl=$('ctx')._tbl;if(tbl){const btn=tbl.closest('.tbl-wrap')?.querySelector('.cp-tbl');if(btn)copyTable(btn);else{const ta=document.createElement('textarea');ta.value=tbl.innerText;document.body.appendChild(ta);ta.select();document.execCommand('copy');document.body.removeChild(ta);showCopyOk('Table copied')}}}
function copySelection(){const t=_ctxSel||window.getSelection().toString();if(t){if(navigator.clipboard&&navigator.clipboard.writeText){navigator.clipboard.writeText(t).catch(()=>fallbackCopy(t))}else{fallbackCopy(t)}showCopyOk()}}
async function pasteToInput(){try{if(navigator.clipboard&&navigator.clipboard.readText){const t=await navigator.clipboard.readText();if(t){$('uinp').value+=(($('uinp').value?' ':'')+t);autoR($('uinp'))}}else{fallbackPaste(function(t){if(t){$('uinp').value+=(($('uinp').value?' ':'')+t);autoR($('uinp'))}})}}catch(e){}}

async function branchC(i){const nc=await pywebview.api.branch_chat(cid,i);if(nc){cid=nc.id;renderMsgs(nc.messages);refreshCL()}}
function restoreStreamingUI(){const a=$('msgs');if(!a||a.querySelector('#smsg'))return;a.innerHTML+=stPH();const rblk=$('rblk'),rcnt=$('rcnt'),scnt=$('scnt');if(rblk&&rcnt&&streamBuffer[cid]?.reasoning){rblk.style.display='block';rcnt.textContent=streamBuffer[cid].reasoning}if(scnt&&streamBuffer[cid]?.content){scnt.innerHTML=rmd(streamBuffer[cid].content);scnt._r=streamBuffer[cid].content}$('stopbtn').style.display='flex';$('stopbtn').disabled=false;if(!userScrolled)a.scrollTop=a.scrollHeight}
async function regen(){if(stm||!cid)return;$('msgs').innerHTML+=stPH();$('msgs').scrollTop=$('msgs').scrollHeight;stm=true;activeStreams.add(cid);$('sbtn').disabled=true;$('stopbtn').style.display='flex';$('stopbtn').disabled=false;userScrolled=false;await pywebview.api.stream_regenerate(cid)}
async function editM(i){const c=await pywebview.api.get_chat(cid);if(!c||!c.messages[i])return;const ct=c.messages[i].content;const t=typeof ct==='string'?ct:(Array.isArray(ct)?ct.filter(x=>x.type==='text').map(x=>x.text).join('\n'):String(ct));const el=$('mc-'+i);if(!el)return;el.innerHTML=`<textarea class="ea" id="eta-${i}">${esc(t)}</textarea><div class="ebs"><button class="esv" onclick="subEdit(${i})">Save & Resend</button><button class="ecn" onclick="loadChat('${cid}')">Cancel</button></div>`;$('eta-'+i).focus()}
async function subEdit(i){const ta=$('eta-'+i);if(!ta)return;const nt=ta.value.trim();if(!nt)return;const a=$('msgs');const msgDivs=a.querySelectorAll('.msg');for(const md of msgDivs){const di=parseInt(md.dataset.idx);if(di>=i){const wrap=md.parentElement;if(wrap&&wrap.classList.contains('user-wrap'))wrap.remove();else md.remove()}}const wsMatch=nt.match(/^\[Web Search Results\]\n([\s\S]*?)\n\n\[User Question\]\n/);const dt=wsMatch?nt.replace(wsMatch[0],''):nt;const wsData=wsMatch?wsMatch[1]:'';let umHtml=`<div class="msg user"><div class="mh"><span class="mr">You</span></div>`;if(wsData)umHtml+=`<div class="ws-toggle" onclick="this.nextElementSibling.classList.toggle('show')"><span class="i i12"><svg viewBox="0 0 24 24"><line x1="3" y1="6" x2="21" y2="6"/><line x1="3" y1="12" x2="21" y2="12"/><line x1="3" y1="18" x2="21" y2="18"/></svg></span> Web results</div><div class="ws-content">${esc(wsData)}</div>`;umHtml+=`<div class="mc">${rmd(dt)}</div></div>`;a.innerHTML+=umHtml;a.innerHTML+=stPH();stm=true;activeStreams.add(cid);$('sbtn').disabled=true;$('stopbtn').style.display='flex';$('stopbtn').disabled=false;userScrolled=false;await pywebview.api.stream_edit(cid,i,nt)}
function _renderPFBadges(){
 const list=$('fb-list');if(!list)return;
 list.innerHTML='';
 pf.forEach((r,i)=>{
  const ic=r.type==='image'?I.img:I.file;
  const badge=document.createElement('div');
  badge.className='fb';
  badge.setAttribute('data-pf-idx',i);
  badge.innerHTML=`${ic} ${esc(r.filename)} <button class="pf-rm-btn" data-pf-idx="${i}" title="Remove" style="background:none;border:none;color:var(--dg);cursor:pointer;font-size:16px;line-height:1;padding:0 2px;pointer-events:all">×</button>`;
  list.appendChild(badge);
 });
}
async function uploadFile(){
 const results=await pywebview.api.upload_file(cid||'');
 if(!results||!Array.isArray(results)||results.length===0)return;
 if(!Array.isArray(pf))pf=[];
 for(const r of results){pf.push(r);}
 _renderPFBadges();
}
function removePF(idx){
 if(!Array.isArray(pf))return;
 const badge=document.querySelector('[data-pf-idx="'+idx+'"]');
 if(badge){badge.style.display='none';badge.remove();}
 pf.splice(idx,1);
 _renderPFBadges();
}
function clearPF(){pf=[];const list=$('fb-list');if(list)list.innerHTML=''}
function toggleSearch(){webS=!webS;updatePlusState()}
function updatePlusState(){
 const pb=$('plusbtn');const wo=$('wsopt');if(pb)pb.classList.toggle('active',webS);if(wo)wo.classList.toggle('active',webS);
 const cpb=$('cmp-plusbtn');const cwo=$('cmp-wsopt');if(cpb)cpb.classList.toggle('active',webS);if(cwo)cwo.classList.toggle('active',webS);
}
function togglePlusMenu(e){e.stopPropagation();const m=$('plusmenu');const pb=$('plusbtn');const open=!m.classList.contains('show');m.classList.toggle('show',open);pb.classList.toggle('open',open)}
function closePlusMenu(){$('plusmenu').classList.remove('show');$('plusbtn').classList.remove('open')}
function toggleCmpPlusMenu(e){e.stopPropagation();const m=$('cmp-plusmenu');const pb=$('cmp-plusbtn');const open=!m.classList.contains('show');m.classList.toggle('show',open);pb.classList.toggle('open',open)}
function closeCmpPlusMenu(){const m=$('cmp-plusmenu');const pb=$('cmp-plusbtn');if(m)m.classList.remove('show');if(pb)pb.classList.remove('open')}
document.addEventListener('click',e=>{if(!e.target.closest('.plus-wrap'))closePlusMenu();if(!e.target.closest('.cmp-ib .plus-wrap'))closeCmpPlusMenu()})
async function toggleVoice(){
 const b=$('vbtn');const cb=$('cmp-vbtn');
 if(isRec){
  if(b)b.classList.remove('recording');if(cb)cb.classList.remove('recording');isRec=false;
  const r=await pywebview.api.stop_voice_record();
  if(r&&r.text){$('uinp').value=$('uinp').value+(($('uinp').value?' ':'')+r.text);autoR($('uinp'))}
  else if(r&&r.error){showAlert(r.error)}
  else{showAlert('No speech detected.')}
  return}
 const res=await pywebview.api.start_voice_record();
 if(res&&res.error){showAlert(res.error);return}
 if(b)b.classList.add('recording');if(cb)cb.classList.add('recording');isRec=true}
function toggleMD(){const d=$('mdd'),b=$('mbtn');d.classList.toggle('show');b.classList.toggle('open');if(d.classList.contains('show'))renderMD()}
function renderMD(){$('mdd').innerHTML=models.map(m=>`<div class="mdi ${m===curM?'active':''}" onclick="selM('${esc(m)}')"><span>${esc(m)}</span><button class="rm" onclick="event.stopPropagation();remM('${esc(m)}')" title="Remove">×</button></div>`).join('')+`<div class="ma"><input type="text" id="nmi" placeholder="org/model-name" onkeydown="if(event.key==='Enter')addM()"><button onclick="addM()">Add</button></div>`}
async function selM(m){curM=m;$('mname').textContent=short(m);await pywebview.api.set_current_model(m);toggleMD()}
async function addM(){const inp=$('nmi');const n=inp.value.trim();if(!n)return;models=await pywebview.api.add_model(n);inp.value='';renderMD()}
async function remM(m){models=await pywebview.api.remove_model(m);if(curM===m&&models.length){curM=models[0];await pywebview.api.set_current_model(curM);$('mname').textContent=short(curM)}renderMD()}
document.addEventListener('click',e=>{if(!e.target.closest('.ms')){$('mdd').classList.remove('show');$('mbtn').classList.remove('open')}});
function closeM(id){$(id).classList.remove('show')}
function openSet(){$('setm').classList.add('show');upTh();const cfg=window._cfg||{}}
async function saveSet(){const c={api_key:$('sak').value.trim(),brave_api_key:$('sbk').value.trim(),base_url:$('sbu').value.trim(),max_tokens:parseInt($('smt').value),max_memory:parseInt($('smm').value),theme:selTh};maxTk=c.max_tokens;await pywebview.api.save_settings(c);closeM('setm')}

// Keyboard shortcuts
document.addEventListener('keydown',function(e){
 const inInput=e.target.closest('#uinp,input,textarea');
 if(e.ctrlKey||e.metaKey){
  if(e.key==='n'){e.preventDefault();createNewChat()}
  if(e.key==='/'&&!e.shiftKey){e.preventDefault();openSet()}
  if(e.key==='f'){e.preventDefault();$('sch').focus()}
  if(e.key==='b'){e.preventDefault();toggleSidebar()}
  if(e.key==='e'){e.preventDefault();openExp()}
  if(e.key==='w'&&!inInput){e.preventDefault();toggleSearch();updatePlusState()}
  if(e.key==='u'){e.preventDefault();uploadFile()}
  if(e.key==='m'){e.preventDefault();openCompare()}
 }
 if(e.key==='Escape'){closeAllModals();closeCompare()}
 if(e.key==='?'&&!inInput){e.preventDefault();$('kbm').classList.add('show')}
});

function closeAllModals(){document.querySelectorAll('.mo').forEach(m=>m.classList.remove('show'))}

/* ── THEME EDITOR ── */
function openThemeEditor(){
 restoreThemeOverrides();
 const r=getComputedStyle(document.documentElement);
 const acc=rgbToHex(r.getPropertyValue('--ac').trim())||'#e88a2a';
 $('te-accent').value=acc;$('te-accent-val').textContent=acc;
 const fs=parseInt(r.getPropertyValue('--chat-font-size')||'17');
 $('te-font').value=fs;$('te-font-val').textContent=fs+'px';
 const br=parseInt(_teDynRadius||'80');
 $('te-radius').value=br;$('te-radius-val').textContent=br+'%';
 $('them').classList.add('show');
}
function rgbToHex(rgb){
 if(!rgb||rgb.startsWith('#'))return rgb||null;
 const m=rgb.match(/rgb\((\d+),\s*(\d+),\s*(\d+)\)/);
 if(!m)return null;
 return '#'+[m[1],m[2],m[3]].map(x=>parseInt(x).toString(16).padStart(2,'0')).join('');
}
function applyTeAccent(v){
 document.documentElement.style.setProperty('--ac',v);
 document.documentElement.style.setProperty('--acd',adjustColor(v,-20));
 document.documentElement.style.setProperty('--acg',v+'26');
 $('te-accent-val').textContent=v;
}
let _teDynFont='',_teDynRadius='';
function _updateDynStyle(){
 let s=document.getElementById('te-dynamic');
 if(!s){s=document.createElement('style');s.id='te-dynamic';document.head.appendChild(s);}
 let css='';
 if(_teDynFont)css+='.msg,.mc{font-size:'+_teDynFont+' !important}';
 if(_teDynRadius)css+='.msg.assistant{max-width:'+_teDynRadius+' !important}';
 s.textContent=css;
}
function applyTeFont(v){
 _teDynFont=v+'px';_updateDynStyle();
 document.documentElement.style.setProperty('--chat-font-size',v+'px');
 $('te-font-val').textContent=v+'px';
}
function applyTeRadius(v){
 _teDynRadius=v+'%';_updateDynStyle();
 document.querySelectorAll('.msg.assistant').forEach(el=>el.style.setProperty('max-width',v+'%','important'));
 $('te-radius-val').textContent=v+'%';
}
function adjustColor(hex,amt){
 let c=parseInt(hex.slice(1),16);
 let r=Math.min(255,Math.max(0,((c>>16)&255)+amt));
 let g=Math.min(255,Math.max(0,((c>>8)&255)+amt));
 let b=Math.min(255,Math.max(0,(c&255)+amt));
 return '#'+[r,g,b].map(x=>x.toString(16).padStart(2,'0')).join('');
}
function saveThemeEditor(){
 const overrides={};
 const ac=document.documentElement.style.getPropertyValue('--ac');if(ac)overrides.ac=ac;
 const acd=document.documentElement.style.getPropertyValue('--acd');if(acd)overrides.acd=acd;
 const acg=document.documentElement.style.getPropertyValue('--acg');if(acg)overrides.acg=acg;
 const fs=document.documentElement.style.getPropertyValue('--chat-font-size');if(fs)overrides.chatFontSize=fs;
 if(_teDynFont)overrides.dynFont=_teDynFont;
 if(_teDynRadius)overrides.dynRadius=_teDynRadius;
 try{localStorage.setItem('themeOverrides',JSON.stringify(overrides));}catch(e){}
 try{pywebview.api.update_theme_overrides(JSON.stringify(overrides));}catch(e){}
 showCopyOk('Theme applied ✓');closeM('them');
}
function resetThemeEditor(){
 document.documentElement.style.removeProperty('--ac');
 document.documentElement.style.removeProperty('--acd');
 document.documentElement.style.removeProperty('--acg');
 document.documentElement.style.removeProperty('--chat-font-size');
 document.querySelectorAll('.msg.assistant').forEach(el=>el.style.removeProperty('max-width'));
 _teDynFont='';_teDynRadius='';_updateDynStyle();
 try{localStorage.removeItem('themeOverrides');}catch(e){}
 try{pywebview.api.update_theme_overrides('{}');}catch(e){}
 openThemeEditor();
}
function restoreThemeOverrides(){
 try{
  const raw=localStorage.getItem('themeOverrides');
  if(!raw)return;
  const o=JSON.parse(raw);
  if(o.ac){document.documentElement.style.setProperty('--ac',o.ac);}
  if(o.acd){document.documentElement.style.setProperty('--acd',o.acd);}
  if(o.acg){document.documentElement.style.setProperty('--acg',o.acg);}
  if(o.chatFontSize){document.documentElement.style.setProperty('--chat-font-size',o.chatFontSize);}
  if(o.dynFont){_teDynFont=o.dynFont;}
  if(o.dynRadius){_teDynRadius=o.dynRadius;}
  _updateDynStyle();
 }catch(e){}
}

/* ── MODEL COMPARE ── */
let cmpActive=false;
function openCompare(){
 if(!models.length){showAlert('No models configured.');return}
 const sel1=$('cmp-m1'),sel2=$('cmp-m2');
 sel1.innerHTML=sel2.innerHTML=models.map(m=>`<option value="${esc(m)}">${esc(m.split('/').pop())}</option>`).join('');
 if(models.length>1)sel2.selectedIndex=1;
 $('cmp-label1').textContent=short(models[0]);
 $('cmp-label2').textContent=short(models[Math.min(1,models.length-1)]);
 sel1.onchange=()=>$('cmp-label1').textContent=short(sel1.value);
 sel2.onchange=()=>$('cmp-label2').textContent=short(sel2.value);
 $('cmp-body1').innerHTML='<div style="color:var(--t3);text-align:center;padding:40px 20px;font-size:13px">Enter a prompt below to compare</div>';
 $('cmp-body2').innerHTML='<div style="color:var(--t3);text-align:center;padding:40px 20px;font-size:13px">Enter a prompt below to compare</div>';
 $('compare-overlay').classList.add('show');cmpActive=true;
 const ta=$('cmp-inp');if(ta){ta.value='';requestAnimationFrame(()=>ta.focus())}
}
function closeCompare(){$('compare-overlay').classList.remove('show');cmpActive=false}
async function runCompare(){
 const ta=$('cmp-inp');if(!ta)return;
 const ut=ta.value.trim();if(!ut){showAlert('Enter a prompt first.');return}
 ta.value='';autoR(ta);
 const m1=$('cmp-m1').value,m2=$('cmp-m2').value;
 if(!m1||!m2){showAlert('Select two models.');return}
 if(!cid)await createNewChat();
 // Build file content from cmpPf (same logic as main send)
 let fc=null;
 if(Array.isArray(cmpPf)&&cmpPf.length>0){
  if(cmpPf.length===1){fc=cmpPf[0];}
  else{const parts=cmpPf.map(f=>f.type==='text'?`[File: ${f.filename}]\n${f.content}`:`[Image: ${f.filename}]`);fc={type:'text',filename:`${cmpPf.length} files`,content:parts.join('\n\n---\n\n')};}
 }
 const fcj=fc?JSON.stringify(fc):null;
 cmpPf=[];_renderCmpPFBadges();
 const spinHTML='<div style="color:var(--t3);padding:20px;font-size:13px">⏳ Generating…</div>';
 $('cmp-body1').innerHTML=spinHTML;$('cmp-body2').innerHTML=spinHTML;
 $('cmp-body1')._r='';$('cmp-body2')._r='';
 await pywebview.api.stream_compare(cid,ut,fcj,m1,m2);
}
function onCompareToken(d){
 const pane=$(`cmp-body${d.slot}`);if(!pane)return;
 if(!pane._r)pane._r='';
 pane._r+=d.token;
 pane.innerHTML=`<div class="mc">${rmd(pane._r)}</div>`;
 pane.scrollTop=pane.scrollHeight;
}
function onCompareDone(d){
 const pane=$(`cmp-body${d.slot}`);if(!pane)return;
 pane.innerHTML=`<div class="mc">${rmd(pane._r||d.text||'')}</div>`;
}
function onCompareError(d){
 const pane=$(`cmp-body${d.slot}`);if(!pane)return;
 pane.innerHTML=`<div style="color:var(--dg);padding:12px;font-size:13px">Error: ${esc(d.error)}</div>`;
}

/* ── COMPARE FILE ATTACHMENTS (global scope) ── */
var cmpPf=[];
var _tplForCompare=false;
function _renderCmpPFBadges(){
 const list=$('cmp-fb-list');if(!list)return;
 list.innerHTML='';
 cmpPf.forEach((r,i)=>{
  const ic=r.type==='image'?I.img:I.file;
  const badge=document.createElement('div');
  badge.className='fb';
  badge.innerHTML=`${ic} ${esc(r.filename)} <button onclick="removeCmpPF(${i})" style="background:none;border:none;color:var(--dg);cursor:pointer;font-size:16px;line-height:1;padding:0 2px">×</button>`;
  list.appendChild(badge);
 });
}
function removeCmpPF(idx){cmpPf.splice(idx,1);_renderCmpPFBadges();}
async function uploadFileForCompare(){
 const results=await pywebview.api.upload_file('');
 if(!results||!Array.isArray(results)||results.length===0)return;
 for(const r of results)cmpPf.push(r);
 _renderCmpPFBadges();
}
function openTemplatesForCompare(){_tplForCompare=true;openTemplates();}

// Chat filtering
let allChats=[];
async function filterChats(query){const el=$('chatlist');if(!query){refreshCL();return}const chats=await pywebview.api.list_chats();allChats=chats;let filtered=chats.filter(c=>(titleOverrides[c.id]||c.title).toLowerCase().includes(query.toLowerCase()));if(activeFolderTab!=='All')filtered=filtered.filter(c=>c.folder===activeFolderTab);if(!filtered.length){el.innerHTML='<div style="text-align:center;padding:20px;color:var(--t3);font-size:12px">No chats found</div>';return}el.innerHTML=filtered.map(c=>{const title=titleOverrides[c.id]||c.title;const fc=folderColors[c.folder]||'#555566';const fdot=c.folder&&c.folder!=='General'?`<span class="ci-fdot" style="background:${fc}"></span>`:'';return`<div class="ci ${c.id===cid?'active':''}" onclick="loadChat('${c.id}')" data-tip="${esc(title)}" data-id="${c.id}">${c.pinned?`<button class="cip" onclick="event.stopPropagation();togglePin('${c.id}')" title="Unpin">`+I.pin+'</button>':''}${c.branched?'<span class="ci-branch">'+I.branch+'</span>':''}${fdot}<span class="cit">${esc(title)}</span><span class="cia"><div class="chat-menu"><button onclick="event.stopPropagation();toggleChatMenu('${c.id}')" title="Actions">${I.menu}</button><div class="chat-menu-dropdown" id="cm-${c.id}"><div class="chat-menu-item" onclick="event.stopPropagation();togglePin('${c.id}');hideChatMenu('${c.id}')">${c.pinned?I.pin:I.pinO}<span>${c.pinned?'Unpin':'Pin'}</span></div><div class="chat-menu-item" onclick="event.stopPropagation();renameChat('${c.id}');hideChatMenu('${c.id}')">${I.edit}<span>Rename</span></div><div class="chat-menu-item" onclick="event.stopPropagation();moveChatToFolder('${c.id}');hideChatMenu('${c.id}')"><span class="i i12"><svg viewBox="0 0 24 24"><path d="M22 19a2 2 0 01-2 2H4a2 2 0 01-2-2V5a2 2 0 012-2h5l2 3h9a2 2 0 012 2z"/></svg></span><span>Move to folder</span></div><div class="chat-menu-item delete" onclick="event.stopPropagation();deleteChat('${c.id}');hideChatMenu('${c.id}')">${I.del}<span>Delete</span></div></div></div></span></div>`}).join('')}

// Event listeners for range inputs - set up after DOM ready
function setupRangeListeners(){const smt=$('smt'),smm=$('smm');if(smt)smt.addEventListener('input',function(){$('smtv').textContent=this.value});if(smm)smm.addEventListener('input',function(){$('smmv').textContent=this.value})}
async function openSP(){if(!cid){showAlert('Open a chat first.');return}const c=await pywebview.api.get_chat(cid);if(!c)return;$('spi').value=c.system_prompt||c.messages[0]?.content||'';$('spm').classList.add('show')}
async function saveSP(){const v=$('spi').value.trim();if(cid&&v)await pywebview.api.update_system_prompt(cid,v);closeM('spm')}
function openExp(){if(!cid){showAlert('Open a chat first.');return}$('expm').classList.add('show')}
async function doExp(f){if(!cid)return;const r=await pywebview.api.save_export(cid,f);if(r){closeM('expm');showAlert('Exported to: '+r)}else{showAlert('Export failed. Make sure the chat has messages and try again.')}}
const ta=$('uinp');ta.addEventListener('input',function(){autoR(this)});ta.addEventListener('keydown',function(e){if(e.key==='Enter'&&!e.shiftKey){e.preventDefault();sendMessage()}});ta.addEventListener('paste',function(e){e.stopPropagation();e.preventDefault();const txt=(e.clipboardData||window.clipboardData).getData('text');if(txt){const cleaned=txt.replace(/\n\s*\n+/g,'\n').replace(/^\s+|\s+$/g,'');const start=ta.selectionStart,end=ta.selectionEnd;ta.value=ta.value.substring(0,start)+cleaned+ta.value.substring(end);ta.selectionStart=ta.selectionEnd=start+cleaned.length;autoR(ta)}});
const cmpTa=$('cmp-inp');if(cmpTa){cmpTa.addEventListener('paste',function(e){e.stopPropagation();e.preventDefault();const txt=(e.clipboardData||window.clipboardData).getData('text');if(txt){const cleaned=txt.replace(/\n\s*\n+/g,'\n').replace(/^\s+|\s+$/g,'');const start=cmpTa.selectionStart,end=cmpTa.selectionEnd;cmpTa.value=cmpTa.value.substring(0,start)+cleaned+cmpTa.value.substring(end);cmpTa.selectionStart=cmpTa.selectionEnd=start+cleaned.length;autoR(cmpTa)}});}
function autoR(el){el.style.height='auto';const sh=el.scrollHeight;const h=Math.min(sh,200);el.style.height=h+'px';el.style.overflow=sh>200?'auto':'hidden'}

// Global tooltip system - fixed position, never clipped, no blinking
(function(){
const tip=document.createElement('div');tip.id='gtip';document.body.appendChild(tip);
let hideTimer=null;
function showTip(el){
 if(hideTimer){clearTimeout(hideTimer);hideTimer=null}
 const t=el.getAttribute('data-tip');if(!t)return;
 tip.textContent=t;tip.style.display='block';
 const r=el.getBoundingClientRect();
 let top=r.top-tip.offsetHeight-6;let left=r.left+r.width/2-tip.offsetWidth/2;
 if(top<4)top=r.bottom+6;
 if(left<4)left=4;if(left+tip.offsetWidth>window.innerWidth-4)left=window.innerWidth-tip.offsetWidth-4;
 tip.style.left=left+'px';tip.style.top=top+'px';tip.style.opacity='1';
}
function hideTip(){hideTimer=setTimeout(()=>{tip.style.opacity='0';setTimeout(()=>{tip.style.display='none'},120)},80)}
document.addEventListener('mouseover',function(e){
 const el=e.target.closest('[title]');
 if(el&&el.getAttribute('title')){el.setAttribute('data-tip',el.getAttribute('title'));el.removeAttribute('title')}
 const te=e.target.closest('[data-tip]');
 if(te)showTip(te);
});
document.addEventListener('mouseout',function(e){
 const te=e.target.closest('[data-tip]');
 if(te){const rt=e.relatedTarget;if(!rt||!te.contains(rt))hideTip()}
});
})();
/* TTS */
const ttsVoices=[{id:'af_bella',name:'Bella',gender:'F 🇺🇸'},{id:'af_sarah',name:'Sarah',gender:'F 🇺🇸'},{id:'am_adam',name:'Adam',gender:'M 🇺🇸'},{id:'am_michael',name:'Michael',gender:'M 🇺🇸'},{id:'bf_isabella',name:'Isabella',gender:'F 🇬🇧'},{id:'bf_emma',name:'Emma',gender:'F 🇬🇧'}];
let ttsVoice='af_bella';let ttsAudio=null;let ttsBtnEl=null;let ttsIdx=null;let ttsLoading=false;
function openVoicePicker(e,idx){
 e.stopPropagation();
 document.querySelectorAll('.voice-picker').forEach(p=>p.classList.remove('show'));
 const vp=$(`vp-${idx}`);if(!vp)return;
 vp.innerHTML=ttsVoices.map(v=>`<div class="voice-opt${v.id===ttsVoice?' active':''}" onclick="selectVoice('${v.id}',${idx})">${v.name}<span class="vgender">${v.gender}</span></div>`).join('');
 vp.classList.add('show');
}
function selectVoice(id,idx){
 ttsVoice=id;
 document.querySelectorAll('.voice-picker').forEach(p=>p.classList.remove('show'));
 speakMsg(idx,$(`ttsb-${idx}`));
}
function _stopTTS(){
 if(ttsAudio){try{ttsAudio.onended=null;ttsAudio.onerror=null;ttsAudio.pause();ttsAudio.src='';}catch(e){}ttsAudio=null}
 if(ttsBtnEl){ttsBtnEl.classList.remove('playing');ttsBtnEl=null}
 ttsLoading=false;
}
async function speakMsg(idx,btn){
 document.querySelectorAll('.voice-picker').forEach(p=>p.classList.remove('show'));
 // If this button is already playing/loading, stop immediately
 if(btn.classList.contains('playing')||ttsBtnEl===btn){_stopTTS();return}
 // Stop any other playback
 _stopTTS();
 const mc=$(`mc-${idx}`);if(!mc)return;
 const text=mc.innerText||mc.textContent;
 btn.classList.add('playing');ttsBtnEl=btn;ttsIdx=idx;ttsLoading=true;
 const r=await pywebview.api.tts_speak(text,ttsVoice);
 // Check if cancelled while loading
 if(!ttsLoading||ttsBtnEl!==btn){return}
 ttsLoading=false;
 if(r&&r.audio){
  const src=`data:audio/${r.format||'mp3'};base64,${r.audio}`;
  ttsAudio=new Audio(src);
  ttsAudio.onended=()=>{btn.classList.remove('playing');if(ttsBtnEl===btn)ttsBtnEl=null};
  ttsAudio.onerror=()=>{btn.classList.remove('playing');if(ttsBtnEl===btn)ttsBtnEl=null;showAlert('Audio playback error.')};
  ttsAudio.play();
 } else {
  btn.classList.remove('playing');if(ttsBtnEl===btn)ttsBtnEl=null;
  showAlert(r&&r.error?r.error:'TTS failed.');
 }
}
document.addEventListener('click',e=>{if(!e.target.closest('.tts-wrap'))document.querySelectorAll('.voice-picker').forEach(p=>p.classList.remove('show'))});

/* ── FOLDER FILTER DROPDOWN ── */
function renderFolderTabs(){
 // Update the filter button label/dot to reflect current selection
 const lbl=$('folder-filter-label'),dot=$('folder-filter-dot');
 if(!lbl)return;
 if(activeFolderTab==='All'){lbl.textContent='All chats';if(dot)dot.style.background='var(--t3)'}
 else{lbl.textContent=activeFolderTab;if(dot)dot.style.background=folderColors[activeFolderTab]||'#555566'}
}
function toggleFolderDrop(e){
 e.stopPropagation();
 const dd=$('folder-drop');if(!dd)return;
 const open=dd.classList.contains('show');
 if(open){dd.classList.remove('show');return}
 renderFolderDropdown();
 const btn=$('folder-filter-btn');const rect=btn.getBoundingClientRect();
 dd.style.left=rect.left+'px';dd.style.top=(rect.bottom+4)+'px';dd.style.minWidth=rect.width+'px';
 dd.classList.add('show');
}
document.addEventListener('click',e=>{
 if(!e.target.closest('#folder-filter-btn')&&!e.target.closest('#folder-drop')){
  const dd=$('folder-drop');if(dd)dd.classList.remove('show');
 }
});
function renderFolderDropdown(){
 const dd=$('folder-drop');if(!dd)return;
 const tabs=['All',...folderList];
 let html=tabs.map(t=>{
  const color=t==='All'?'var(--t3)':(folderColors[t]||'#555566');
  const active=activeFolderTab===t?' active':'';
  const dot=`<span class="fdi-dot" style="background:${color}"></span>`;
  const editBtn=t!=='All'?`<button class="fda-btn" onclick="event.stopPropagation();startFolderRename('${esc(t)}')" title="Edit">✎</button><button class="fda-btn del" onclick="event.stopPropagation();deleteFolderDrop('${esc(t)}')" title="Delete">×</button>`:'';
  return`<div class="folder-drop-item${active}" id="fdi-${esc(t)}" onclick="setFolderTab('${esc(t)}')">${dot}<span style="flex:1">${esc(t)}</span>${editBtn}</div>`;
 }).join('');
 html+=`<div class="folder-drop-sep"></div><div class="folder-drop-add"><input type="color" id="fd-color-inp" class="fd-color-btn" value="#6366f1" title="Pick folder color" oninput="document.getElementById('fd-preview-dot').style.background=this.value"><span id="fd-preview-dot" style="width:8px;height:8px;border-radius:50%;background:#6366f1;display:inline-block;flex-shrink:0;transition:background .1s"></span><input type="text" id="fd-new-inp" placeholder="New folder…" onkeydown="if(event.key==='Enter'){event.stopPropagation();addFolderFromDrop()}"><button onclick="addFolderFromDrop()">Add</button></div>`;
 dd.innerHTML=html;
}
function startFolderRename(name){
 const item=document.getElementById(`fdi-${name}`);if(!item)return;
 const color=folderColors[name]||'#555566';
 item.innerHTML=`<input type="color" id="fdr-color" value="${color}" title="Change color" style="width:22px;height:22px;border-radius:5px;border:1.5px solid rgba(255,255,255,.15);padding:0;cursor:pointer;flex-shrink:0;background:none"><div class="folder-drop-edit"><input class="folder-drop-edit-inp" id="fdr-inp" value="${esc(name)}" onkeydown="if(event.key==='Enter'){event.stopPropagation();commitFolderRename('${esc(name)}')}else if(event.key==='Escape'){renderFolderDropdown()}"></div><div class="folder-drop-actions"><button class="fda-btn" onclick="event.stopPropagation();commitFolderRename('${esc(name)}')" title="Save">✓</button><button class="fda-btn" onclick="event.stopPropagation();renderFolderDropdown()" title="Cancel">✕</button></div>`;
 // Disable the outer setFolderTab onclick so clicks inside (color picker, inputs) don't trigger it
 item.onclick=(e)=>e.stopPropagation();
 const inp=document.getElementById('fdr-inp');if(inp){inp.focus();inp.select()}
}
async function commitFolderRename(old){
 const inp=document.getElementById('fdr-inp');if(!inp)return;
 const nv=inp.value.trim();
 const colorInp=document.getElementById('fdr-color');
 const newColor=colorInp?colorInp.value:null;
 if(newColor&&newColor!==(folderColors[old]||'#555566')){
  await pywebview.api.set_folder_color(old,newColor);
  folderColors[old]=newColor;
 }
 if(!nv||nv===old){renderFolderDropdown();return}
 const ok=await pywebview.api.rename_folder(old,nv);
 if(ok){
  if(activeFolderTab===old)activeFolderTab=nv;
  const cfg=await pywebview.api.get_config();folderList=cfg.folders||[];folderColors=cfg.folder_colors||{};
  renderFolderTabs();renderFolderDropdown();refreshCL();
 }
}
async function deleteFolderDrop(name){
 const dd=$('folder-drop');if(dd)dd.classList.remove('show');
 showConfirm(`Delete folder "${name}"? Chats will move to General.`,async()=>{
  const fl=await pywebview.api.delete_folder(name);
  folderList=fl;if(activeFolderTab===name)activeFolderTab='All';
  renderFolderTabs();refreshCL();
 });
}
function updateFolderColorLive(name, color){
 folderColors[name]=color;
 const dot=document.querySelector(`#fdi-${name} .fdi-dot`);
 if(dot)dot.style.background=color;
 // also update sidebar tab dot
 renderFolderTabs();
}
async function saveFolderColor(name, color){
 folderColors[name]=color;
 await pywebview.api.set_folder_color(name,color);
 renderFolderTabs();refreshCL();
}
async function addFolderFromDrop(){
 const inp=document.getElementById('fd-new-inp');if(!inp)return;
 const name=inp.value.trim();if(!name)return;
 const colorInp=document.getElementById('fd-color-inp');
 const color=colorInp?colorInp.value:'#6366f1';
 folderColors[name]=color;  // optimistic
 folderList=await pywebview.api.add_folder(name);
 await pywebview.api.set_folder_color(name,color);
 const cfg=await pywebview.api.get_config();folderColors=cfg.folder_colors||{};
 renderFolderDropdown();refreshCL();
}
function setFolderTab(name){
 activeFolderTab=name;
 const dd=$('folder-drop');if(dd)dd.classList.remove('show');
 renderFolderTabs();refreshCL();
}
let _dragCid=null;
function onChatDragStart(e,id){_dragCid=id;e.target.classList.add('dragging');e.dataTransfer.effectAllowed='move'}
function onChatDragEnd(e){e.target.classList.remove('dragging')}
async function onDropToFolder(e,folder){
 e.preventDefault();
 if(!_dragCid)return;
 await pywebview.api.set_chat_folder(_dragCid,folder);
 _dragCid=null;refreshCL();
}

/* ── MOVE TO FOLDER MODAL ── */
let _moveFolderCid=null;
async function moveChatToFolder(id){
 _moveFolderCid=id;
 const folders=await pywebview.api.get_folders();
 const currentFolder=(await pywebview.api.get_chat_folder(id))||'General';
 const el=$('foldr-list');if(!el)return;
 el.innerHTML=folders.map(f=>{
  const color=folderColors[f]||'#555566';
  const isCur=f===currentFolder;
  return`<div class="foldr-btn${isCur?' current':''}" onclick="applyMoveToFolder('${esc(f)}')"><span class="fdot2" style="background:${color}"></span>${esc(f)}${isCur?'<span class="fcur">current</span>':''}</div>`;
 }).join('');
 $('foldr-new-inp').value='';
 $('foldrm').classList.add('show');
}
async function applyMoveToFolder(folder){
 if(!_moveFolderCid)return;
 await pywebview.api.set_chat_folder(_moveFolderCid,folder);
 _moveFolderCid=null;closeM('foldrm');refreshCL();
}
async function addFolderFromModal(){
 const inp=$('foldr-new-inp');const name=inp.value.trim();if(!name)return;
 const colorInp=$('foldr-color-inp');
 const color=colorInp?colorInp.value:'#6366f1';
 folderList=await pywebview.api.add_folder(name);
 await pywebview.api.set_folder_color(name,color);
 const cfg=await pywebview.api.get_config();folderColors=cfg.folder_colors||{};
 inp.value='';
 // re-render the folder list in modal
 await moveChatToFolder(_moveFolderCid);
}

/* ── PROMPT TEMPLATES ── */
let _templates=[],_tplQuery='';
async function openTemplates(){
 _templates=await pywebview.api.get_templates();
 _tplQuery='';
 const si=$('tpl-search');if(si)si.value='';
 renderTemplateList(_templates);
 $('tpl-name-inp').value='';$('tpl-text-inp').value='';
 $('tplm').classList.add('show');
 setTimeout(()=>{const s=$('tpl-search');if(s)s.focus()},80);
}
function filterTemplates(q){
 _tplQuery=q.toLowerCase();
 const filtered=q?_templates.filter(t=>t.name.toLowerCase().includes(_tplQuery)||t.text.toLowerCase().includes(_tplQuery)):_templates;
 renderTemplateList(filtered);
}
function clearTplSearch(){
 const s=$('tpl-search');if(s)s.value='';
 _tplQuery='';renderTemplateList(_templates);
 if(s)s.focus();
}
function renderTemplateList(list){
 const el=$('tpl-list');if(!el)return;
 if(!list.length){el.innerHTML=`<div class="tpl-empty">${_tplQuery?'No templates match your search.':'No templates saved yet. Create one below.'}</div>`;return}
 el.innerHTML=list.map(t=>`<div class="tpl-item" onclick="applyTemplate('${t.id}')"><div class="tpl-item-text"><div class="tpl-name">${esc(t.name)}</div><div class="tpl-preview">${esc(t.text)}</div></div><button class="tpl-del" onclick="event.stopPropagation();delTemplate('${t.id}')" title="Delete">×</button></div>`).join('');
}
function applyTemplate(id){
 const t=_templates.find(x=>x.id===id);if(!t)return;
 const inp=_tplForCompare?$('cmp-inp'):$('uinp');
 _tplForCompare=false;
 inp.value=t.text;autoR(inp);inp.focus();inp.setSelectionRange(t.text.length,t.text.length);
 closeM('tplm');
}
async function saveNewTemplate(){
 const name=$('tpl-name-inp').value.trim();const text=$('tpl-text-inp').value.trim();
 if(!name||!text){showAlert('Enter both a name and prompt text.');return}
 _templates=await pywebview.api.save_template(name,text);
 $('tpl-name-inp').value='';$('tpl-text-inp').value='';
 filterTemplates(_tplQuery);
}
async function delTemplate(id){
 _templates=await pywebview.api.delete_template(id);filterTemplates(_tplQuery);
}

/* ── CHAT SUMMARY ── */
let _summaryText='';
async function summarizeChat(){
 if(!cid){showAlert('Open a chat first.');return}
 $('summ-content').innerHTML='<div class="summ-loading">⏳ Generating summary…</div>';
 $('summ-copy-btn').style.display='none';
 $('summm').classList.add('show');
 const r=await pywebview.api.summarize_chat(cid);
 if(r&&r.summary){
  _summaryText=r.summary;
  $('summ-content').innerHTML=`<div id="summ-result">${rmd(r.summary)}</div>`;
  $('summ-copy-btn').style.display='inline-flex';
 } else {
  $('summ-content').innerHTML=`<div class="summ-loading" style="color:var(--dg)">Error: ${esc(r?.error||'Unknown error')}</div>`;
 }
}
function copySummary(){
 if(!_summaryText)return;
 if(navigator.clipboard&&navigator.clipboard.writeText){navigator.clipboard.writeText(_summaryText).then(()=>showAlert('Summary copied!')).catch(()=>fallbackCopy(_summaryText))}
 else fallbackCopy(_summaryText);
}

window.addEventListener('pywebviewready',init);

/* Chat Search Functionality */
(function(){
let searchQuery='';let searchMatches=[];let currentMatchIndex=-1;let originalContent=new Map();
const searchInput=$('chat-search');const searchNav=$('search-nav');const searchCount=$('search-count');const searchPrev=$('search-prev');const searchNext=$('search-next');const searchClose=$('search-close');

function saveOriginalContent(){originalContent.clear();document.querySelectorAll('.mc').forEach((el,idx)=>{originalContent.set(el.id||idx,el.innerHTML)});}
function restoreOriginalContent(){originalContent.forEach((html,id)=>{const el=document.getElementById(id);if(el)el.innerHTML=html;else{const els=document.querySelectorAll('.mc');const idx=parseInt(id);if(els[idx])els[idx].innerHTML=html}});}

function escapeRegExp(string){return string.replace(/[.*+?^${}()|[\]\\]/g,'\\$&');}

function performSearch(){searchQuery=searchInput.value.trim();if(!searchQuery){clearSearch();return;}searchMatches=[];currentMatchIndex=-1;const regex=new RegExp('\\b'+escapeRegExp(searchQuery)+'\\b','gi');document.querySelectorAll('.msg').forEach((msg,msgIdx)=>{const mc=msg.querySelector('.mc');if(!mc)return;let html=mc.innerHTML;const textContent=mc.textContent;const matches=[...textContent.matchAll(regex)];if(matches.length>0){html=html.replace(regex,match=>`<mark class="search-highlight">${match}</mark>`);mc.innerHTML=html;const highlights=mc.querySelectorAll('.search-highlight');highlights.forEach((hl,hlIdx)=>{searchMatches.push({element:hl,msgElement:msg,index:searchMatches.length});});}});updateSearchUI();if(searchMatches.length>0){currentMatchIndex=0;highlightCurrent();scrollToMatch(0);}}

function updateSearchUI(){if(searchMatches.length===0){searchCount.textContent='';searchPrev.disabled=true;searchNext.disabled=true;searchNav.style.display='none';if(searchQuery){searchInput.classList.add('no-match');setTimeout(()=>searchInput.classList.remove('no-match'),900);}}else{searchCount.textContent=`${currentMatchIndex+1}/${searchMatches.length}`;searchPrev.disabled=currentMatchIndex<=0;searchNext.disabled=currentMatchIndex>=searchMatches.length-1;searchNav.style.display='flex';}}

function highlightCurrent(){searchMatches.forEach((m,idx)=>{m.element.className=idx===currentMatchIndex?'search-highlight search-highlight-current':'search-highlight';});}

function scrollToMatch(idx){if(idx<0||idx>=searchMatches.length)return;const match=searchMatches[idx];const msgContainer=$('msgs');if(!msgContainer||!match.msgElement)return;const msgRect=match.msgElement.getBoundingClientRect();const containerRect=msgContainer.getBoundingClientRect();const relativeTop=msgRect.top-containerRect.top+msgContainer.scrollTop;msgContainer.scrollTo({top:relativeTop-100,behavior:'smooth'});}

function nextMatch(){if(searchMatches.length===0)return;currentMatchIndex=(currentMatchIndex+1)%searchMatches.length;highlightCurrent();updateSearchUI();scrollToMatch(currentMatchIndex);}
function prevMatch(){if(searchMatches.length===0)return;currentMatchIndex=(currentMatchIndex-1+searchMatches.length)%searchMatches.length;highlightCurrent();updateSearchUI();scrollToMatch(currentMatchIndex);}
function clearSearch(){searchQuery='';searchInput.value='';searchMatches=[];currentMatchIndex=-1;restoreOriginalContent();originalContent.clear();searchNav.style.display='none';}

searchInput.addEventListener('keydown',e=>{if(e.key==='Enter'){e.preventDefault();if(!originalContent.size)saveOriginalContent();restoreOriginalContent();performSearch();setTimeout(()=>searchInput.focus(),0);}if(e.key==='Escape'){e.preventDefault();clearSearch();}});
searchPrev.addEventListener('click',prevMatch);
searchNext.addEventListener('click',nextMatch);
searchClose.addEventListener('click',clearSearch);


/* ── PROFILES ── */
window.openProfileModal=async function(){
 const list=$('prof-list');
 list.innerHTML='<p style="font-size:12px;color:var(--t3);text-align:center;padding:10px 0">Loading…</p>';
 $('profm').classList.add('show');
 try{
  const profiles=await pywebview.api.get_profiles();
  let activePid=null;
  try{const cfg=await pywebview.api.get_config();activePid=(cfg&&cfg.current_profile_id)||null;}catch(e){}
  const logoutBtn=$('prof-logout-btn');
  if(logoutBtn){logoutBtn.classList.toggle('visible',!!activePid);}
  list.innerHTML='';
  if(!profiles.length){list.innerHTML='<p style="font-size:13px;color:var(--t3);text-align:center;padding:10px 0">No profiles yet. Create one below.</p>';return;}
  profiles.forEach(p=>{
   const div=document.createElement('div');
   div.className='prof-item'+(p.id===activePid?' active-prof':'');
   div.innerHTML=`<span class="prof-name">${p.name}</span>`
    +(p.id===activePid?'<span class="prof-badge">ACTIVE</span>':
     `<button class="prof-del-btn" onclick="deleteProfile('${p.id}','${p.name}')" title="Delete profile">✕</button>
      <div class="prof-pw-row" style="margin:6px 0 0;width:100%;flex:unset">
       <input type="password" placeholder="Password" id="pw-load-${p.id}" style="flex:1;min-width:0" onkeydown="if(event.key==='Enter')loadProfile('${p.id}')">
       <button class="prof-switch-btn" onclick="loadProfile('${p.id}')"><svg viewBox="0 0 24 24" width="13" height="13" fill="none" stroke="currentColor" stroke-width="2.5" stroke-linecap="round" stroke-linejoin="round"><path d="M15 3h4a2 2 0 012 2v14a2 2 0 01-2 2h-4"/><polyline points="10 17 15 12 10 7"/><line x1="15" y1="12" x2="3" y2="12"/></svg>Switch</button>
      </div>`)
    +(p.id===activePid?`<button class="prof-del-btn" onclick="deleteProfile('${p.id}','${p.name}')" title="Delete profile">✕</button>`:'');
   list.appendChild(div);
  });
 }catch(e){list.innerHTML='<p style="font-size:13px;color:var(--dg);text-align:center;padding:10px 0">Error loading profiles.</p>';}
}

window.createProfile=function(){
 const name=($('prof-new-name').value||'').trim();
 const pw=$('prof-new-pw').value;
 const pw2=$('prof-new-pw2').value;
 if(!name){showAlert('Please enter a profile name.');return;}
 if(!pw){showAlert('Please enter a password.');return;}
 if(pw!==pw2){showAlert('Passwords do not match.');return;}
 pywebview.api.create_profile(name,pw).then(r=>{
  if(r&&r.error){showAlert(r.error);return;}
  $('prof-new-name').value='';$('prof-new-pw').value='';$('prof-new-pw2').value='';
  showCopyOk('Profile "'+r.name+'" created ✓');
  openProfileModal();
 });
}

window.loadProfile=function(pid){
 const pwEl=document.getElementById('pw-load-'+pid);
 const pw=pwEl?pwEl.value:'';
 if(!pw){showAlert('Enter the password for this profile.');return;}
 pywebview.api.load_profile(pid,pw).then(r=>{
  if(r&&r.error){showAlert(r.error);return;}
  if(r.theme_overrides){
   try{
    const o=JSON.parse(r.theme_overrides);
    localStorage.setItem('themeOverrides',JSON.stringify(o));
   }catch(e){}
  }
  closeM('profm');
  showCopyOk('Loaded profile: '+r.name+' ✓');
  setTimeout(()=>reinitializeApp(),400);
 });
}

async function reinitializeApp(){
 // Reset all state
 cid=null;pf=[];stm=false;webS=false;
 activeStreams=new Set();streamBuffer={};
 sessionTokens=0;lastTotalTokens=0;
 titleOverrides={};chatTokens={};chatDrafts={};
 activeFolderTab='All';
 // Clear UI
 const msgs=$('msgs');if(msgs)msgs.innerHTML='';
 const cl=$('chatlist');if(cl)cl.innerHTML='';
 // Re-run full init
 await init();
 showE();
}

window.logoutProfile=async function(){
 await pywebview.api.logout_profile();
 closeM('profm');
 showCopyOk('Logged out');
 setTimeout(()=>reinitializeApp(),400);
}

let _delProfId=null,_delProfName=null;
window.deleteProfile=function(pid,name){
 _delProfId=pid;_delProfName=name;
 $('del-prof-msg').textContent='Delete profile "'+name+'"? This cannot be undone. Enter the profile password to confirm.';
 $('del-prof-pw').value='';
 $('del-prof-m').classList.add('show');
 setTimeout(()=>$('del-prof-pw').focus(),80);
}
window.confirmDeleteProfile=function(){
 const pw=$('del-prof-pw').value;
 if(!pw){showAlert('Enter a password.');return;}
 pywebview.api.delete_profile(_delProfId,pw).then(r=>{
  if(r&&r.error){showAlert(r.error);return;}
  closeM('del-prof-m');
  showCopyOk('Profile deleted');
  openProfileModal();
 });
}

window.addEventListener('beforeunload',()=>{
 try{
  const raw=localStorage.getItem('themeOverrides');
  if(raw)pywebview.api.update_theme_overrides(raw);
  pywebview.api.save_current_profile();
 }catch(e){}
});

const originalRenderMsgs=window.renderMsgs;
window.renderMsgs=function(msgs){if(originalRenderMsgs)originalRenderMsgs(msgs);setTimeout(()=>{saveOriginalContent();if(searchQuery)performSearch();},10);};
})();
</script></body></html>
"""

def _wipe_guest_state(cfg):
    """Wipe all guest-session data: chats, folders, theme, templates.
    
    Optimized to handle errors gracefully and ensure clean state.
    """
    # Wipe chat files with better error handling
    try:
        if os.path.exists(CHATS_DIR):
            for fn in os.listdir(CHATS_DIR):
                if fn.endswith('.json'):
                    try:
                        os.remove(os.path.join(CHATS_DIR, fn))
                    except OSError:
                        pass  # File may be locked or already deleted
    except OSError:
        pass  # Directory may not exist
    
    # Reset config fields to defaults (atomic updates)
    cfg.update({
        'folders': list(DEFAULT_CONFIG['folders']),
        'chat_folders': {},
        'folder_colors': dict(DEFAULT_CONFIG['folder_colors']),
        'prompt_templates': [],
        '_theme_overrides': '{}',
        'theme': DEFAULT_CONFIG['theme']
    })


def _start_e2ee_proxy():
    """Start the Chutes E2EE proxy Docker container and wait until it is ready.

    Returns the container ID string on success, or None if Docker is unavailable.
    """
    import shutil
    if not shutil.which("docker"):
        print("Warning: docker not found, E2EE proxy not started. Falling back to direct API.")
        return None
    try:
        # Pull image silently if not present
        subprocess.run(
            ["docker", "pull", "parachutes/e2ee-proxy:latest"],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=60
        )
        result = subprocess.run(
            [
                "docker", "run", "-d", "--rm",
                "-p", "8443:443",
                "-e", "ALLOW_NON_CONFIDENTIAL=true",
                "parachutes/e2ee-proxy:latest"
            ],
            capture_output=True, text=True, timeout=15
        )
        if result.returncode != 0:
            print(f"Warning: E2EE proxy failed to start: {result.stderr.strip()}")
            return None
        cid = result.stdout.strip()
        print(f"E2EE proxy started: {cid[:12]}, waiting for readiness...")
        # Poll /health until the proxy is accepting connections (up to 15s)
        import urllib.request
        for _ in range(30):
            try:
                with urllib.request.urlopen(
                    "https://e2ee-local-proxy.chutes.dev:8443/health",
                    timeout=1
                ):
                    print("E2EE proxy ready.")
                    return cid
            except Exception:
                time.sleep(0.5)
        print("Warning: E2EE proxy did not become ready in time. Falling back to direct API.")
        _stop_e2ee_proxy(cid)
        return None
    except Exception as e:
        print(f"Warning: E2EE proxy error: {e}")
        return None


def _stop_e2ee_proxy(cid):
    """Stop the E2EE proxy Docker container by ID."""
    if not cid:
        return
    try:
        subprocess.run(
            ["docker", "stop", cid],
            stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL, timeout=10
        )
    except Exception:
        pass


def main():
    api = Api()
    icon_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "bittensor-chat.png")
    if not os.path.exists(icon_path): icon_path = None

    # ── Start E2EE proxy ──────────────────────────────────────────────
    _e2ee_cid = _start_e2ee_proxy()
    if _e2ee_cid:
        # Proxy is up — always point at it regardless of saved config
        api.config["base_url"] = "https://e2ee-local-proxy.chutes.dev:8443/v1"
    else:
        # Proxy failed — fall back to direct API
        api.config["base_url"] = "https://llm.chutes.ai/v1"
    api._client = None  # force client reinit with new base_url

    # ── Startup guest cleanup ──────────────────────────────────────────
    # If no profile is active (guest mode), guarantee a clean slate even
    # after a crash or force-close that skipped the normal close handler.
    if not api.config.get('current_profile_id'):
        _wipe_guest_state(api.config)
        save_config(api.config)

    window = webview.create_window(APP_TITLE, html=HTML, js_api=api,
        width=1100, height=740, min_size=(700,500), background_color="#06060b", text_select=True)
    api._window = window

    def on_window_closed():
        """On close: stop E2EE proxy, save active profile, reset guest state."""
        _stop_e2ee_proxy(_e2ee_cid)
        # Always save the active profile first so nothing is lost
        try:
            api.save_current_profile()
        except Exception:
            pass
        # Clear the active profile marker
        try:
            api.config.pop('current_profile_id', None)
            save_config(api.config)
        except Exception:
            pass
        # Wipe guest state so the next launch is always clean
        try:
            _wipe_guest_state(api.config)
            save_config(api.config)
        except Exception:
            pass

    window.events.closed += on_window_closed

    # Webview configuration
    start_kwargs = {"debug": False}
    if icon_path:
        start_kwargs["icon"] = icon_path
    
    try:
        webview.start(**start_kwargs)
    except Exception:
        # Fallback to basic configuration
        try:
            webview.start(debug=False)
        except Exception as e:
            print(f"Failed to start webview: {e}")
            sys.exit(1)

if __name__ == "__main__":
    main()
